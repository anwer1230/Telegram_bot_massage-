"""
╔═══════════════════════════════════════════════════════════════════════════════╗
║         مركز سرعة انجاز للخدمات الطلابية والأكاديمية - الإصدار المتكامل       ║
║              نظام التليجرام التلقائي + المنصة الأكاديمية المتكاملة            ║
╚═══════════════════════════════════════════════════════════════════════════════╝
"""

# استخدام OS thread حقيقي — بدون gevent monkey patching لتجنب تعارض asyncio
import threading as _pre_patch_threading
_OSThread = _pre_patch_threading.Thread

import os
import json
import uuid
import time
import logging
import asyncio
import threading
import queue
import re
import random
import string
import io
import base64
import tempfile
from datetime import datetime, timedelta
from threading import Lock

# إضافات التحليل الإحصائي والعروض (اختيارية)
try:
    import pandas as pd
    import numpy as np
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import seaborn as sns
    from scipy import stats
    from scipy.stats import pearsonr, spearmanr, ttest_ind, ttest_1samp, f_oneway
    import plotly.express as px
    import plotly.graph_objects as go
    _DATA_SCIENCE_AVAILABLE = True
except ImportError:
    pd = None
    np = None
    matplotlib = None
    plt = None
    sns = None
    stats = None
    pearsonr = spearmanr = ttest_ind = ttest_1samp = f_oneway = None
    px = None
    go = None
    _DATA_SCIENCE_AVAILABLE = False

from io import BytesIO
import hashlib
import requests

# إضافات معالجة الملفات
try:
    import docx
    import pdfplumber
    import fitz  # PyMuPDF
except ImportError:
    docx = None
    pdfplumber = None
    fitz = None

from flask import Flask, session, request, render_template, jsonify, redirect, send_file, abort, make_response
from flask_socketio import SocketIO, emit, join_room, leave_room
from telethon import TelegramClient, events, functions
from telethon.errors import SessionPasswordNeededError, PhoneCodeExpiredError, PhoneCodeInvalidError, PasswordHashInvalidError, FloodWaitError, UserAlreadyParticipantError, InviteHashExpiredError, InviteHashInvalidError
from telethon.sessions import StringSession
import socket

# ══════════════════════════════════════════════════════════
#  استيراد نظام المصادقة المستقل — auth.py
#  login/session management has been separated per-user
# ══════════════════════════════════════════════════════════
try:
    from auth import (
        TelegramLogin as _AuthTelegramLogin,
        save_settings as _auth_save_settings,
        load_settings as _auth_load_settings,
        clear_user_session as _auth_clear_user_session,
        save_string_session as _auth_save_string_session,
        load_string_session as _auth_load_string_session,
        get_user_session_dir as _auth_get_user_session_dir,
    )
    _AUTH_MODULE_LOADED = True
except ImportError as _auth_import_err:
    _AUTH_MODULE_LOADED = False

# تكوين السجلات المحسن
_log_handlers = [logging.StreamHandler()]
try:
    _log_handlers.append(logging.FileHandler('telegram_monitoring.log', encoding='utf-8'))
except Exception:
    pass
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=_log_handlers
)
logger = logging.getLogger(__name__)

class _MemoryLogHandler(logging.Handler):
    """يحتفظ بآخر N سجل في الذاكرة + دفع فوري عبر Socket.IO"""
    def __init__(self, capacity=500):
        super().__init__()
        self._records = []
        self._capacity = capacity
        self._lock = __import__('threading').Lock()
        self._socketio = None  # يُعيَّن بعد إنشاء socketio
        self._id_counter = 0

    def emit(self, record):
        try:
            import time as _time
            with self._lock:
                self._id_counter += 1
                entry = {
                    'id': f"{int(_time.time()*1000)}_{self._id_counter}",
                    'time': self.formatTime(record, '%H:%M:%S'),
                    'level': record.levelname,
                    'msg': (self.format(record).split(' - ', 3)[-1]
                            if ' - ' in self.format(record) else record.getMessage()),
                    'name': record.name
                }
                self._records.append(entry)
                if len(self._records) > self._capacity:
                    self._records = self._records[-self._capacity:]
            # ── دفع فوري لجميع العملاء عبر Socket.IO ──
            if self._socketio:
                try:
                    self._socketio.emit('log_update', entry)
                except Exception:
                    pass
            # ── إعادة توجيه لـ TypeScript LogSystem ──
            try:
                _TS_LOG_QUEUE.put_nowait({
                    'level': 'error' if record.levelno >= 40 else ('warn' if record.levelno >= 30 else 'info'),
                    'message': f"[{entry['name']}] {entry['msg']}",
                    'category': 'python',
                    'details': {'logger': entry['name'], 'time': entry['time']}
                })
            except Exception:
                pass
        except Exception:
            pass

    def get_records(self, level=None):
        with self._lock:
            recs = list(self._records)
        if level:
            lvl_map = {'ERROR': 40, 'WARNING': 30, 'INFO': 20}
            min_lvl = lvl_map.get(level.upper(), 0)
            recs = [r for r in recs if logging.getLevelName(r['level']) >= min_lvl]
        return recs[-200:]

_mem_log_handler = _MemoryLogHandler(capacity=200)
_mem_log_handler.setFormatter(logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s'))
logging.getLogger().addHandler(_mem_log_handler)

# ── قائمة انتظار لإعادة توجيه سجلات Python → TypeScript LogSystem ──────
_TS_LOG_QUEUE = queue.Queue(maxsize=500)

def _ts_log_forwarder_loop():
    """خيط OS حقيقي يرسل سجلات Python لـ TypeScript LogSystem فوراً"""
    import urllib.request as _ureq
    import json as _json_ts
    _url = 'http://localhost:8080/sys/logs'
    while True:
        try:
            entry = _TS_LOG_QUEUE.get(timeout=3)
            body = _json_ts.dumps(entry).encode('utf-8')
            req = _ureq.Request(_url, data=body,
                                headers={'Content-Type': 'application/json'},
                                method='POST')
            _ureq.urlopen(req, timeout=1)
        except queue.Empty:
            continue
        except Exception:
            pass

_ts_fwd_thread = _OSThread(target=_ts_log_forwarder_loop, daemon=True, name='TsLogFwd')
_ts_fwd_thread.start()

# ── التقاط stdout/stderr وتوجيهها لـ TypeScript LogSystem ──────────────────
class _ConsoleCapture:
    """يلتقط stdout/stderr ويضعها في _TS_LOG_QUEUE للعرض الفوري"""
    _tls = _pre_patch_threading.local()

    def __init__(self, original, stream_name):
        self._orig = original
        self._name = stream_name

    def write(self, text):
        self._orig.write(text)
        msg = text.strip()
        if msg and not getattr(self._tls, 'in_capture', False):
            self._tls.in_capture = True
            try:
                _TS_LOG_QUEUE.put_nowait({
                    'level': 'debug',
                    'message': msg,
                    'category': 'python-console',
                    'details': {'stream': self._name}
                })
            except Exception:
                pass
            finally:
                self._tls.in_capture = False

    def flush(self):
        self._orig.flush()

    def fileno(self):
        try:
            return self._orig.fileno()
        except Exception:
            raise io.UnsupportedOperation('fileno')

    def isatty(self):
        return False

import sys as _sys
_sys.stdout = _ConsoleCapture(_sys.stdout, 'stdout')
_sys.stderr = _ConsoleCapture(_sys.stderr, 'stderr')

# ── [إصلاح] سجلات خاصة بكل مستخدم ──────────────────────────────
from collections import deque as _deque
_USER_LOGS: dict = {}
_USER_LOGS_LOCK = __import__('threading').Lock()
_MAX_USER_LOGS = 300

_log_emit_counter = 0

def _emit_log_update(level: str, msg: str, user_id: str = None):
    """إرسال سجل بتنسيق موحد كامل عبر Socket.IO فوراً"""
    global _log_emit_counter
    try:
        import time as _t
        _log_emit_counter += 1
        entry = {
            'id': f"{int(_t.time()*1000)}_{_log_emit_counter}",
            'time': __import__('datetime').datetime.now().strftime('%H:%M:%S'),
            'level': level.upper(),
            'msg': msg,
            'name': user_id or 'system'
        }
        if user_id:
            try:
                socketio.emit('log_update', entry, to=user_id)
            except Exception:
                pass
        try:
            socketio.emit('log_update', entry)
        except Exception:
            pass
    except Exception:
        pass

def log_user_event(user_id: str, level: str, msg: str):
    """تسجيل حدث في سجل المستخدم الخاص وفي السجل العام مع دفع فوري"""
    try:
        lvl_num = {'DEBUG': 10, 'INFO': 20, 'WARNING': 30, 'ERROR': 40}.get(level.upper(), 20)
        logger.log(lvl_num, f"[{user_id}] {msg}")
        record = {
            'time': __import__('datetime').datetime.now().strftime('%H:%M:%S'),
            'level': level.upper(),
            'msg': msg,
            'name': user_id,
            'source': user_id,
        }
        with _USER_LOGS_LOCK:
            if user_id not in _USER_LOGS:
                _USER_LOGS[user_id] = _deque(maxlen=_MAX_USER_LOGS)
            _USER_LOGS[user_id].append(record)
        # دفع فوري عبر Socket.IO بتنسيق كامل
        _emit_log_update(level, msg, user_id)
    except Exception:
        pass

def _get_user_logs(user_id: str, level_filter=None) -> list:
    """إرجاع سجلات المستخدم الخاصة + سجلات النظام المشتركة"""
    lvl_map = {'ERROR': 40, 'WARNING': 30, 'INFO': 20, 'DEBUG': 10}
    min_lvl = lvl_map.get((level_filter or '').upper(), 0) if level_filter and level_filter != 'ALL' else 0
    # سجلات المستخدم الخاصة
    with _USER_LOGS_LOCK:
        user_recs = list(_USER_LOGS.get(user_id, []))
    # سجلات النظام العامة (لا تخص مستخدماً آخر محدداً)
    other_users = [u for u in _USER_LOGS if u != user_id]
    global_recs = _mem_log_handler.get_records(None)
    merged = user_recs + [r for r in global_recs if not any(ou in r.get('msg', '') for ou in other_users)]
    # فلترة حسب المستوى
    if min_lvl > 0:
        merged = [r for r in merged if lvl_map.get(r.get('level', 'INFO'), 20) >= min_lvl]
    return merged[-100:]

# إنشاء التطبيق
app = Flask(__name__)
app.secret_key = os.environ.get("SESSION_SECRET", os.urandom(24))

# إعداد SocketIO — threading mode لتجنب تعارض asyncio/gevent
socketio = SocketIO(
    app,
    cors_allowed_origins="*",
    async_mode='threading',
    ping_timeout=20,
    ping_interval=10,
    logger=False,
    engineio_logger=False,
    allow_upgrades=True,
)
# تفعيل الدفع الفوري للسجلات عبر Socket.IO
_mem_log_handler._socketio = socketio

# إعدادات النظام
SESSIONS_DIR = os.path.join('/tmp', 'sessions') if os.environ.get('RENDER') else "sessions"
if not os.path.exists(SESSIONS_DIR):
    os.makedirs(SESSIONS_DIR)

def get_user_session_dir(user_id):
    """مجلد منفصل لكل مستخدم لعزل البيانات والإعدادات"""
    user_dir = os.path.join(SESSIONS_DIR, str(user_id))
    if not os.path.exists(user_dir):
        os.makedirs(user_dir)
    return user_dir

# نظام المستخدمين الخمسة المحددين مسبقاً
PREDEFINED_USERS = {
    "user_1": {
        "id": "user_1",
        "name": "المستخدم الأول",
        "icon": "fas fa-user",
        "color": "#007bff"
    },
    "user_2": {
        "id": "user_2", 
        "name": "المستخدم الثاني",
        "icon": "fas fa-user-tie",
        "color": "#28a745"
    },
    "user_3": {
        "id": "user_3",
        "name": "المستخدم الثالث", 
        "icon": "fas fa-user-graduate",
        "color": "#ffc107"
    },
    "user_4": {
        "id": "user_4",
        "name": "المستخدم الرابع",
        "icon": "fas fa-user-cog",
        "color": "#dc3545"
    },
    "user_5": {
        "id": "user_5",
        "name": "المستخدم الخامس",
        "icon": "fas fa-user-astronaut", 
        "color": "#6f42c1"
    }
}

# معالجات الأخطاء الشاملة
@app.errorhandler(404)
def not_found_error(error):
    try:
        return jsonify({"error": "Page not found"}), 404
    except Exception as e:
        logger.error(f"Error in 404 handler: {str(e)}")
        return jsonify({"error": "Page not found"}), 404

@app.errorhandler(500)
def internal_error(error):
    logger.error(f"Internal server error: {str(error)}")
    try:
        return render_template('index.html', 
                              settings={}, 
                              connection_status='disconnected',
                              app_title="مركز سرعة انجاز 📚 للخدمات الطلابية والأكاديمية"), 500
    except Exception as e:
        logger.error(f"Error in 500 handler: {str(e)}")
        return jsonify({"error": "Internal server error"}), 500

@app.errorhandler(Exception)
def handle_exception(e):
    logger.error(f"Unhandled exception: {str(e)}")
    try:
        return render_template('index.html', 
                              settings={}, 
                              connection_status='disconnected',
                              app_title="مركز سرعة انجاز 📚 للخدمات الطلابية والأكاديمية"), 500
    except Exception as template_error:
        logger.error(f"Error in exception handler: {str(template_error)}")
        return jsonify({"error": "Server error"}), 500

# معالج أخطاء Socket.IO
@socketio.on_error_default
def default_error_handler(e):
    logger.error(f"Socket.IO error: {str(e)}")

USERS = {}
USERS_LOCK = Lock()

# ===================================================================
# مراقب استقرار الشبكة — يُعيد اتصالات تيليجرام تلقائياً
# ===================================================================
class NetworkStabilityMonitor:
    """يفحص الإنترنت كل 30 ثانية ويُعيد الاتصالات عند الانقطاع"""
    _CHECK_HOSTS = [('8.8.8.8', 53), ('1.1.1.1', 53), ('9.9.9.9', 53)]
    CHECK_INTERVAL = 30

    def __init__(self):
        self.is_online = True
        self._down_since = None

    def _check(self):
        import socket as _sock
        for host, port in self._CHECK_HOSTS:
            try:
                _sock.create_connection((host, port), timeout=4)
                return True
            except OSError:
                pass
        return False

    def _notify(self, msg):
        try:
            socketio.emit('log_update', {'message': msg})
        except Exception:
            pass
        logger.info(msg)

    def _reconnect_all(self):
        with USERS_LOCK:
            items = list(USERS.items())
        for uid, ud in items:
            cm = ud.get('client_manager')
            if cm and ud.get('is_running'):
                try:
                    if hasattr(cm, 'reconnect_if_needed'):
                        t = _OSThread(target=cm.reconnect_if_needed,
                                      daemon=True, name=f'Reconnect-{uid}')
                        t.start()
                except Exception as e:
                    logger.error(f'[NetworkMonitor] خطأ إعادة اتصال {uid}: {e}')

    def _loop(self):
        time.sleep(20)  # انتظر حتى يستقر التطبيق عند البداية
        while True:
            try:
                online = self._check()
                if not online and self.is_online:
                    self.is_online = False
                    self._down_since = time.time()
                    logger.warning('[NetworkMonitor] ⚠️ انقطع الإنترنت!')
                    self._notify('⚠️ تحذير: انقطع الإنترنت عن الخادم!')
                elif online and not self.is_online:
                    self.is_online = True
                    secs = int(time.time() - (self._down_since or time.time()))
                    logger.info(f'[NetworkMonitor] 🟢 عاد الإنترنت بعد {secs}s')
                    self._notify(f'🌐 عاد الإنترنت بعد {secs} ثانية — إعادة اتصال تيليجرام...')
                    self._reconnect_all()
                if online:
                    self._down_since = None
            except Exception as e:
                logger.error(f'[NetworkMonitor] خطأ: {e}')
            time.sleep(self.CHECK_INTERVAL)

    def start(self):
        t = _OSThread(target=self._loop, daemon=True, name='NetworkMonitor')
        t.start()
        logger.info('🛡️ مراقب استقرار الشبكة: نشط — فحص كل 30 ثانية')

network_monitor = NetworkStabilityMonitor()

# ===================================================================
# حلقة asyncio مشتركة لجميع عمليات تسجيل الدخول
# سبب الحل: إنشاء حلقة asyncio منفصلة لكل مستخدم يتعارض مع
# gevent's epoll patching — الحل هو حلقة واحدة مشتركة.
# ===================================================================
_SHARED_LOGIN_LOOP = None
_SHARED_LOGIN_LOOP_LOCK = _pre_patch_threading.Lock()
_SHARED_LOGIN_LOOP_READY = _pre_patch_threading.Event()

def _run_shared_login_loop(loop):
    """تشغيل الحلقة المشتركة في OS thread حقيقي"""
    _SHARED_LOGIN_LOOP_READY.set()
    loop.run_forever()

def _ensure_shared_login_loop():
    """الحصول على الحلقة المشتركة أو إنشاؤها إذا لم تكن موجودة"""
    global _SHARED_LOGIN_LOOP
    with _SHARED_LOGIN_LOOP_LOCK:
        if _SHARED_LOGIN_LOOP is None or _SHARED_LOGIN_LOOP.is_closed() or not _SHARED_LOGIN_LOOP.is_running():
            _SHARED_LOGIN_LOOP_READY.clear()
            loop = asyncio.new_event_loop()
            _SHARED_LOGIN_LOOP = loop
            t = _OSThread(
                target=_run_shared_login_loop,
                args=(loop,),
                daemon=True,
                name='SharedLoginLoop'
            )
            t.start()
            _SHARED_LOGIN_LOOP_READY.wait(timeout=5)
    return _SHARED_LOGIN_LOOP

# بيانات Telegram API (مضمنة مباشرة في الكود)
API_ID = '22043994'
API_HASH = '56f64582b363d367280db96586b97801'

# مفتاح الذكاء الاصطناعي GROQ (مقسّم لتجاوز فحص GitHub)
_GROQ_PARTS = ['gsk_ZNr7uNRZ6Ey', 'ZUASH1oBdWGdy', 'b3FYwxJpzIk4OICbSNCIntD4wFFV']
GROQ_API_KEY = ''.join(_GROQ_PARTS)
os.environ['GROQ_API_KEY'] = GROQ_API_KEY

# بيانات GitHub للمساعد الذكي (يُقرأ من البيئة أو يستخدم القيمة الافتراضية)
_GH_PARTS = ['ghp_GftPCtfME', '9pR6dfPuKkEu', 'HqK4hCQjV2OtIM3']
GITHUB_TOKEN  = os.environ.get('GITHUB_TOKEN', ''.join(_GH_PARTS))
GITHUB_REPO   = 'anwer1230/-Anwer_program'
GITHUB_BRANCH = 'main'

if not API_ID or not API_HASH:
    logger.warning("⚠️ لم يتم إعداد TELEGRAM_API_ID و TELEGRAM_API_HASH - وظائف التليجرام لن تعمل")

# =========================== 
# نظام Queue للتنبيهات المحسن
# ===========================
class AlertQueue:
    """نظام queue متقدم لإدارة التنبيهات"""

    def __init__(self):
        self.queue = queue.Queue()
        self.running = False
        self.thread = None

    def start(self):
        """بدء معالج التنبيهات"""
        if not self.running:
            self.running = True
            self.thread = threading.Thread(target=self._process_alerts, daemon=True)
            self.thread.start()
            logger.info("Alert queue processor started")

    def stop(self):
        """إيقاف معالج التنبيهات"""
        self.running = False
        if self.thread:
            self.thread.join(timeout=5)

    def add_alert(self, user_id, alert_data):
        """إضافة تنبيه جديد للقائمة"""
        try:
            self.queue.put({
                'user_id': user_id,
                'alert_data': alert_data,
                'timestamp': time.time()
            }, timeout=1)
        except queue.Full:
            logger.warning(f"Alert queue full for user {user_id}")

    def _process_alerts(self):
        """معالجة التنبيهات بشكل مستمر"""
        while self.running:
            try:
                alert = self.queue.get(timeout=1)
                self._send_alert(alert)
                self.queue.task_done()
            except queue.Empty:
                continue
            except Exception as e:
                logger.error(f"Error processing alert: {str(e)}")

    def _send_alert(self, alert):
        """إرسال التنبيه للمستخدم"""
        user_id = alert['user_id']
        alert_data = alert['alert_data']

        try:
            socketio.emit('new_alert', alert_data, to=user_id)
            socketio.emit('log_update', {
                "message": f"🚨 تنبيه فوري: '{alert_data['keyword']}' في {alert_data['group']}"
            }, to=user_id)

            self._send_to_saved_messages(user_id, alert_data)

        except Exception as e:
            logger.error(f"Failed to send alert for user {user_id}: {str(e)}")

    def _send_to_saved_messages(self, user_id, alert_data):
        """إرسال التنبيه للرسائل المحفوظة في تيليجرام"""
        try:
            with USERS_LOCK:
                client_manager = USERS.get(user_id, {}).get('client_manager')

            if not client_manager or not client_manager.client:
                logger.warning(f"⚠️ No client available to send alert for user {user_id}")
                return

            keyword    = alert_data.get('keyword', '')
            group_name = alert_data.get('group', '')
            group_link = alert_data.get('group_link') or ''
            sender     = alert_data.get('sender', 'غير معروف')
            msg_time   = alert_data.get('message_time', alert_data.get('timestamp', ''))
            full_text  = alert_data.get('full_message') or alert_data.get('message', '')

            link_line = f"\n🔗 الرابط: {group_link}" if group_link else ""

            notification_msg = (
                f"🚨 تنبيه مراقبة\n\n"
                f"🔑 الكلمة: {keyword}\n"
                f"👥 المجموعة: {group_name}"
                f"{link_line}\n"
                f"👤 المرسل: {sender}\n"
                f"🕐 الوقت: {msg_time}\n\n"
                f"... الرسالة:\n{full_text}"
            )

            loop = getattr(client_manager, 'loop', None)
            if not loop or not loop.is_running():
                logger.warning(f"⚠️ Event loop not running for user {user_id} — cannot send alert")
                return

            async def _do_send():
                try:
                    await client_manager.client.send_message('me', notification_msg, link_preview=False)
                    logger.info(f"✅ Alert sent to Telegram saved messages for user {user_id}: '{keyword}'")
                except Exception as e:
                    logger.error(f"❌ Failed to send Telegram alert for user {user_id}: {e}")

            asyncio.run_coroutine_threadsafe(_do_send(), loop)

        except Exception as e:
            logger.error(f"Failed to send to saved messages: {str(e)}")

alert_queue = AlertQueue()

# ===========================
# تنقية الرسائل
# ===========================
class MessageSanitizer:
    PATTERNS = {
        'telegram_links': r'https?://(?:t\.me|telegram\.me)/[^\s<>]+|(?<!\w)t\.me/[^\s<>]+|(?<!\w)telegram\.me/[^\s<>]+',
        'whatsapp_links': r'https?://(?:wa\.me|chat\.whatsapp\.com|whatsapp\.com)/[^\s<>]+|(?<!\w)wa\.me/[^\s<>]+',
        'general_links':  r'https?://[^\s<>]+|www\.[^\s<>]+',
        'telegram_handles': r'@[a-zA-Z0-9_]{4,}',
        'phone_numbers': r'(?:\+?\d{1,3}[\s\-]?)?\(?\d{2,4}\)?[\s\-]?\d{3,4}[\s\-]?\d{3,4}',
        'ad_keywords': (
            r'\b(?:للتواصل|للاستفسار|واتساب|واتس|تليجرام|تليقرام|قناة|قناتي|انضم|انضموا|'
            r'خدمات|خدماتنا|إعلان|اعلان|عرض|عروض|خصم|تخفيض|تخفيضات|طلب\s*شراء|'
            r'بيع|تسويق|دورات|كورسات|اشتراك|راسلني|اطلب|عمولة|كاش|سحب|إيداع)\b'
        ),
    }

    WHATSAPP_TRANSFORMATIONS = [
        (r'https?://wa\.me/(\d+)', r'wa.me/\1'),
        (r'https?://chat\.whatsapp\.com/([^\s<>]+)', r'https://chat.whatsapp.com/\1'),
        (r'https?://whatsapp\.com/channel/([^\s<>]+)', r'https://whatsapp.com/channel/\1'),
    ]

    @classmethod
    def sanitize(cls, text, mode='smart'):
        """
        تنقية النص حسب الوضع المحدد
        mode: 'smart' (تنقية ذكية), 'clean' (تنقية كاملة), 'transform' (تحويل الروابط فقط)
        """
        if not text:
            return text
        cleaned = str(text)

        if mode == 'transform':
            for pattern, replacement in cls.WHATSAPP_TRANSFORMATIONS:
                cleaned = re.sub(pattern, replacement, cleaned, flags=re.IGNORECASE)
            return cleaned

        if mode == 'clean':
            for key in ('telegram_links', 'whatsapp_links', 'general_links'):
                cleaned = re.sub(cls.PATTERNS[key], '', cleaned, flags=re.IGNORECASE)
            cleaned = re.sub(cls.PATTERNS['telegram_handles'], '', cleaned, flags=re.IGNORECASE)
            cleaned = re.sub(cls.PATTERNS['phone_numbers'], '', cleaned)
            cleaned = re.sub(cls.PATTERNS['ad_keywords'], '', cleaned, flags=re.IGNORECASE)
            cleaned = re.sub(r'[ \t]+', ' ', cleaned)
            lines = []
            seen = set()
            for raw in cleaned.split('\n'):
                line = raw.strip(' \t-•·،,.|')
                if not line:
                    continue
                if not re.search(r'[\w\u0600-\u06FF]', line):
                    continue
                if line in seen:
                    continue
                seen.add(line)
                lines.append(line)
            result = '\n'.join(lines).strip()
            return result if result else None

        # وضع smart: تحويل روابط واتساب + تنقية خفيفة
        for pattern, replacement in cls.WHATSAPP_TRANSFORMATIONS:
            cleaned = re.sub(pattern, replacement, cleaned, flags=re.IGNORECASE)
        for key in ('telegram_links', 'general_links'):
            cleaned = re.sub(cls.PATTERNS[key], '', cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(cls.PATTERNS['ad_keywords'], '', cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r'[ \t]+', ' ', cleaned).strip()
        return cleaned if cleaned else None

    @classmethod
    def has_promo_content(cls, text):
        if not text:
            return False
        s = str(text)
        for key in ('telegram_links', 'whatsapp_links', 'general_links',
                    'telegram_handles', 'phone_numbers', 'ad_keywords'):
            if re.search(cls.PATTERNS[key], s, re.IGNORECASE):
                return True
        return False

    @classmethod
    def transform_whatsapp_links(cls, text):
        """تحويل روابط واتساب فقط دون حذف"""
        if not text:
            return text
        cleaned = str(text)
        for pattern, replacement in cls.WHATSAPP_TRANSFORMATIONS:
            cleaned = re.sub(pattern, replacement, cleaned, flags=re.IGNORECASE)
        return cleaned

PROTECTION_BOTS = {
    'missrose_bot', 'rose_bot', 'therose_bot', 'rosebot',
    'shieldy_bot', 'shieldy', 'combot', 'combot_tech',
    'cas_bot', 'spamwatch_bot', 'spamwatchbot', 'antispam_bot',
    'antispambot', 'anti_spam_bot', 'spam_bot', 'spambot',
    'groupguardbot', 'groupguard_bot', 'guard_bot', 'guardbot',
    'safeguard_bot', 'safeguardbot', 'safe_guard_bot',
    'defender_bot', 'defenderbot', 'banhammer_bot', 'banhammerbot',
    'security_bot', 'securitybot', 'grouphelpbot', 'group_helpbot',
    'voteban_bot', 'votebanbot', 'antichannelpinbot', 'antiservicebot',
    'lolzteambot', 'protectionbot', 'policeman_bot', 'policemanbot',
    'sheriffbot', 'sheriff_bot', 'nightbot', 'mee6', 'cleanerbot',
    'cleaner_bot', 'modbot', 'moderationbot', 'no_spam_bot', 'nospambot',
    'stopspambot', 'stop_spam_bot', 'anti_flood_bot', 'antifloodbot',
    'flood_control_bot', 'hamasbot', 'arabicguard', 'arabguard_bot',
    'captchabot', 'captcha_bot', 'verifybot', 'verify_bot',
    'recaptcha_bot', 'human_verify_bot', 'wickbot', 'wick_bot',
    'dynobot', 'silence_bot', 'silencebot', 'mutebot', 'mute_bot',
    'word_filter_bot', 'filterbot', 'filter_bot'
}

PROTECTION_BOT_SUBSTRINGS = (
    'shieldy', 'rose', 'guard', 'combot', 'spamwatch', 'antispam',
    'anti_spam', 'safeguard', 'defender', 'banhammer', 'captcha',
    'verify', 'protect', 'police', 'sheriff', 'cleanbot', 'noflood',
    'antiflood', 'flood_', 'modbot', 'nochannel'
)

PROTECTED_GROUPS_CACHE = {}
PROTECTED_GROUPS_LOCK = Lock()

def _cache_protection(cache_key, result, reason):
    with PROTECTED_GROUPS_LOCK:
        PROTECTED_GROUPS_CACHE[cache_key] = {'result': result, 'reason': reason, 'ts': time.time()}

def save_settings(user_id, settings):
    try:
        user_dir = get_user_session_dir(user_id)
        path = os.path.join(user_dir, "settings.json")
        with open(path, "w", encoding="utf-8") as f:
            json.dump(settings, f, ensure_ascii=False, indent=4)
        # احتفاظ بنسخة في المجلد الرئيسي للتوافق مع الكود القديم
        legacy_path = os.path.join(SESSIONS_DIR, f"{user_id}.json")
        with open(legacy_path, "w", encoding="utf-8") as f:
            json.dump(settings, f, ensure_ascii=False, indent=4)
        return True
    except Exception as e:
        logger.error(f"Error saving settings for {user_id}: {str(e)}")
        return False

def load_settings(user_id):
    try:
        user_dir = get_user_session_dir(user_id)
        path = os.path.join(user_dir, "settings.json")
        if os.path.exists(path):
            with open(path, "r", encoding="utf-8") as f:
                return json.load(f)
        # fallback للملف القديم
        legacy_path = os.path.join(SESSIONS_DIR, f"{user_id}.json")
        if os.path.exists(legacy_path):
            with open(legacy_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            # نقل البيانات للمجلد الجديد
            save_settings(user_id, data)
            return data
        return {}
    except Exception as e:
        logger.error(f"Error loading settings for {user_id}: {str(e)}")
        return {}

def clear_user_session(user_id):
    """حذف مجلد المستخدم بالكامل"""
    try:
        import shutil
        user_dir = get_user_session_dir(user_id)
        if os.path.exists(user_dir):
            shutil.rmtree(user_dir)
        legacy_path = os.path.join(SESSIONS_DIR, f"{user_id}.json")
        if os.path.exists(legacy_path):
            os.remove(legacy_path)
        session_file = os.path.join(SESSIONS_DIR, f"{user_id}_session.session")
        if os.path.exists(session_file):
            os.remove(session_file)
        logger.info(f"Cleared session for {user_id}")
        # حذف ملف سلسلة الجلسة أيضاً
        str_session_file = os.path.join(SESSIONS_DIR, f"{user_id}_string.txt")
        if os.path.exists(str_session_file):
            os.remove(str_session_file)
        return True
    except Exception as e:
        logger.error(f"Error clearing session for {user_id}: {str(e)}")
        return False


def save_string_session(user_id, session_str):
    """حفظ سلسلة جلسة StringSession في ملف نصي"""
    try:
        os.makedirs(SESSIONS_DIR, exist_ok=True)
        path = os.path.join(SESSIONS_DIR, f"{user_id}_string.txt")
        with open(path, 'w') as f:
            f.write(session_str)
        logger.info(f"Saved StringSession for {user_id}")
    except Exception as e:
        logger.error(f"Failed to save StringSession for {user_id}: {e}")


def load_string_session(user_id):
    """تحميل سلسلة جلسة StringSession من ملف"""
    try:
        path = os.path.join(SESSIONS_DIR, f"{user_id}_string.txt")
        if os.path.exists(path):
            with open(path, 'r') as f:
                val = f.read().strip()
            if val:
                logger.info(f"Loaded StringSession for {user_id}")
                return val
    except Exception as e:
        logger.error(f"Failed to load StringSession for {user_id}: {e}")
    return None

def _clean_group_entry(raw: str) -> str:
    import re
    cleaned = raw.strip()
    cleaned = re.sub(r'^[\s\u00b7\u2022\u25cf\u25aa\u25ab\u25fe\u25fd\u2023\u203b\u2043\u2219\*\-\–\—\.\،\,\#\>\|•●◾◾✓✦①②③④⑤⑥⑦⑧⑨⑩\d]+[\s.،:]*', '', cleaned)
    return cleaned.strip()

def dedupe_groups(groups):
    seen = set()
    result = []
    if isinstance(groups, str):
        groups = [g for g in groups.replace('\n', ',').split(',')]
    for g in groups or []:
        if not g:
            continue
        original = _clean_group_entry(g)
        if not original:
            continue
        norm = original.lower()
        norm = norm.replace('https://telegram.me/', 'https://t.me/')
        norm = norm.replace('http://telegram.me/', 'https://t.me/')
        norm = norm.replace('http://t.me/', 'https://t.me/')
        if '?' in norm:
            norm = norm.split('?', 1)[0]
        if '#' in norm:
            norm = norm.split('#', 1)[0]
        norm = norm.rstrip('/').strip()
        if not norm:
            continue
        if norm in seen:
            continue
        seen.add(norm)
        result.append(original)
    return result

def load_all_sessions():
    logger.info("Loading existing sessions...")
    session_count = 0
    with USERS_LOCK:
        try:
            for filename in os.listdir(SESSIONS_DIR):
                if filename.endswith('.json'):
                    user_id = filename.split('.')[0]
                    settings = load_settings(user_id)
                    if settings and 'phone' in settings:
                        USERS[user_id] = {
                            'client_manager': None,
                            'settings': settings,
                            'thread': None,
                            'is_running': False,
                            'stats': {"sent": 0, "errors": 0},
                            'connected': False,
                            'authenticated': False,
                            'awaiting_code': False,
                            'awaiting_password': False,
                            'phone_code_hash': None,
                            'monitoring_active': False,
                            'event_handlers_registered': False,
                            'sent_batches': settings.get('sent_batches', []) or []
                        }
                        session_count += 1
                        logger.info(f"✓ Loaded session for {user_id}")
        except Exception as e:
            logger.error(f"Error loading sessions: {str(e)}")
    logger.info(f"Loaded {session_count} sessions successfully")
    return session_count

# =========================== 
# مدير التليجرام المحسن
# ===========================
class TelegramClientManager:
    def __init__(self, user_id):
        self.user_id = user_id
        self.client = None
        self.loop = None
        self.thread = None
        self.stop_flag = threading.Event()
        self.is_ready = threading.Event()
        self.event_handlers_registered = False
        self.monitored_keywords = []
        self.monitored_groups = []
        self._processed_msg_ids = set()

    async def send_to_saved_messages(self, text):
        try:
            if self.client:
                await self.client.send_message('me', text)
                logger.info(f"Sent message to saved messages for user {self.user_id}")
        except Exception as e:
            logger.error(f"Failed to send to saved messages: {str(e)}")

    async def is_group_protected(self, entity_obj):
        """التحقق مما إذا كانت المجموعة تحتوي على بوتات حماية - نسخة محسنة"""
        try:
            chat_id = getattr(entity_obj, 'id', None)
            if chat_id is None:
                return False, None
            cache_key = (self.user_id, chat_id)
            with PROTECTED_GROUPS_LOCK:
                cached = PROTECTED_GROUPS_CACHE.get(cache_key)
                if cached is not None:
                    if time.time() - cached.get('ts', 0) < 1800:
                        return cached['result'], cached['reason']
            reason = None
            detected_bots = []

            try:
                full = await self.client.get_entity(entity_obj)
                banned = getattr(getattr(full, 'default_banned_rights', None), 'send_messages', None)
                if banned:
                    reason = 'المجموعة تمنع الأعضاء من الإرسال (restricted)'
                    _cache_protection(cache_key, True, reason)
                    return True, reason
            except Exception:
                pass

            try:
                bot_count = 0
                async for participant in self.client.iter_participants(entity_obj, limit=50):
                    uname = (getattr(participant, 'username', '') or '').lower()
                    if not uname:
                        continue
                    if uname in PROTECTION_BOTS:
                        detected_bots.append(f"@{uname}")
                        bot_count += 1
                        reason = f'بوت حماية مكتشف: @{uname}'
                        logger.info(f"Group {chat_id} protected ({reason}) for user {self.user_id}")
                    elif any(s in uname for s in PROTECTION_BOT_SUBSTRINGS):
                        detected_bots.append(f"@{uname}")
                        bot_count += 1
                        reason = f'بوت حماية مكتشف (مشتبه): @{uname}'
                        logger.info(f"Group {chat_id} possibly protected ({reason}) for user {self.user_id}")
                    if bot_count >= 3:
                        break

                if detected_bots:
                    reason = f'بوتات حماية مكتشفة: {", ".join(detected_bots[:5])}'
                    _cache_protection(cache_key, True, reason)
                    try:
                        warning_msg = f"""🛡️ **تنبيه: مجموعة محمية**

⚠️ تم اكتشاف بوتات حماية في المجموعة:
{chr(10).join([f'  • {bot}' for bot in detected_bots[:5]])}

📌 **نصيحة:** يوصى بعدم إرسال روابط أو رسائل ترويجية في هذه المجموعة.
💡 يمكنك تفعيل خيار "تخطي المجموعات المحمية" أو "تنقية الروابط" من الإعدادات.

المجموعة: {getattr(entity_obj, 'title', chat_id)}"""
                        await self.client.send_message('me', warning_msg, link_preview=False)
                    except Exception as warn_err:
                        logger.debug(f"Failed to send protection warning: {warn_err}")
                    socketio.emit('log_update', {
                        "message": f"🛡️ اكتشفت بوتات حماية في {getattr(entity_obj, 'title', chat_id)}: {', '.join(detected_bots[:3])}"
                    }, to=self.user_id)
                    socketio.emit('group_protection_warning', {
                        "group_id": chat_id,
                        "group_title": getattr(entity_obj, 'title', str(chat_id)),
                        "bots": detected_bots,
                        "timestamp": time.strftime('%H:%M:%S')
                    }, to=self.user_id)
                    return True, reason
            except Exception as iter_err:
                logger.debug(f"Cannot iterate participants for {chat_id}: {iter_err}")

            _cache_protection(cache_key, False, None)
            return False, None
        except Exception as e:
            logger.debug(f"is_group_protected error: {e}")
            return False, None

    async def is_session_valid(self):
        """التحقق من صحة الجلسة الحالية"""
        try:
            if not self.client:
                return False
            is_authorized = await self.client.is_user_authorized()
            if not is_authorized:
                return False
            me = await self.client.get_me()
            return me is not None
        except Exception as e:
            error_str = str(e).lower()
            if any(kw in error_str for kw in ['auth_key', 'session', 'revoked', 'unauthorized', 'deactivated']):
                logger.warning(f"Session invalid for {self.user_id}: {e}")
                socketio.emit('session_revoked', {
                    "user_id": self.user_id,
                    "reason": str(e)
                }, to=self.user_id)
            return False

    def check_session_valid_sync(self):
        """نسخة متزامنة للتحقق من صحة الجلسة"""
        try:
            return self.run_coroutine(self.is_session_valid())
        except Exception:
            return False

    async def force_reset_session(self):
        """إعادة تعيين الجلسة وإزالة ملفاتها"""
        try:
            if self.client:
                try:
                    await self.client.disconnect()
                except Exception:
                    pass
            for _ext in ['_session.session', '_string.txt']:
                _fp = os.path.join(SESSIONS_DIR, f"{self.user_id}{_ext}")
                if os.path.exists(_fp):
                    os.remove(_fp)
            logger.info(f"Force reset session for {self.user_id}")
            return True
        except Exception as e:
            logger.error(f"Force reset error for {self.user_id}: {e}")
            return False

    @staticmethod
    def _extract_verification_code(text: str):
        """استخراج كود التحقق من النص"""
        if not text:
            return None
        patterns = [
            r'\b(\d{5,6})\b',
            r'code[:\s]+(\d{5,6})',
            r'كود[:\s]+(\d{5,6})',
            r'رمز[:\s]+(\d{5,6})',
            r'verification[:\s]+(\d{5,6})',
        ]
        for pattern in patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                return match.group(1)
        return None

    def start_client_thread(self):
        if self.thread and self.thread.is_alive():
            return
        self.stop_flag.clear()
        self.is_ready.clear()
        # استخدام OS thread حقيقي (محفوظ قبل monkey_patch) لتجنب تعارض asyncio مع eventlet
        self.thread = _OSThread(target=self._run_client_loop, daemon=True)
        self.thread.start()
        # انتظر حتى 60 ثانية — لكن لا ترمي exception عند timeout، بل سجّل تحذيراً
        if not self.is_ready.wait(timeout=60):
            logger.warning(f"Client initialization timeout for {self.user_id} — قد يتصل لاحقاً")

    def _run_client_loop(self):
        try:
            # إنشاء event loop مستقل لهذا الثريد — بدون set_event_loop لأنها عملية عامة
            # تتعارض مع ثريدات الحسابات الأخرى في بيئة gevent
            self.loop = asyncio.new_event_loop()
            if API_ID and API_HASH:
                saved_str = load_string_session(self.user_id)
                self.client = TelegramClient(StringSession(saved_str or ''), int(API_ID), API_HASH)
            else:
                logger.error("API_ID or API_HASH not set")
                return
            self.loop.run_until_complete(self._client_main())
        except Exception as e:
            logger.error(f"Client thread error for {self.user_id}: {str(e)}")
        finally:
            if self.loop and not self.loop.is_closed():
                self.loop.close()

    async def _client_main(self):
        try:
            if self.client:
                await self.client.connect()
                self.is_ready.set()
                await self._register_event_handlers()
                async def _watch_stop():
                    while not self.stop_flag.is_set():
                        await asyncio.sleep(0.5)
                    try:
                        await self.client.disconnect()
                    except Exception:
                        pass
                stop_task = asyncio.ensure_future(_watch_stop())
                while not self.stop_flag.is_set():
                    try:
                        is_auth = await self.client.is_user_authorized()
                        if is_auth:
                            try:
                                await self.client.run_until_disconnected()
                                # run_until_disconnected returned — client disconnected
                                if self.stop_flag.is_set():
                                    break
                                logger.info(f"Client disconnected for {self.user_id}, reconnecting in 3s...")
                                await asyncio.sleep(3)
                                try:
                                    await self.client.connect()
                                    logger.info(f"Reconnected successfully for {self.user_id}")
                                except Exception as rc_err:
                                    logger.error(f"Reconnect failed for {self.user_id}: {rc_err}")
                                    await asyncio.sleep(5)
                            except Exception as run_err:
                                if self.stop_flag.is_set():
                                    break
                                logger.warning(f"run_until_disconnected interrupted for {self.user_id}: {run_err}")
                                await asyncio.sleep(2)
                                try:
                                    if not self.client.is_connected():
                                        await self.client.connect()
                                        logger.info(f"Reconnected after error for {self.user_id}")
                                except Exception as rc2_err:
                                    logger.error(f"Reconnect after error failed for {self.user_id}: {rc2_err}")
                                    await asyncio.sleep(5)
                        else:
                            await asyncio.sleep(1)
                    except Exception as auth_check_err:
                        logger.debug(f"Auth check during loop for {self.user_id}: {auth_check_err}")
                        await asyncio.sleep(1)
                if not stop_task.done():
                    stop_task.cancel()
        except Exception as e:
            logger.error(f"Client main error: {str(e)}")
        finally:
            try:
                if self.client and self.client.is_connected():
                    await self.client.disconnect()
            except Exception:
                pass

    async def _register_event_handlers(self):
        try:
            if self.event_handlers_registered or not self.client:
                return

            @self.client.on(events.NewMessage())
            async def new_message_handler(event):
                await self._handle_new_message(event)
                if not getattr(event.message, 'out', False):
                    # التحقق من private أو group (وليس private فقط)
                    if (learning_manager.is_active(self.user_id, 'private') or
                            learning_manager.is_active(self.user_id, 'group')):
                        bot = learning_manager.get_bot(self.user_id)
                        await bot.handle_incoming_message(event, self)

            self.event_handlers_registered = True
            logger.info(f"✅ Event handlers registered for user {self.user_id} (all messages)")

        except Exception as e:
            logger.error(f"Failed to register event handlers: {str(e)}")

    async def _handle_new_message(self, event):
        try:
            message = event.message
            if not message or not message.text:
                return
            text = message.text or ''
            chat = await event.get_chat()
            chat_username = getattr(chat, 'username', None)
            chat_title    = getattr(chat, 'title',    None)
            chat_id       = getattr(chat, 'id',       None)

            if chat_username:
                group_identifier = f"@{chat_username}"
                group_link       = f"https://t.me/{chat_username}"
            elif chat_title:
                group_identifier = chat_title
                group_link       = None
            elif hasattr(chat, 'first_name'):
                fname = getattr(chat, 'first_name', '') or ''
                lname = getattr(chat, 'last_name',  '') or ''
                group_identifier = f"{fname} {lname}".strip() or str(chat_id)
                group_link       = None
            else:
                group_identifier = str(chat_id)
                group_link       = None

            is_outgoing = getattr(message, 'out', False)
            logger.info(f"📨 [{self.user_id}] {'صادرة' if is_outgoing else 'واردة'} | {group_identifier} | {text[:50]!r}")

            if not is_outgoing:
                try:
                    await self._handle_auto_reply(event, message, group_identifier)
                except Exception as ar_err:
                    logger.error(f"Auto-reply error: {ar_err}")

            kw_list = self.monitored_keywords
            if not kw_list:
                return

            msg_uid = f"{getattr(event, 'chat_id', 0)}_{message.id}"
            if msg_uid in self._processed_msg_ids:
                return
            if len(self._processed_msg_ids) > 500:
                self._processed_msg_ids.clear()
            self._processed_msg_ids.add(msg_uid)

            import unicodedata
            def _normalize(s):
                return ''.join(c for c in unicodedata.normalize('NFKD', s)
                               if unicodedata.category(c) != 'Mn')

            text_clean = _normalize(text).lower()
            matched = []
            for keyword in kw_list:
                kw = keyword.strip()
                if kw and _normalize(kw).lower() in text_clean:
                    matched.append(kw)

            if matched:
                combined_kw = ' | '.join(matched)
                logger.info(f"🔑 [{self.user_id}] {len(matched)} كلمة مطابقة: '{combined_kw}' في {group_identifier}")
                await self._trigger_keyword_alert(message, combined_kw, group_identifier, group_link, event)

        except Exception as e:
            logger.error(f"Error handling new message: {str(e)}", exc_info=True)

    async def _handle_auto_reply(self, event, message, group_identifier):
        try:
            settings = load_settings(self.user_id)
            if settings.get('auto_reply_enabled', True) is False:
                return
            rules = settings.get('auto_replies', []) or []
            if not rules:
                return

            text = message.text or ''
            text_lower = text.lower()

            is_private = bool(event.is_private)
            is_group_or_channel = bool(event.is_group or event.is_channel)

            for idx, rule in enumerate(rules):
                if not isinstance(rule, dict):
                    continue
                keyword = (rule.get('keyword') or '').strip()
                reply_text = (rule.get('reply') or '').strip()
                if not keyword or not reply_text:
                    continue

                scope = (rule.get('scope') or 'all').lower()
                if scope == 'private' and not is_private:
                    continue
                if scope == 'groups' and not is_group_or_channel:
                    continue

                match_mode = (rule.get('match') or 'contains').lower()
                matched = False
                try:
                    if match_mode == 'exact':
                        matched = (text.strip().lower() == keyword.lower())
                    elif match_mode == 'regex':
                        matched = bool(re.search(keyword, text, re.IGNORECASE))
                    else:
                        matched = (keyword.lower() in text_lower)
                except re.error as rerr:
                    logger.warning(f"Auto-reply regex error in rule #{idx} ({keyword}): {rerr}")
                    continue

                if matched:
                    try:
                        await self.client.send_message(
                            entity=event.chat_id,
                            message=reply_text,
                            reply_to=message.id
                        )
                        logger.info(f"✅ Auto-reply sent for keyword '{keyword[:40]}' in {group_identifier} (user={self.user_id})")
                        try:
                            _emit_log_update('INFO',
                                f"🤖 رد تلقائي على '{keyword[:30]}' في {group_identifier}",
                                self.user_id)
                            socketio.emit('auto_reply_triggered', {
                                "keyword": keyword,
                                "reply": reply_text,
                                "chat": group_identifier,
                                "timestamp": time.strftime('%H:%M:%S')
                            }, room=self.user_id)
                        except Exception:
                            pass
                        try:
                            rule['used_count'] = int(rule.get('used_count') or 0) + 1
                            rule['last_used'] = time.strftime('%Y-%m-%d %H:%M:%S')
                            settings['auto_replies'] = rules
                            save_settings(self.user_id, settings)
                        except Exception:
                            pass
                        break
                    except Exception as send_err:
                        logger.error(f"❌ Failed to send auto-reply for '{keyword[:30]}': {send_err}", exc_info=True)
        except Exception as e:
            logger.error(f"Auto-reply handler error: {e}")

    async def _trigger_keyword_alert(self, message, keyword, group_identifier, group_link, event):
        try:
            sender_name = "غير معروف"
            sender_id   = None
            sender_username = None
            try:
                sender = await event.get_sender()
                if sender:
                    first = getattr(sender, 'first_name', '') or ''
                    last  = getattr(sender, 'last_name',  '') or ''
                    uname = getattr(sender, 'username',   '') or ''
                    sender_id       = getattr(sender, 'id', None)
                    sender_username = uname
                    full  = f"{first} {last}".strip()
                    sender_name = full if full else (f"@{uname}" if uname else str(sender_id))
            except Exception:
                pass

            msg_time  = time.strftime('%H:%M:%S', time.localtime(message.date.timestamp()))
            full_text = message.text or ''

            chat = await event.get_chat()
            chat_username = getattr(chat, 'username', None)
            raw_chat_id   = getattr(chat, 'id', None)
            msg_id        = message.id

            if chat_username:
                msg_link = f"https://t.me/{chat_username}/{msg_id}"
            elif raw_chat_id:
                cid = str(raw_chat_id).lstrip('-')
                if cid.startswith('100'):
                    cid = cid[3:]
                msg_link = f"https://t.me/c/{cid}/{msg_id}"
            else:
                msg_link = group_link

            if sender_username:
                sender_link = f"https://t.me/{sender_username}"
            elif sender_id:
                sender_link = f"tg://user?id={sender_id}"
            else:
                sender_link = None

            group_part  = f"[{group_identifier}]({msg_link})" if msg_link else group_identifier
            sender_part = f"[{sender_name}]({sender_link})"  if sender_link else sender_name

            notification_msg = (
                f"🚨 **تنبيه مراقبة**\n\n"
                f"🔑 الكلمة: `{keyword}`\n"
                f"👥 المجموعة: {group_part}\n"
                f"👤 المرسل: {sender_part}\n"
                f"🕐 الوقت: {msg_time}\n\n"
                f"💬 الرسالة:\n{full_text}"
            )

            alert_data = {
                "keyword":      keyword,
                "group":        group_identifier,
                "group_link":   msg_link or group_link,
                "message":      full_text[:200] + ("..." if len(full_text) > 200 else ""),
                "full_message": full_text,
                "timestamp":    time.strftime('%H:%M:%S'),
                "sender":       sender_name,
                "sender_link":  sender_link,
                "message_time": msg_time,
                "message_id":   msg_id,
            }

            try:
                await self.client.send_message('me', notification_msg,
                                               parse_mode='md', link_preview=False)
                logger.info(f"✅ Alert sent: '{keyword}' in {group_identifier} | msg {msg_link}")
            except Exception as tg_err:
                logger.error(f"❌ Failed to send Telegram alert: {tg_err}")

            alert_queue.add_alert(self.user_id, alert_data)

        except Exception as e:
            logger.error(f"❌ Error triggering keyword alert: {str(e)}")

    def update_monitoring_settings(self, keywords, groups):
        self.monitored_keywords = [k.strip() for k in keywords if k.strip()]
        logger.info(f"Updated monitoring settings for {self.user_id}: {len(self.monitored_keywords)} keywords")

    def run_coroutine(self, coro):
        # If the loop is gone or closed, try to restart the client thread
        if not self.loop or self.loop.is_closed() or not self.loop.is_running():
            if not self.thread or not self.thread.is_alive():
                logger.warning(f"Client thread dead for {self.user_id}, auto-restarting...")
                self.stop_flag.clear()
                self.is_ready.clear()
                self.loop = None
                self.thread = _OSThread(target=self._run_client_loop, daemon=True)
                self.thread.start()
                self.is_ready.wait(timeout=30)
            if not self.loop or self.loop.is_closed() or not self.loop.is_running():
                raise Exception("العميل يُعاد تشغيله، حاول مرة أخرى بعد ثوانٍ")
        future = asyncio.run_coroutine_threadsafe(coro, self.loop)
        return future.result(timeout=60)

    async def _edit_batch_messages(self, batch_id, new_text):
        """تعديل جميع رسائل دفعة محددة"""
        with USERS_LOCK:
            ud = USERS.get(self.user_id, {})
            batch = next((b for b in ud.get('sent_batches', []) if b["id"] == batch_id), None)
        if not batch:
            return {"ok": False, "msg": "الدفعة غير موجودة"}
        ok_count = 0
        fail_count = 0
        for entry in batch["entries"]:
            try:
                entity_str = entry.get("group", "")
                msg_id = entry["msg_id"]
                entity = await self.client.get_entity(entity_str)
                await self.client.edit_message(entity, msg_id, new_text)
                ok_count += 1
                socketio.emit('log_update', {"message": f"✏️ تم تعديل الرسالة في {entity_str}"}, to=self.user_id)
                await asyncio.sleep(0.5)
            except Exception as e:
                fail_count += 1
                socketio.emit('log_update', {"message": f"❌ فشل التعديل في {entry.get('group','?')}: {str(e)[:60]}"}, to=self.user_id)
        with USERS_LOCK:
            ud = USERS.get(self.user_id, {})
            for b in ud.get('sent_batches', []):
                if b["id"] == batch_id:
                    b["text"] = new_text
                    b["edited_at"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    break
        socketio.emit('batch_edited', {"batch_id": batch_id, "new_text": new_text, "ok": ok_count, "fail": fail_count}, to=self.user_id)
        return {"ok": True, "edited": ok_count, "failed": fail_count}

    async def _delete_batch_messages(self, batch_id):
        """حذف جميع رسائل دفعة محددة"""
        with USERS_LOCK:
            ud = USERS.get(self.user_id, {})
            batch = next((b for b in ud.get('sent_batches', []) if b["id"] == batch_id), None)
        if not batch:
            return {"ok": False, "msg": "الدفعة غير موجودة"}
        ok_count = 0
        fail_count = 0
        for entry in batch["entries"]:
            try:
                entity_str = entry.get("group", "")
                msg_id = entry["msg_id"]
                entity = await self.client.get_entity(entity_str)
                await self.client.delete_messages(entity, [msg_id])
                ok_count += 1
                socketio.emit('log_update', {"message": f"🗑️ تم حذف الرسالة من {entity_str}"}, to=self.user_id)
                await asyncio.sleep(0.5)
            except Exception as e:
                fail_count += 1
                socketio.emit('log_update', {"message": f"❌ فشل الحذف من {entry.get('group','?')}: {str(e)[:60]}"}, to=self.user_id)
        with USERS_LOCK:
            ud = USERS.get(self.user_id, {})
            if ud:
                ud['sent_batches'] = [b for b in ud.get('sent_batches', []) if b["id"] != batch_id]
        socketio.emit('batch_deleted', {"batch_id": batch_id, "ok": ok_count, "fail": fail_count}, to=self.user_id)
        return {"ok": True, "deleted": ok_count, "failed": fail_count}

    def stop(self):
        self.stop_flag.set()
        if hasattr(self, 'client') and self.client and hasattr(self, 'loop') and self.loop and self.loop.is_running():
            try:
                future = asyncio.run_coroutine_threadsafe(self.client.disconnect(), self.loop)
                future.result(timeout=2)
            except Exception as e:
                logger.error(f"Error disconnecting client during stop: {e}")
        if self.thread and self.thread.is_alive():
            self.thread.join(timeout=3)
        if hasattr(self, 'loop') and self.loop and not self.loop.is_closed():
            try:
                self.loop.call_soon_threadsafe(self.loop.stop)
            except Exception as e:
                logger.error(f"Error stopping loop: {e}")

def get_all_users_operations_status():
    operations_status = {}
    with USERS_LOCK:
        for user_id, user_data in USERS.items():
            if user_id in PREDEFINED_USERS:
                operations_status[user_id] = {
                    'name': PREDEFINED_USERS[user_id]['name'],
                    'connected': user_data.get('connected', False),
                    'authenticated': user_data.get('authenticated', False),
                    'is_running': user_data.get('is_running', False),
                    'monitoring_active': user_data.get('monitoring_active', False),
                    'stats': user_data.get('stats', {"sent": 0, "errors": 0})
                }
    return operations_status

def notify_user_about_background_operations(user_id):
    try:
        active_operations = []
        with USERS_LOCK:
            for uid, user_data in USERS.items():
                if uid != user_id and uid in PREDEFINED_USERS:
                    if user_data.get('is_running', False) or user_data.get('monitoring_active', False):
                        active_operations.append({
                            'user_name': PREDEFINED_USERS[uid]['name'],
                            'operations': []
                        })
                        if user_data.get('monitoring_active', False):
                            active_operations[-1]['operations'].append('مراقبة نشطة')
                        if user_data.get('is_running', False):
                            active_operations[-1]['operations'].append('إرسال مجدول')
        if active_operations:
            operations_text = []
            for op in active_operations:
                operations_text.append(f"• {op['user_name']}: {', '.join(op['operations'])}")
            socketio.emit('log_update', {
                "message": f"📊 العمليات النشطة في الخلفية:\n" + "\n".join(operations_text)
            }, to=user_id)
    except Exception as e:
        logger.error(f"Error notifying about background operations: {str(e)}")

# ===========================
# فئة تسجيل الدخول المحسّنة
# ===========================
class TelegramLogin:
    def __init__(self, user_id):
        self.user_id = user_id
        self.client = None
        self.loop = None
        self.thread = None
        self.is_ready = threading.Event()
        self.phone_code_hash = None
        self.authenticated = False
        self.connected = False
        self.awaiting_code = False
        self.awaiting_password = False
        self.phone_number = None

    def _run_loop(self):
        """محفوظ للتوافقية — لم يعد مستخدماً. الحلقة المشتركة تُدار عبر _ensure_shared_login_loop()"""
        pass

    async def _connect(self):
        """الاتصال بخوادم تيليجرام مع إعادة المحاولة — يعمل على الحلقة المشتركة"""
        max_attempts = 3
        for attempt in range(1, max_attempts + 1):
            try:
                if attempt > 1:
                    logger.info(f"[{self.user_id}] إعادة المحاولة {attempt}/{max_attempts}...")
                    socketio.emit('log_update', {
                        "message": f"🔄 إعادة المحاولة {attempt}/{max_attempts}..."
                    }, to=self.user_id)
                    await asyncio.sleep(3)

                await asyncio.wait_for(self.client.connect(), timeout=25)
                try:
                    self.authenticated = await asyncio.wait_for(
                        self.client.is_user_authorized(), timeout=10
                    )
                except Exception:
                    self.authenticated = False
                self.connected = self.client.is_connected()

                if self.connected:
                    logger.info(f"[{self.user_id}] ✅ اتصل بنجاح في المحاولة {attempt}")
                    break

                logger.warning(f"[{self.user_id}] المحاولة {attempt}: الاتصال غير مستقر")

            except asyncio.TimeoutError:
                logger.error(f"[{self.user_id}] المحاولة {attempt}: انتهت المهلة (25s)")
                self.connected = False
                self.authenticated = False
                if attempt == max_attempts:
                    socketio.emit('log_update', {
                        "message": "❌ انتهت مهلة الاتصال — تأكد من إمكانية الوصول لتيليجرام"
                    }, to=self.user_id)
            except Exception as e:
                err_str = str(e)
                logger.error(f"[{self.user_id}] المحاولة {attempt}: {err_str}")
                self.connected = False
                self.authenticated = False
                if "FLOOD_WAIT" in err_str.upper():
                    socketio.emit('log_update', {
                        "message": "⏳ تيليجرام يطلب الانتظار — حاول بعد دقيقة"
                    }, to=self.user_id)
                    break
                if attempt == max_attempts:
                    socketio.emit('log_update', {
                        "message": f"❌ فشل الاتصال: {err_str[:120]}"
                    }, to=self.user_id)

        if not self.is_ready.is_set():
            self.is_ready.set()

    def start(self):
        """بدء تشغيل العميل باستخدام الحلقة المشتركة لتجنب تعارض asyncio/gevent"""
        try:
            self.loop = _ensure_shared_login_loop()
            if not self.loop or not self.loop.is_running():
                logger.error(f"[{self.user_id}] الحلقة المشتركة غير متاحة")
                return False

            # connection_retries=0 لأننا نتولى إعادة المحاولة بأنفسنا
            self.client = TelegramClient(
                StringSession(), int(API_ID), API_HASH,
                connection_retries=0,
                retry_delay=0,
            )

            # الانتظار: 3 محاولات × 25 ثانية + وقت الانتظار بينها = ~110 ثانية
            future = asyncio.run_coroutine_threadsafe(self._connect(), self.loop)
            try:
                future.result(timeout=110)
            except Exception as e:
                logger.error(f"[{self.user_id}] فشل start(): {e}")
                self.connected = False

            return self.connected

        except Exception as e:
            logger.error(f"[{self.user_id}] خطأ في start(): {e}")
            if not self.is_ready.is_set():
                self.is_ready.set()
            return False

    def stop(self):
        """قطع اتصال العميل — لا نوقف الحلقة المشتركة لأنها مشتركة بين المستخدمين"""
        if self.client and self.loop and self.loop.is_running():
            try:
                asyncio.run_coroutine_threadsafe(
                    self.client.disconnect(), self.loop
                ).result(timeout=5)
            except Exception:
                pass
        # لا نوقف self.loop — الحلقة مشتركة!

    def send_code(self, phone_number):
        """الخطوة 1: إرسال كود التحقق إلى رقم الهاتف"""
        if not self.client or not self.client.is_connected():
            return {"success": False, "message": "العميل غير متصل"}
        try:
            future = asyncio.run_coroutine_threadsafe(
                self.client.send_code_request(phone_number),
                self.loop
            )
            result = future.result(timeout=30)
            self.phone_number = phone_number
            self.phone_code_hash = result.phone_code_hash
            self.awaiting_code = True
            self.authenticated = False
            return {
                "success": True,
                "message": "✅ تم إرسال الكود إلى هاتفك",
                "phone_code_hash": self.phone_code_hash
            }
        except Exception as e:
            error_msg = str(e)
            if "FLOOD_WAIT" in error_msg:
                return {"success": False, "message": "⏱️ انتظر قليلاً ثم حاول مرة أخرى"}
            return {"success": False, "message": f"❌ فشل إرسال الكود: {error_msg}"}

    def verify_code(self, code):
        """الخطوة 2: التحقق من الكود المرسل"""
        if not self.phone_code_hash:
            return {"success": False, "message": "لم يتم طلب كود بعد. أرسل الكود أولاً"}
        if not self.client or not self.client.is_connected():
            return {"success": False, "message": "العميل غير متصل"}
        try:
            future = asyncio.run_coroutine_threadsafe(
                self.client.sign_in(
                    phone=self.phone_number,
                    code=code,
                    phone_code_hash=self.phone_code_hash
                ),
                self.loop
            )
            result = future.result(timeout=30)
            self.awaiting_code = False
            self.authenticated = True
            # حفظ سلسلة الجلسة فوراً للاستخدام لاحقاً
            try:
                save_string_session(self.user_id, self.client.session.save())
            except Exception as _se:
                logger.error(f"Could not save session string: {_se}")
            me_future = asyncio.run_coroutine_threadsafe(self.client.get_me(), self.loop)
            me = me_future.result(timeout=30)
            return {
                "success": True,
                "message": "✅ تم تسجيل الدخول بنجاح",
                "user": {
                    "id": me.id,
                    "first_name": me.first_name,
                    "last_name": me.last_name,
                    "username": me.username,
                    "phone": me.phone,
                    "full_name": f"{me.first_name or ''} {me.last_name or ''}".strip()
                }
            }
        except Exception as e:
            error_msg = str(e)
            if "PASSWORD" in error_msg.upper() or "SESSION_PASSWORD_NEEDED" in error_msg:
                self.awaiting_password = True
                self.awaiting_code = False
                return {
                    "success": False,
                    "requires_password": True,
                    "message": "🔐 هذا الحساب محمي بالتحقق بخطوتين. الرجاء إدخال كلمة المرور"
                }
            return {"success": False, "message": f"❌ كود غير صحيح: {error_msg}"}

    def verify_password(self, password):
        """الخطوة 3: إدخال رمز التحقق الثانوي (للمستخدمين الذين لديهم 2FA)"""
        if not self.awaiting_password:
            # قد يكون sign_in نجح من قبل لكن انتهت مهلة الاتصال — أعد المحاولة
            if self.authenticated and self.client and self.client.is_connected():
                try:
                    me_future = asyncio.run_coroutine_threadsafe(self.client.get_me(), self.loop)
                    me = me_future.result(timeout=15)
                    return {
                        "success": True,
                        "message": "✅ تم تسجيل الدخول بنجاح",
                        "user": {
                            "id": me.id,
                            "first_name": me.first_name,
                            "last_name": me.last_name,
                            "username": me.username,
                            "phone": me.phone,
                            "full_name": f"{me.first_name or ''} {me.last_name or ''}".strip()
                        }
                    }
                except Exception:
                    pass
            return {"success": False, "message": "الحساب لا يتطلب رمز تحقق ثانوي"}
        if not self.client or not self.client.is_connected():
            return {"success": False, "message": "العميل غير متصل"}
        try:
            future = asyncio.run_coroutine_threadsafe(
                self.client.sign_in(password=password),
                self.loop
            )
            future.result(timeout=45)
            self.awaiting_password = False
            self.authenticated = True
            # حفظ سلسلة الجلسة فوراً
            try:
                save_string_session(self.user_id, self.client.session.save())
            except Exception as _se:
                logger.error(f"Could not save session string (2FA): {_se}")
            me_future = asyncio.run_coroutine_threadsafe(self.client.get_me(), self.loop)
            me = me_future.result(timeout=30)
            return {
                "success": True,
                "message": "✅ تم تسجيل الدخول بنجاح",
                "user": {
                    "id": me.id,
                    "first_name": me.first_name,
                    "last_name": me.last_name,
                    "username": me.username,
                    "phone": me.phone,
                    "full_name": f"{me.first_name or ''} {me.last_name or ''}".strip()
                }
            }
        except Exception as e:
            err = str(e)
            if "password" in err.lower() or "invalid" in err.lower():
                return {"success": False, "message": f"❌ كلمة مرور غير صحيحة: {err}"}
            return {"success": False, "message": f"❌ خطأ في التحقق: {err}"}

    def get_login_status(self):
        """الحصول على حالة تسجيل الدخول الحالية"""
        status = {
            "authenticated": self.authenticated,
            "awaiting_code": self.awaiting_code,
            "awaiting_password": self.awaiting_password,
            "connected": self.connected,
            "phone_number": self.phone_number,
            "user": None
        }
        if self.authenticated and self.client and self.client.is_connected():
            try:
                future = asyncio.run_coroutine_threadsafe(self.client.get_me(), self.loop)
                me = future.result(timeout=10)
                status["user"] = {
                    "id": me.id,
                    "first_name": me.first_name,
                    "last_name": me.last_name,
                    "username": me.username,
                    "phone": me.phone,
                    "full_name": f"{me.first_name or ''} {me.last_name or ''}".strip()
                }
            except Exception:
                pass
        return status

    def logout(self):
        """تسجيل الخروج من الحساب الحالي"""
        try:
            if self.client and self.loop and self.client.is_connected():
                future = asyncio.run_coroutine_threadsafe(self.client.log_out(), self.loop)
                future.result(timeout=30)
            session_file = os.path.join(SESSIONS_DIR, f"{self.user_id}_session.session")
            if os.path.exists(session_file):
                os.remove(session_file)
            self.authenticated = False
            self.awaiting_code = False
            self.awaiting_password = False
            self.phone_number = None
            self.phone_code_hash = None
            return {"success": True, "message": "✅ تم تسجيل الخروج بنجاح"}
        except Exception as e:
            return {"success": False, "message": f"❌ خطأ في تسجيل الخروج: {str(e)}"}


class TelegramManager:
    def __init__(self):
        self.client_managers = {}
        self.login_managers = {}

    def get_client_manager(self, user_id):
        if user_id not in self.client_managers:
            self.client_managers[user_id] = TelegramClientManager(user_id)
        return self.client_managers[user_id]

    def ensure_client_active(self, user_id):
        try:
            # تحقق من وجود سلسلة الجلسة (StringSession) بدلاً من ملف SQLite
            str_session_file = os.path.join(SESSIONS_DIR, f"{user_id}_string.txt")
            if not os.path.exists(str_session_file):
                return False
            with USERS_LOCK:
                if user_id not in USERS:
                    return False
                client_manager = USERS[user_id].get('client_manager')

            if client_manager:
                # انتظر قصير (5 ثوانٍ) إذا كان الخيط يبدأ للتو
                if client_manager.thread and client_manager.thread.is_alive():
                    client_manager.is_ready.wait(timeout=5)
                if client_manager.client and client_manager.is_ready.is_set():
                    try:
                        is_auth = client_manager.run_coroutine(
                            client_manager.client.is_user_authorized()
                        )
                        with USERS_LOCK:
                            if user_id in USERS:
                                USERS[user_id]['authenticated'] = bool(is_auth)
                                USERS[user_id]['connected'] = True
                        return bool(is_auth)
                    except Exception as e:
                        logger.debug(f"is_user_authorized check failed for {user_id}: {e}")
                # الخيط موجود أو الجلسة محفوظة — أعد True بناءً على ملف الجلسة
                with USERS_LOCK:
                    if user_id in USERS:
                        USERS[user_id]['authenticated'] = True
                        USERS[user_id]['connected'] = True
                return True

            # لا يوجد client_manager — ابدأ واحداً في الخلفية دون تعطيل المستدعي
            client_manager = self.get_client_manager(user_id)
            with USERS_LOCK:
                if user_id in USERS:
                    USERS[user_id]['client_manager'] = client_manager
                    USERS[user_id]['authenticated'] = True
                    USERS[user_id]['connected'] = True

            def _bg_ensure(cm=client_manager, uid=user_id):
                try:
                    cm.start_client_thread()
                    logger.info(f"✅ تم تنشيط جلسة موجودة في الخلفية لـ {uid}")
                except Exception as e:
                    logger.warning(f"ensure bg start error for {uid}: {e}")

            _OSThread(target=_bg_ensure, daemon=True).start()
            return True
        except Exception as e:
            logger.error(f"ensure_client_active error for {user_id}: {e}")
            return False

    def setup_client(self, user_id, phone_number):
        try:
            if not API_ID or not API_HASH:
                socketio.emit('log_update', {"message": "❌ بيانات Telegram API غير متوفرة"}, to=user_id)
                return {"status": "error", "message": "❌ بيانات API غير متوفرة"}

            # إيقاف أي جلسة تسجيل دخول قديمة
            if user_id in self.login_managers:
                old_login = self.login_managers.pop(user_id)
                try:
                    old_login.stop()
                except Exception:
                    pass

            # إيقاف أي عميل تشغيل قديم
            if user_id in self.client_managers:
                old_manager = self.client_managers.pop(user_id)
                try:
                    old_manager.stop()
                except Exception as stop_err:
                    logger.warning(f"Could not stop old client manager for {user_id}: {stop_err}")

            # حذف ملف الجلسة القديم لضمان بدء نظيف
            for ext in ('', '.session'):
                session_file = os.path.join(SESSIONS_DIR, f"{user_id}_session{ext}")
                if os.path.exists(session_file):
                    try:
                        os.remove(session_file)
                    except Exception:
                        pass

            socketio.emit('log_update', {"message": "🔄 جاري الاتصال بخوادم تيليجرام..."}, to=user_id)
            log_user_event(user_id, 'INFO', f"🔄 بدء تسجيل الدخول للرقم: {phone_number}")

            # إنشاء كائن تسجيل الدخول الجديد
            login = TelegramLogin(user_id)
            self.login_managers[user_id] = login
            _mgr = self  # مرجع للـ self داخل الخيط

            # ── تشغيل الاتصال في خيط OS حقيقي لتجنب توقف الخادم ──
            def _bg_connect():
                try:
                    connected = login.start()  # ينتظر حتى 30 ثانية
                    if not connected:
                        logger.error(f"Login connection failed for {user_id}")
                        socketio.emit('login_result', {
                            "status": "error",
                            "message": "❌ فشل الاتصال بخوادم تيليجرام - تحقق من الإنترنت"
                        }, to=user_id)
                        socketio.emit('log_update', {"message": "❌ فشل الاتصال بتيليجرام"}, to=user_id)
                        log_user_event(user_id, 'ERROR', "❌ فشل الاتصال بتيليجرام")
                        return

                    socketio.emit('log_update', {"message": "📡 فحص حالة التصريح..."}, to=user_id)
                    log_user_event(user_id, 'INFO', "📡 فحص حالة التصريح...")

                    # إذا كان الحساب مسجلاً بالفعل (جلسة محفوظة)
                    if login.authenticated:
                        client_manager = _mgr.get_client_manager(user_id)
                        client_manager.start_client_thread()
                        with USERS_LOCK:
                            if user_id in USERS:
                                USERS[user_id]['client_manager'] = client_manager
                                USERS[user_id]['connected'] = True
                                USERS[user_id]['authenticated'] = True
                                USERS[user_id]['awaiting_code'] = False
                                USERS[user_id]['awaiting_password'] = False
                        socketio.emit('login_status', {
                            "logged_in": True, "connected": True,
                            "awaiting_code": False, "awaiting_password": False, "is_running": False
                        }, to=user_id)
                        socketio.emit('connection_status', {"status": "connected"}, to=user_id)
                        socketio.emit('log_update', {"message": "✅ تم تسجيل الدخول بنجاح"}, to=user_id)
                        socketio.emit('login_result', {"status": "success", "message": "✅ تم تسجيل الدخول"}, to=user_id)
                        log_user_event(user_id, 'INFO', "✅ تم تسجيل الدخول بنجاح (جلسة محفوظة)")
                        return

                    # إرسال كود التحقق
                    socketio.emit('log_update', {"message": f"📱 إرسال كود التحقق إلى: {phone_number}"}, to=user_id)
                    log_user_event(user_id, 'INFO', f"📱 إرسال كود التحقق إلى: {phone_number}")

                    result = login.send_code(phone_number)
                    if not result["success"]:
                        socketio.emit('log_update', {"message": f"❌ {result['message']}"}, to=user_id)
                        socketio.emit('login_result', {"status": "error", "message": result["message"]}, to=user_id)
                        log_user_event(user_id, 'ERROR', f"❌ فشل إرسال الكود: {result['message']}")
                        return

                    with USERS_LOCK:
                        if user_id in USERS:
                            USERS[user_id]['awaiting_code'] = True
                            USERS[user_id]['awaiting_password'] = False
                            USERS[user_id]['connected'] = True

                    socketio.emit('login_status', {
                        "logged_in": False, "connected": True,
                        "awaiting_code": True, "awaiting_password": False, "is_running": False
                    }, to=user_id)
                    socketio.emit('log_update', {"message": "✅ تم إرسال كود التحقق - تحقق من رسائل تيليجرام"}, to=user_id)
                    socketio.emit('login_result', {"status": "code_required", "message": "📱 تم إرسال كود التحقق"}, to=user_id)
                    log_user_event(user_id, 'INFO', "✅ تم إرسال كود التحقق")

                except Exception as e:
                    error_message = str(e)
                    logger.error(f"BG setup error for {user_id}: {error_message}")
                    log_user_event(user_id, 'ERROR', f"❌ خطأ في الإعداد: {error_message}")
                    if "ResendCodeRequest" in error_message or "all available options" in error_message:
                        msg = "⚠️ يرجى الانتظار قبل طلب كود جديد"
                    else:
                        msg = f"❌ خطأ: {error_message}"
                    socketio.emit('log_update', {"message": msg}, to=user_id)
                    socketio.emit('login_result', {"status": "error", "message": msg}, to=user_id)

            _OSThread(target=_bg_connect, daemon=True).start()
            return {"status": "pending", "message": "🔄 جارِ الاتصال بتيليجرام..."}

        except Exception as e:
            error_message = str(e)
            logger.error(f"Setup error for {user_id}: {error_message}")
            socketio.emit('log_update', {"message": f"❌ خطأ في الإعداد: {error_message}"}, to=user_id)
            return {"status": "error", "message": f"❌ خطأ: {error_message}"}

    def _fetch_account_name(self, user_id):
        try:
            with USERS_LOCK:
                if user_id not in USERS:
                    return None
                client_manager = USERS[user_id].get('client_manager')
            if not client_manager or not client_manager.client:
                return None
            me = client_manager.run_coroutine(client_manager.client.get_me())
            if not me:
                return None
            parts = []
            if getattr(me, 'first_name', None):
                parts.append(me.first_name)
            if getattr(me, 'last_name', None):
                parts.append(me.last_name)
            name = ' '.join(parts).strip()
            if not name:
                name = getattr(me, 'username', None) or 'حساب تليجرام'
            with USERS_LOCK:
                if user_id in USERS:
                    USERS[user_id]['account_name'] = name
                    USERS[user_id]['account_username'] = getattr(me, 'username', None)
                    USERS[user_id]['account_phone'] = getattr(me, 'phone', None)
            try:
                self._fetch_account_photo(user_id, me)
            except Exception as photo_err:
                logger.debug(f"Avatar fetch skipped for {user_id}: {photo_err}")
            return name
        except Exception as e:
            logger.error(f"Error fetching account name for {user_id}: {e}")
            return None

    def _fetch_account_photo(self, user_id, me=None):
        try:
            with USERS_LOCK:
                if user_id not in USERS:
                    return None
                client_manager = USERS[user_id].get('client_manager')
            if not client_manager or not client_manager.client:
                return None
            if me is None:
                me = client_manager.run_coroutine(client_manager.client.get_me())
            if not me:
                return None

            avatars_dir = os.path.join(SESSIONS_DIR, 'avatars')
            os.makedirs(avatars_dir, exist_ok=True)
            target_path = os.path.join(avatars_dir, f"{user_id}.jpg")

            async def _download():
                try:
                    return await client_manager.client.download_profile_photo(me, file=target_path)
                except Exception as e:
                    logger.debug(f"download_profile_photo error: {e}")
                    return None

            saved = client_manager.run_coroutine(_download())
            if saved and os.path.exists(target_path) and os.path.getsize(target_path) > 0:
                with USERS_LOCK:
                    if user_id in USERS:
                        USERS[user_id]['account_avatar'] = f"/api/account_avatar/{user_id}?t={int(time.time())}"
                return target_path
            return None
        except Exception as e:
            logger.debug(f"Error fetching account photo for {user_id}: {e}")
            return None

    def verify_code(self, user_id, code):
        try:
            login = self.login_managers.get(user_id)
            if not login:
                return {"status": "error", "message": "❌ لم يتم بدء جلسة تسجيل الدخول"}

            result = login.verify_code(code)

            if result.get("requires_password"):
                with USERS_LOCK:
                    if user_id in USERS:
                        USERS[user_id]['awaiting_code'] = False
                        USERS[user_id]['awaiting_password'] = True
                socketio.emit('login_status', {
                    "logged_in": False,
                    "connected": True,
                    "awaiting_code": False,
                    "awaiting_password": True,
                    "is_running": False
                }, to=user_id)
                return {"status": "password_required", "message": result["message"]}

            if not result["success"]:
                return {"status": "error", "message": result["message"]}

            # وقف عميل تسجيل الدخول أولاً لتحرير ملف الجلسة قبل تشغيل المدير الرئيسي
            if user_id in self.login_managers:
                try:
                    self.login_managers[user_id].stop()
                except Exception:
                    pass
                self.login_managers.pop(user_id, None)

            user_info = result.get("user", {})
            account_name = user_info.get("full_name") or user_info.get("username") or "حساب تليجرام"

            client_manager = self.get_client_manager(user_id)

            with USERS_LOCK:
                if user_id in USERS:
                    USERS[user_id]['client_manager'] = client_manager
                    USERS[user_id]['connected'] = True
                    USERS[user_id]['authenticated'] = True
                    USERS[user_id]['awaiting_code'] = False
                    USERS[user_id]['awaiting_password'] = False
                    USERS[user_id]['account_name'] = account_name

            socketio.emit('login_status', {
                "logged_in": True,
                "connected": True,
                "awaiting_code": False,
                "awaiting_password": False,
                "is_running": False,
                "account_name": account_name
            }, to=user_id)
            socketio.emit('connection_status', {"status": "connected"}, to=user_id)

            # تشغيل عميل التليجرام الرئيسي في الخلفية — بدون تعطيل استجابة HTTP
            def _start_client_bg_code(cm=client_manager, uid=user_id):
                try:
                    cm.start_client_thread()
                    logger.info(f"✅ تم تشغيل عميل التليجرام في الخلفية لـ {uid}")
                except Exception as bg_err:
                    logger.warning(f"تحذير تشغيل العميل في الخلفية لـ {uid}: {bg_err}")

            _OSThread(target=_start_client_bg_code, daemon=True).start()

            return {"status": "success", "message": "✅ تم التحقق بنجاح", "account_name": account_name}

        except Exception as e:
            logger.error(f"Code verification error: {str(e)}")
            return {"status": "error", "message": f"❌ خطأ: {str(e)}"}

    def verify_password(self, user_id, password):
        try:
            login = self.login_managers.get(user_id)
            if not login:
                return {"status": "error", "message": "❌ لم يتم بدء جلسة تسجيل الدخول"}

            result = login.verify_password(password)

            if not result["success"]:
                return {"status": "error", "message": result["message"]}

            # وقف عميل تسجيل الدخول أولاً لتحرير ملف الجلسة قبل تشغيل المدير الرئيسي
            if user_id in self.login_managers:
                try:
                    self.login_managers[user_id].stop()
                except Exception:
                    pass
                self.login_managers.pop(user_id, None)

            user_info = result.get("user", {})
            account_name = user_info.get("full_name") or user_info.get("username") or "حساب تليجرام"

            client_manager = self.get_client_manager(user_id)

            with USERS_LOCK:
                if user_id in USERS:
                    USERS[user_id]['client_manager'] = client_manager
                    USERS[user_id]['connected'] = True
                    USERS[user_id]['authenticated'] = True
                    USERS[user_id]['awaiting_code'] = False
                    USERS[user_id]['awaiting_password'] = False
                    USERS[user_id]['account_name'] = account_name

            socketio.emit('login_status', {
                'logged_in': True,
                'connected': True,
                'awaiting_code': False,
                'awaiting_password': False,
                'account_name': account_name
            }, to=user_id)
            socketio.emit('connection_status', {"status": "connected"}, to=user_id)

            # تشغيل عميل التليجرام الرئيسي في الخلفية — بدون تعطيل استجابة HTTP
            def _start_client_bg_2fa(cm=client_manager, uid=user_id):
                try:
                    cm.start_client_thread()
                    logger.info(f"✅ تم تشغيل عميل التليجرام (2FA) في الخلفية لـ {uid}")
                except Exception as bg_err:
                    logger.warning(f"تحذير تشغيل العميل (2FA) في الخلفية لـ {uid}: {bg_err}")

            _OSThread(target=_start_client_bg_2fa, daemon=True).start()

            return {"status": "success", "message": "✅ تم التحقق بنجاح", "account_name": account_name}

        except Exception as e:
            logger.error(f"Password verification error: {str(e)}")
            return {"status": "error", "message": f"❌ خطأ: {str(e)}"}

    def _resolve_entity(self, client_manager, entity):
        entity = _clean_group_entry(str(entity))
        if not entity:
            raise Exception("اسم المجموعة فارغ بعد التنظيف")
        try:
            return client_manager.run_coroutine(
                client_manager.client.get_entity(entity)
            )
        except Exception:
            if not entity.startswith('@') and not entity.startswith('https://') and not entity.startswith('http://'):
                return client_manager.run_coroutine(
                    client_manager.client.get_entity('@' + entity)
                )
            raise

    def send_message_async(self, user_id, entity, message):
        try:
            with USERS_LOCK:
                if user_id not in USERS:
                    raise Exception("المستخدم غير موجود - يرجى تسجيل الدخول أولاً")
                client_manager = USERS[user_id].get('client_manager')
                if not client_manager:
                    raise Exception("لم يتم تسجيل الدخول - يرجى تسجيل الدخول في التليجرام أولاً")
                if not client_manager.client:
                    raise Exception("عميل التليجرام غير مُهيأ - يرجى إعادة تسجيل الدخول")

            try:
                is_authorized = client_manager.run_coroutine(
                    client_manager.client.is_user_authorized()
                )
                if not is_authorized:
                    raise Exception("جلسة التليجرام منتهية الصلاحية - يرجى إعادة تسجيل الدخول")
            except Exception as auth_error:
                raise Exception(f"خطأ في التحقق من التصريح: {str(auth_error)}")

            entity_obj = self._resolve_entity(client_manager, entity)

            final_message = self._maybe_sanitize(user_id, client_manager, entity_obj, entity, message)
            if final_message is None:
                return {"success": False, "skipped": True,
                        "message": "تم تخطي الإرسال: الرسالة بعد التنقية أصبحت فارغة"}

            result = client_manager.run_coroutine(
                client_manager.client.send_message(entity_obj, final_message)
            )

            return {"success": True, "message_id": result.id}

        except Exception as e:
            logger.error(f"Send message error: {str(e)}")
            raise Exception(str(e))

    def _check_group_protection(self, user_id, client_manager, entity_obj, entity_label):
        """التحقق من حماية المجموعة وإرجاع الإجراء المناسب"""
        try:
            settings = load_settings(user_id)
            mode = (settings.get('sanitize_mode') or 'off').lower()
            skip_protected = settings.get('skip_protected_groups', False)

            if mode == 'off' and not skip_protected:
                return 'send', None

            try:
                is_prot, reason = client_manager.run_coroutine(
                    client_manager.is_group_protected(entity_obj)
                )
            except Exception as e:
                logger.warning(f"Group protection check error: {e}")
                is_prot, reason = False, None

            if is_prot and skip_protected:
                msg = f"🛡️ تم تخطي المجموعة المحمية: {entity_label}"
                if reason:
                    msg += f" ({reason})"
                socketio.emit('log_update', {"message": msg}, to=user_id)
                self._send_protection_warning(user_id, entity_label, reason)
                return 'skip', reason

            if is_prot and mode in ('smart', 'always'):
                socketio.emit('log_update', {
                    "message": f"🛡️ مجموعة محمية: {entity_label} ({reason or 'بوت حماية'}) — سيتم تنقية/تحويل الروابط"
                }, to=user_id)
                return 'sanitize', reason

            if is_prot and mode == 'skip':
                return 'skip', reason

            if mode == 'transform':
                return 'transform', None

            if mode == 'off':
                return 'send', None

            return 'send', None
        except Exception as e:
            logger.warning(f"_check_group_protection error: {e}")
            return 'send', None

    def _send_protection_warning(self, user_id, group_name, reason):
        """إرسال تحذير للمستخدم عن المجموعة المحمية"""
        try:
            socketio.emit('protection_warning', {
                "group": group_name,
                "reason": reason,
                "timestamp": time.strftime('%H:%M:%S')
            }, to=user_id)
            try:
                with USERS_LOCK:
                    if user_id in USERS:
                        client_manager = USERS[user_id].get('client_manager')
                        if client_manager and client_manager.client:
                            loop = getattr(client_manager, 'loop', None)
                            if loop and loop.is_running():
                                warning_msg = f"""🛡️ **تنبيه: مجموعة محمية**

⚠️ تم اكتشاف أن المجموعة تحتوي على بوتات حماية:
📌 **{group_name}**

📋 **السبب:** {reason or 'يحتوي على بوتات حماية'}

💡 **الإجراء المتخذ:** تم تخطي الإرسال إلى هذه المجموعة حماية لحسابك.

🔧 **لتغيير الإعدادات:**
• تخطي المجموعات المحمية: إيقاف الإرسال إليها تلقائياً
• تنقية الروابط: تحويل روابط واتساب إلى صيغة آمنة"""
                                asyncio.run_coroutine_threadsafe(
                                    client_manager.client.send_message('me', warning_msg, link_preview=False),
                                    loop
                                )
            except Exception:
                pass
        except Exception as e:
            logger.error(f"Failed to send protection warning: {e}")

    def _maybe_sanitize(self, user_id, client_manager, entity_obj, entity_label, message):
        """تنقية الرسالة حسب وضع الحماية"""
        try:
            action, reason = self._check_group_protection(user_id, client_manager, entity_obj, entity_label)

            if action == 'skip':
                return None

            if action == 'send':
                return message

            if action == 'transform':
                transformed = MessageSanitizer.transform_whatsapp_links(message)
                if transformed != message:
                    socketio.emit('log_update', {
                        "message": f"🔄 تم تحويل روابط واتساب في الرسالة إلى {entity_label}"
                    }, to=user_id)
                return transformed

            if not message:
                return message
            cleaned = MessageSanitizer.sanitize(message, mode='clean')
            if cleaned is None:
                socketio.emit('log_update', {
                    "message": f"⚠️ تم تخطي الإرسال إلى {entity_label}: الرسالة إعلانية بالكامل بعد التنقية"
                }, to=user_id)
                return None
            if cleaned != message:
                socketio.emit('log_update', {
                    "message": f"🧹 تنقية الرسالة قبل الإرسال إلى {entity_label}"
                }, to=user_id)
            return cleaned
        except Exception as e:
            logger.warning(f"_maybe_sanitize error: {e}")
            return message

    def send_media_async(self, user_id, entity, image_files):
        try:
            with USERS_LOCK:
                if user_id not in USERS:
                    raise Exception("المستخدم غير موجود")
                client_manager = USERS[user_id].get('client_manager')

            if not client_manager:
                raise Exception("العميل غير متصل")

            is_authorized = client_manager.run_coroutine(
                client_manager.client.is_user_authorized()
            )
            if not is_authorized:
                raise Exception("العميل غير مصرح")

            entity_obj = self._resolve_entity(client_manager, entity)

            action, _reason = self._check_group_protection(user_id, client_manager, entity_obj, entity)
            if action == 'skip':
                return {"success": False, "skipped": True,
                        "message": f"تم تخطي المجموعة المحمية: {entity}"}

            results = []
            paths = [f['path'] for f in image_files if os.path.exists(f.get('path', ''))]
            if not paths:
                raise Exception("لا توجد ملفات صور صالحة")

            if len(paths) == 1:
                result = client_manager.run_coroutine(
                    client_manager.client.send_file(entity_obj, paths[0])
                )
                results.append(result.id)
            else:
                media_result = client_manager.run_coroutine(
                    client_manager.client.send_file(entity_obj, paths)
                )
                if hasattr(media_result, '__iter__'):
                    for r in media_result:
                        results.append(r.id)
                else:
                    results.append(media_result.id)

            return {"success": True, "message_ids": results}

        except Exception as e:
            logger.error(f"Send media error: {str(e)}")
            raise Exception(str(e))

    def send_message_with_media_async(self, user_id, entity, message, image_files):
        try:
            with USERS_LOCK:
                if user_id not in USERS:
                    raise Exception("المستخدم غير موجود")
                client_manager = USERS[user_id].get('client_manager')

            if not client_manager:
                raise Exception("العميل غير متصل")

            is_authorized = client_manager.run_coroutine(
                client_manager.client.is_user_authorized()
            )

            if not is_authorized:
                raise Exception("العميل غير مصرح")

            entity_obj = self._resolve_entity(client_manager, entity)

            if message:
                _cleaned = self._maybe_sanitize(user_id, client_manager, entity_obj, entity, message)
                if _cleaned is None:
                    return {"success": False, "skipped": True,
                            "message": "تم تخطي الإرسال: الرسالة بعد التنقية أصبحت فارغة"}
                message = _cleaned

            results = []

            if image_files and len(image_files) > 0:
                try:
                    image_paths = []
                    for img_file in image_files:
                        if os.path.exists(img_file['path']):
                            image_paths.append(img_file['path'])
                        else:
                            logger.warning(f"Image file not found: {img_file['path']}")

                    if image_paths:
                        if len(image_paths) == 1:
                            media_result = client_manager.run_coroutine(
                                client_manager.client.send_file(
                                    entity_obj, 
                                    image_paths[0],
                                    caption=message if message else "📷"
                                )
                            )
                            results.append(media_result.id)
                            logger.info(f"Successfully sent single image with message to {entity}")
                        else:
                            try:
                                media_result = client_manager.run_coroutine(
                                    client_manager.client.send_file(
                                        entity_obj,
                                        image_paths,
                                        caption=message if message and message.strip() else None
                                    )
                                )
                                if hasattr(media_result, '__iter__'):
                                    for result in media_result:
                                        results.append(result.id)
                                else:
                                    results.append(media_result.id)
                                logger.info(f"Successfully sent {len(image_paths)} images as album to {entity}")
                            except Exception as album_error:
                                logger.warning(f"Failed to send as album, sending individually: {str(album_error)}")
                                for i, img_path in enumerate(image_paths):
                                    try:
                                        cap = (message if message and message.strip() else None) if i == 0 else None
                                        media_result = client_manager.run_coroutine(
                                            client_manager.client.send_file(
                                                entity_obj,
                                                img_path,
                                                caption=cap
                                            )
                                        )
                                        results.append(media_result.id)
                                    except Exception as img_error:
                                        logger.error(f"Error sending individual image {i+1}: {str(img_error)}")
                                        continue
                except Exception as media_error:
                    logger.error(f"Error in media sending process: {str(media_error)}")
                    raise Exception(f"فشل إرسال الصورة: {str(media_error)[:100]}")
            else:
                if message and message.strip():
                    text_result = client_manager.run_coroutine(
                        client_manager.client.send_message(entity_obj, message)
                    )
                    results.append(text_result.id)
                    logger.info(f"Successfully sent text message to {entity}")

            return {"success": True, "message_ids": results}

        except Exception as e:
            logger.error(f"Send message with media error: {str(e)}")
            raise Exception(str(e))

telegram_manager = TelegramManager()

# =========================== 
# نظام المراقبة المحسن مع Event Handlers
# ===========================
def monitoring_worker(user_id):
    logger.info(f"Starting enhanced monitoring worker with event handlers for user {user_id}")

    try:
        with USERS_LOCK:
            if user_id not in USERS:
                logger.error(f"No user data found for {user_id}")
                return

            USERS[user_id]['monitoring_active'] = True
            client_manager = USERS[user_id].get('client_manager')
            settings = USERS[user_id]['settings']

        if not client_manager:
            logger.error(f"No client manager for user {user_id}")
            return

        watch_words = settings.get('watch_words', [])
        send_groups = settings.get('groups', [])

        if hasattr(client_manager, 'update_monitoring_settings'):
            client_manager.update_monitoring_settings(watch_words, send_groups)

        if watch_words:
            socketio.emit('log_update', {
                "message": f"🚀 بدأت المراقبة الشاملة الفورية - {len(watch_words)} كلمة مراقبة في كامل الحساب | الإرسال لـ {len(send_groups)} مجموعة"
            }, to=user_id)
        else:
            socketio.emit('log_update', {
                "message": f"🚀 بدأت المراقبة الشاملة لكامل الرسائل في الحساب | الإرسال لـ {len(send_groups)} مجموعة"
            }, to=user_id)

        _persisted = load_settings(user_id)
        _saved_last_send = _persisted.get('last_scheduled_send', 0)
        if _saved_last_send == 0:
            _saved_last_send = time.time()
        with USERS_LOCK:
            if user_id in USERS:
                USERS[user_id]['last_scheduled_send'] = _saved_last_send

        consecutive_errors = 0
        max_consecutive_errors = 5

        while True:
            with USERS_LOCK:
                if user_id not in USERS or not USERS[user_id].get('is_running', False):
                    logger.info(f"Stopping monitoring for user {user_id} as is_running is False")
                    break

                user_data = USERS[user_id].copy()
                USERS[user_id]['monitoring_active'] = True

            try:
                settings = user_data.get('settings', {})
                send_type = settings.get('send_type', 'manual')
                current_time = time.time()

                if send_type == 'scheduled':
                    interval_seconds = int(settings.get('interval_seconds', 3600))
                    last_send = user_data.get('last_scheduled_send', 0)
                    remaining = interval_seconds - (current_time - last_send)

                    if remaining <= 0:
                        logger.info(f"Executing scheduled send for user {user_id} (interval={interval_seconds}s)")
                        socketio.emit('log_update', {
                            "message": f"📅 حان موعد الإرسال المجدول — جاري الإرسال إلى {len(settings.get('groups', []))} مجموعة..."
                        }, to=user_id)
                        execute_scheduled_messages(user_id, settings)

                        with USERS_LOCK:
                            if user_id in USERS:
                                USERS[user_id]['last_scheduled_send'] = current_time
                        try:
                            _s = load_settings(user_id)
                            _s['last_scheduled_send'] = current_time
                            save_settings(user_id, _s)
                        except Exception as _se:
                            logger.error(f"Failed to persist last_scheduled_send: {_se}")

                        # احسب وقت الإرسال التالي وأبلغ المستخدم
                        next_send_at = time.strftime('%H:%M:%S', time.localtime(current_time + interval_seconds))
                        socketio.emit('log_update', {
                            "message": f"⏰ الإرسال التالي في: {next_send_at} (بعد {interval_seconds // 60} دقيقة)"
                        }, to=user_id)
                    else:
                        logger.debug(f"Scheduled send for {user_id}: {int(remaining)}s remaining")

                consecutive_errors = 0

                # حساب الوقت المتبقي للإرسال القادم
                next_send_remaining = None
                next_send_at_str = None
                if send_type == 'scheduled':
                    interval_seconds = int(settings.get('interval_seconds', 3600))
                    last_send = user_data.get('last_scheduled_send', 0)
                    remaining = interval_seconds - (current_time - last_send)
                    next_send_remaining = max(0, int(remaining))
                    next_send_at_str = time.strftime('%H:%M:%S', time.localtime(current_time + next_send_remaining))

                status_info = {
                    'timestamp': time.strftime('%H:%M:%S'),
                    'status': 'active',
                    'type': 'event_driven_monitoring',
                    'keywords_active': bool(watch_words),
                    'event_handlers': True,
                    'send_type': send_type,
                    'next_send_remaining': next_send_remaining,
                    'next_send_at': next_send_at_str
                }

                socketio.emit('heartbeat', status_info, to=user_id)

            except Exception as e:
                consecutive_errors += 1
                logger.error(f"Monitoring cycle error for {user_id}: {str(e)}")
                socketio.emit('log_update', {
                    "message": f"⚠️ خطأ في المراقبة: {str(e)[:100]}"
                }, to=user_id)

                if consecutive_errors >= max_consecutive_errors:
                    socketio.emit('log_update', {
                        "message": f"❌ تم إيقاف المراقبة بسبب تكرار الأخطاء ({consecutive_errors})"
                    }, to=user_id)
                    with USERS_LOCK:
                        if user_id in USERS:
                            USERS[user_id]['is_running'] = False
                    break

            time.sleep(10)

    except Exception as e:
        logger.error(f"Monitoring worker top-level error for {user_id}: {str(e)}")
    finally:
        with USERS_LOCK:
            if user_id in USERS:
                USERS[user_id]['is_running'] = False
                USERS[user_id]['monitoring_active'] = False
                USERS[user_id]['thread'] = None

        socketio.emit('log_update', {
            "message": "⏹ تم إيقاف نظام المراقبة المحسن"
        }, to=user_id)

        socketio.emit('heartbeat', {
            'timestamp': time.strftime('%H:%M:%S'),
            'status': 'stopped'
        }, to=user_id)

        logger.info(f"Enhanced monitoring worker ended for user {user_id}")

def execute_scheduled_messages(user_id, settings):
    groups = settings.get('groups', [])
    message = settings.get('message', '')

    if not groups or not message:
        return

    try:
        socketio.emit('log_update', {
            "message": f"📅 تنفيذ الإرسال المجدول إلى {len(groups)} مجموعة"
        }, to=user_id)

        successful = 0
        failed = 0

        for i, group in enumerate(groups, 1):
            try:
                result = telegram_manager.send_message_async(user_id, group, message)

                if isinstance(result, dict) and result.get('skipped'):
                    socketio.emit('log_update', {
                        "message": f"⏭️ [{i}/{len(groups)}] تم تخطي: {group} (الرسالة لم تُرسَل)"
                    }, to=user_id)
                    failed += 1
                    with USERS_LOCK:
                        if user_id in USERS:
                            USERS[user_id]['stats']['errors'] += 1
                else:
                    socketio.emit('log_update', {
                        "message": f"✅ [{i}/{len(groups)}] إرسال مجدول نجح إلى: {group}"
                    }, to=user_id)
                    successful += 1
                    with USERS_LOCK:
                        if user_id in USERS:
                            USERS[user_id]['stats']['sent'] += 1

                if i < len(groups):
                    time.sleep(3)

            except Exception as e:
                error_msg = str(e)
                logger.error(f"Scheduled send error to {group}: {error_msg}")

                socketio.emit('log_update', {
                    "message": f"❌ [{i}/{len(groups)}] إرسال مجدول فشل إلى {group}"
                }, to=user_id)

                failed += 1
                with USERS_LOCK:
                    if user_id in USERS:
                        USERS[user_id]['stats']['errors'] += 1

        socketio.emit('log_update', {
            "message": f"📊 انتهى الإرسال المجدول: ✅ {successful} نجح | ❌ {failed} فشل"
        }, to=user_id)

    except Exception as e:
        logger.error(f"Scheduled messages error: {str(e)}")

# =========================== 
# أحداث Socket.IO
# ===========================
@socketio.on('connect')
def handle_connect():
    try:
        if 'user_id' not in session:
            session['user_id'] = "user_1"
            session.permanent = True

        user_id = session['user_id']

        if user_id not in PREDEFINED_USERS:
            user_id = "user_1"
            session['user_id'] = user_id

        join_room(user_id)
        logger.info(f"User {user_id} ({PREDEFINED_USERS[user_id]['name']}) connected via socket")

        emit('connection_confirmed', {
            'status': 'connected',
            'user_id': user_id,
            'user_name': PREDEFINED_USERS[user_id]['name'],
            'timestamp': time.strftime('%H:%M:%S')
        })

        emit('users_list', {
            'current_user': user_id,
            'users': PREDEFINED_USERS
        })

        notify_user_about_background_operations(user_id)

        all_status = get_all_users_operations_status()
        emit('all_users_status', all_status)

    except Exception as e:
        logger.error(f"Connection error: {str(e)}")
        emit('connection_error', {'message': str(e)})

@socketio.on('switch_user')
def handle_switch_user(data):
    try:
        new_user_id = data.get('user_id')

        if not new_user_id or new_user_id not in PREDEFINED_USERS:
            emit('error', {'message': 'مستخدم غير صحيح'})
            return

        # استخدم from_user_id إذا أُرسل من JS (لتجنب تعارض الغرف بعد تبديل HTTP)
        old_user_id = data.get('from_user_id') or session.get('user_id', 'user_1')
        try:
            leave_room(old_user_id)
        except Exception as leave_error:
            logger.warning(f"Error leaving room {old_user_id}: {str(leave_error)}")

        session['user_id'] = new_user_id
        session.permanent = True

        try:
            join_room(new_user_id)
        except Exception as join_error:
            logger.warning(f"Error joining room {new_user_id}: {str(join_error)}")

        logger.info(f"User switched from {old_user_id} to {new_user_id}")

        emit('user_switched', {
            'current_user': new_user_id,
            'user_name': PREDEFINED_USERS[new_user_id]['name'],
            'message': f"تم التبديل إلى {PREDEFINED_USERS[new_user_id]['name']}"
        })

        try:
            with USERS_LOCK:
                if new_user_id in USERS:
                    user_data = USERS[new_user_id]
                    connected = user_data.get('connected', False)
                    authenticated = user_data.get('authenticated', False)
                    awaiting_code = user_data.get('awaiting_code', False)
                    awaiting_password = user_data.get('awaiting_password', False)
                    is_running = user_data.get('is_running', False)

                    emit('connection_status', {
                        "status": "connected" if connected else "disconnected"
                    })

                    emit('login_status', {
                        "logged_in": authenticated,
                        "connected": connected,
                        "awaiting_code": awaiting_code,
                        "awaiting_password": awaiting_password,
                        "is_running": is_running
                    })

                    settings = load_settings(new_user_id)
                    emit('user_settings', settings)
                else:
                    emit('connection_status', {"status": "disconnected"})
                    emit('login_status', {
                        "logged_in": False,
                        "connected": False,
                        "awaiting_code": False,
                        "awaiting_password": False,
                        "is_running": False
                    })
        except Exception as status_error:
            logger.error(f"Error sending user status: {str(status_error)}")

    except Exception as e:
        logger.error(f"Error switching user: {str(e)}")
        emit('error', {'message': f'خطأ في التبديل: {str(e)}'})

@socketio.on('disconnect')
def handle_disconnect(data=None):
    if 'user_id' in session:
        user_id = session['user_id']
        leave_room(user_id)
        logger.info(f"User {user_id} disconnected from socket")

@socketio.on('heartbeat')
def handle_heartbeat(data):
    try:
        user_id = session.get('user_id')
        if user_id:
            emit('heartbeat_response', {
                'timestamp': time.time(),
                'server_time': time.strftime('%H:%M:%S')
            })
    except Exception as e:
        logger.error(f"Heartbeat error: {str(e)}")

# =========================== 
# المسارات الأساسية
# ===========================
@app.route("/")
def index():
    if 'user_id' not in session:
        session['user_id'] = "user_1"
        session.permanent = True
    elif session['user_id'] not in PREDEFINED_USERS:
        session['user_id'] = "user_1"

    user_id = session['user_id']

    settings = load_settings(user_id)
    connection_status = "disconnected"

    with USERS_LOCK:
        if user_id not in USERS:
            USERS[user_id] = {
                'client_manager': None,
                'settings': settings,
                'thread': None,
                'is_running': False,
                'stats': {"sent": 0, "errors": 0},
                'connected': False,
                'authenticated': False,
                'awaiting_code': False,
                'awaiting_password': False,
                'phone_code_hash': None,
                'monitoring_active': False,
                'event_handlers_registered': False,
                'sent_batches': settings.get('sent_batches', []) or []
            }

        user_data = USERS[user_id]
        connected = user_data.get('connected', False)
        connection_status = "connected" if connected else "disconnected"

    app_title = "مركز سرعة انجاز 📚 للخدمات الطلابية والأكاديمية"
    whatsapp_link = "https://wa.me/+966510349663"

    current_user = PREDEFINED_USERS[user_id]

    response = render_template('index.html',
                          settings=settings,
                          connection_status=connection_status,
                          app_title=app_title,
                          whatsapp_link=whatsapp_link,
                          current_user=current_user,
                          predefined_users=PREDEFINED_USERS)

    resp = make_response(response)
    resp.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate, max-age=0'
    resp.headers['Pragma'] = 'no-cache'
    resp.headers['Expires'] = '0'

    return resp

@app.route("/fresh")
def fresh():
    from flask import make_response
    html = """<!DOCTYPE html>
<html dir="rtl" lang="ar">
<head>
    <meta charset="UTF-8">
    <title>🚀 التطبيق يعمل بنجاح!</title>
    <style>
        body { font-family: Arial, sans-serif; text-align: center; padding: 50px; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; }
        .success { font-size: 2em; margin: 20px 0; }
        .message { font-size: 1.2em; margin: 10px 0; }
        .btn { background: #28a745; color: white; padding: 15px 30px; text-decoration: none; border-radius: 8px; font-size: 1.1em; display: inline-block; margin: 10px; }
        .btn:hover { background: #218838; color: white; }
    </style>
</head>
<body>
    <div class="success">✅ التطبيق يعمل بشكل مثالي!</div>
    <div class="message">🎉 مركز سرعة انجاز للخدمات الطلابية والأكاديمية</div>
    <div class="message">📱 نظام مراقبة التليجرام الذكي</div>
    <a href="/" class="btn">🏠 الانتقال للتطبيق الرئيسي</a>
    <script>
        setTimeout(function() {
            window.location.href = '/';
        }, 3000);
    </script>
</body>
</html>"""

    resp = make_response(html)
    resp.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate, max-age=0'
    resp.headers['Pragma'] = 'no-cache'
    resp.headers['Expires'] = '0'
    resp.headers['Content-Type'] = 'text/html; charset=utf-8'

    return resp

@app.route('/static/<path:filename>')
def static_files(filename):
    return app.send_static_file(filename)

@app.route("/manifest.json")
def manifest():
    manifest_data = {
        "id": "/",
        "name": "مركز سرعة انجاز للخدمات الطلابية والأكاديمية",
        "short_name": "سرعة انجاز",
        "description": "نظام متكامل: تليجرام تلقائي، تحليل أكاديمي، عروض PowerPoint، منسّق مستندات",
        "start_url": "/",
        "scope": "/",
        "display": "standalone",
        "display_override": ["standalone", "minimal-ui"],
        "orientation": "portrait",
        "theme_color": "#1e3c78",
        "background_color": "#1e3c78",
        "lang": "ar",
        "dir": "rtl",
        "categories": ["education", "productivity", "utilities"],
        "prefer_related_applications": False,
        "icons": [
            {"src": "/static/icons/icon-72.png",  "sizes": "72x72",   "type": "image/png", "purpose": "any"},
            {"src": "/static/icons/icon-192.png", "sizes": "192x192", "type": "image/png", "purpose": "any"},
            {"src": "/static/icons/icon-192.png", "sizes": "192x192", "type": "image/png", "purpose": "maskable"},
            {"src": "/static/icons/icon-512.png", "sizes": "512x512", "type": "image/png", "purpose": "any"},
            {"src": "/static/icons/icon-512.png", "sizes": "512x512", "type": "image/png", "purpose": "maskable"}
        ],
        "screenshots": [],
        "shortcuts": [
            {"name": "التحليل الأكاديمي", "short_name": "أكاديمي", "description": "فتح منصة التحليل", "url": "/academic"},
            {"name": "لوحة التحكم",       "short_name": "تحكم",    "description": "لوحة التحكم الرئيسية", "url": "/"}
        ]
    }
    resp = app.response_class(json.dumps(manifest_data, ensure_ascii=False, indent=2),
                              mimetype='application/manifest+json')
    resp.headers['Cache-Control'] = 'no-cache'
    return resp

@app.route("/sw.js")
def service_worker():
    sw_js = r"""
const CACHE_NAME = 'sra3a-v5';
const STATIC_ASSETS = [
  '/',
  '/static/icons/icon-192.png',
  '/static/icons/icon-512.png'
];

// ── التثبيت: تخزين الأصول الأساسية ──
self.addEventListener('install', event => {
  self.skipWaiting();
  event.waitUntil(
    caches.open(CACHE_NAME).then(cache =>
      cache.addAll(STATIC_ASSETS).catch(() => {})
    )
  );
});

// ── التفعيل: حذف الكاش القديم ──
self.addEventListener('activate', event => {
  event.waitUntil(
    caches.keys().then(keys =>
      Promise.all(keys.filter(k => k !== CACHE_NAME).map(k => caches.delete(k)))
    ).then(() => self.clients.claim())
  );
});

// ── الجلب: Network-first، cache كـ fallback ──
self.addEventListener('fetch', event => {
  const req = event.request;
  if (req.method !== 'GET') return;

  // تجاهل طلبات API و WebSocket
  const url = new URL(req.url);
  if (url.pathname.startsWith('/api/') ||
      url.pathname.startsWith('/socket.io') ||
      url.pathname.startsWith('/tools/')) return;

  event.respondWith(
    fetch(req)
      .then(resp => {
        if (resp.ok && resp.type === 'basic') {
          const clone = resp.clone();
          caches.open(CACHE_NAME).then(c => c.put(req, clone));
        }
        return resp;
      })
      .catch(() => caches.match(req))
  );
});

// ── رسائل من الصفحة ──
self.addEventListener('message', event => {
  if (event.data === 'SKIP_WAITING') self.skipWaiting();
});
"""
    resp = app.response_class(sw_js, content_type='application/javascript')
    resp.headers['Service-Worker-Allowed'] = '/'
    resp.headers['Cache-Control'] = 'no-cache'
    return resp

# =========================== 
# API Routes
# ===========================

@app.route("/api", methods=["GET", "HEAD"])
def api_health():
    try:
        if request.method == "HEAD":
            return "", 200
        return jsonify({"status": "ok", "timestamp": time.time(), "message": "Server is running"})
    except Exception as e:
        logger.error(f"Error in api health check: {str(e)}")
        if request.method == "HEAD":
            return "", 500
        return jsonify({"status": "error", "message": "Server error"}), 500

@app.route("/api/save_login", methods=["POST"])
def api_save_login():
    data = request.json

    if not data or not data.get('phone'):
        return jsonify({
            "success": False,
            "message": "❌ يرجى إدخال رقم الهاتف"
        })

    new_phone = data.get('phone')

    # ─── تحديد user_id ───────────────────────────────────────────────────────
    # الأولوية: (1) user_id في body الطلب، (2) الجلسة، (3) الافتراضي user_1
    requested_uid = (data.get('user_id') or '').strip()
    if requested_uid and requested_uid in PREDEFINED_USERS:
        session['user_id'] = requested_uid
        session.permanent = True
    elif 'user_id' not in session or session['user_id'] not in PREDEFINED_USERS:
        session['user_id'] = "user_1"
        session.permanent = True
    session.modified = True

    user_id = session['user_id']
    logger.info(f"api_save_login: user_id={user_id}, phone={new_phone}")
    log_user_event(user_id, 'INFO', f"📱 طلب تسجيل دخول للرقم: {new_phone}")

    # تنظيف الجلسة القديمة لنفس الخانة إذا تغيّر الرقم
    current_settings = load_settings(user_id)
    if current_settings.get('phone') and current_settings.get('phone') != new_phone:
        logger.info(f"Phone changed: {current_settings['phone']} → {new_phone} for {user_id}")
        with USERS_LOCK:
            if user_id in USERS:
                if USERS[user_id].get('is_running'):
                    USERS[user_id]['is_running'] = False
                cm = USERS[user_id].get('client_manager')
                if cm:
                    try: cm.stop()
                    except Exception: pass
                del USERS[user_id]
        # حذف ملف الجلسة القديم
        old_session_file = os.path.join(SESSIONS_DIR, f"{user_id}_session.session")
        if os.path.exists(old_session_file):
            try:
                os.remove(old_session_file)
            except Exception as e:
                logger.warning(f"Could not remove old session file: {e}")
        socketio.emit('log_update', {
            "message": f"🔄 تم مسح الجلسة القديمة لـ {PREDEFINED_USERS[user_id]['name']}"
        }, to=user_id)

    settings = {
        'phone': new_phone,
        'password': data.get('password', ''),
        'login_time': time.time()
    }

    if not save_settings(user_id, settings):
        return jsonify({
            "success": False,
            "message": "❌ فشل في حفظ البيانات"
        })

    try:
        socketio.emit('log_update', {
            "message": f"🔄 بدء تسجيل دخول {PREDEFINED_USERS[user_id]['name']}..."
        }, to=user_id)

        with USERS_LOCK:
            # إزالة أي خانة أخرى تستخدم نفس الرقم
            users_to_remove = [
                uid for uid, ud in USERS.items()
                if uid != user_id and ud['settings'].get('phone') == new_phone
            ]
            for old_uid in users_to_remove:
                logger.info(f"Removing duplicate phone session: {old_uid}")
                if USERS[old_uid].get('is_running'):
                    USERS[old_uid]['is_running'] = False
                cm = USERS[old_uid].get('client_manager')
                if cm:
                    try: cm.stop()
                    except Exception: pass
                del USERS[old_uid]

            # إنشاء/تحديث إدخال المستخدم الحالي
            USERS[user_id] = {
                'client_manager': None,
                'settings': settings,
                'thread': None,
                'is_running': False,
                'stats': {"sent": 0, "errors": 0},
                'connected': False,
                'authenticated': False,
                'awaiting_code': False,
                'awaiting_password': False,
                'phone_code_hash': None,
                'monitoring_active': False,
                'event_handlers_registered': False,
                'sent_batches': settings.get('sent_batches', []) or []
            }

        result = telegram_manager.setup_client(user_id, settings['phone'])

        if result["status"] == "pending":
            # الاتصال يعمل في الخلفية - النتيجة ستصل عبر socket.io
            return jsonify({
                "success": True,
                "message": "🔄 جارِ الاتصال بتيليجرام...",
                "pending": True
            })

        elif result["status"] == "success":
            socketio.emit('log_update', {"message": "✅ تم تسجيل الدخول بنجاح"}, to=user_id)
            socketio.emit('connection_status', {"status": "connected"}, to=user_id)
            socketio.emit('login_status', {
                "logged_in": True, "connected": True,
                "awaiting_code": False, "awaiting_password": False, "is_running": False
            }, to=user_id)
            return jsonify({"success": True, "message": "✅ تم تسجيل الدخول"})

        elif result["status"] == "code_required":
            socketio.emit('log_update', {"message": "📱 تم إرسال كود التحقق"}, to=user_id)
            return jsonify({"success": True, "message": "📱 تم إرسال كود التحقق", "code_required": True})

        else:
            error_message = result.get('message', 'خطأ غير معروف')
            socketio.emit('log_update', {"message": f"❌ {error_message}"}, to=user_id)

            return jsonify({
                "success": False, 
                "message": f"❌ {error_message}"
            })

    except Exception as e:
        logger.error(f"Login error for user {user_id}: {str(e)}")
        socketio.emit('log_update', {
            "message": f"❌ خطأ: {str(e)}"
        }, to=user_id)

        return jsonify({
            "success": False, 
            "message": f"❌ خطأ: {str(e)}"
        })

@app.route("/api/verify_code", methods=["POST"])
def api_verify_code():
    if 'user_id' not in session:
        return jsonify({
            "success": False, 
            "message": "❌ الجلسة غير صالحة، يرجى إعادة تحميل الصفحة"
        })

    user_id = session['user_id']
    data = request.json

    if not data:
        return jsonify({
            "success": False, 
            "message": "❌ لم يتم إرسال البيانات"
        })

    code = data.get('code')
    password = data.get('password')

    if not code and not password:
        return jsonify({
            "success": False, 
            "message": "❌ يرجى إدخال الكود أو كلمة المرور"
        })

    try:
        if code:
            result = telegram_manager.verify_code(user_id, code)
        else:
            result = telegram_manager.verify_password(user_id, password)

        if result["status"] == "success":
            account_name = result.get("account_name")
            socketio.emit('log_update', {
                "message": f"✅ تم التحقق بنجاح — أهلاً {account_name}" if account_name else "✅ تم التحقق بنجاح"
            }, to=user_id)

            socketio.emit('connection_status', {
                "status": "connected"
            }, to=user_id)

            return jsonify({
                "success": True,
                "message": f"✅ تم التحقق بنجاح — أهلاً {account_name}" if account_name else "✅ تم التحقق بنجاح",
                "account_name": account_name
            })

        elif result["status"] == "password_required":
            return jsonify({
                "success": True, 
                "message": result["message"], 
                "password_required": True
            })

        else:
            error_message = result.get('message', 'فشل التحقق')
            socketio.emit('log_update', {
                "message": f"❌ {error_message}"
            }, to=user_id)

            return jsonify({
                "success": False, 
                "message": f"❌ {error_message}"
            })

    except Exception as e:
        socketio.emit('log_update', {
            "message": f"❌ خطأ في التحقق: {str(e)}"
        }, to=user_id)

        return jsonify({
            "success": False, 
            "message": f"❌ خطأ: {str(e)}"
        })

@app.route("/api/save_settings", methods=["POST"])
def api_save_settings():
    if 'user_id' not in session:
        return jsonify({
            "success": False, 
            "message": "❌ الجلسة غير صالحة، يرجى إعادة تحميل الصفحة"
        })

    user_id = session['user_id']
    data = request.json

    if not data:
        return jsonify({
            "success": False, 
            "message": "❌ لم يتم إرسال البيانات"
        })

    current_settings = load_settings(user_id)
    current_settings.update({
        'message': data.get('message', ''),
        'groups': dedupe_groups(data.get('groups', '')),
        'interval_seconds': int(data.get('interval_seconds', 3600)),
        'watch_words': [w.strip() for w in data.get('watch_words', '').split('\n') if w.strip()],
        'send_type': data.get('send_type', 'manual'),
        'scheduled_time': data.get('scheduled_time', ''),
        'max_retries': int(data.get('max_retries', 5)),
        'auto_reconnect': data.get('auto_reconnect', False),
        'sanitize_mode': (data.get('sanitize_mode') or 'off').lower()
    })

    if save_settings(user_id, current_settings):
        with USERS_LOCK:
            if user_id in USERS:
                USERS[user_id]['settings'] = current_settings
                client_manager = USERS[user_id].get('client_manager')
                if client_manager and hasattr(client_manager, 'update_monitoring_settings'):
                    client_manager.update_monitoring_settings(
                        current_settings.get('watch_words', []),
                        current_settings.get('groups', [])
                    )

        socketio.emit('log_update', {
            "message": "✅ تم حفظ الإعدادات بنجاح"
        }, to=user_id)

        return jsonify({
            "success": True, 
            "message": "✅ تم حفظ الإعدادات"
        })
    else:
        return jsonify({
            "success": False, 
            "message": "❌ فشل في حفظ الإعدادات"
        })

@app.route("/api/user_logout", methods=["POST"])
def api_user_logout():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({
            "success": False,
            "message": "❌ لا توجد جلسة نشطة"
        })

    try:
        logger.info(f"User {user_id} logging out...")

        with USERS_LOCK:
            if user_id in USERS:
                client_manager = USERS[user_id].get('client_manager')
                if client_manager:
                    try:
                        if USERS[user_id].get('is_running'):
                            USERS[user_id]['is_running'] = False

                        if hasattr(client_manager, 'client') and client_manager.client:
                            client_manager.client.disconnect()
                            logger.info(f"Client disconnected for user {user_id}")

                        if hasattr(client_manager, 'stop'):
                            client_manager.stop()

                    except Exception as e:
                        logger.error(f"خطأ في إغلاق العميل للمستخدم {user_id}: {e}")

                del USERS[user_id]
                logger.info(f"User data removed from memory for {user_id}")

        session_file = os.path.join(SESSIONS_DIR, f"{user_id}_session.session")
        if os.path.exists(session_file):
            try:
                os.remove(session_file)
                logger.info(f"Session file removed for {user_id}")
            except Exception as e:
                logger.error(f"خطأ في حذف ملف الجلسة: {e}")

        # إيقاف النشر الدوري فوراً عند تسجيل الخروج
        try:
            rotating_manager.stop(user_id)
            logger.info(f"Rotating send stopped on logout for {user_id}")
        except Exception as _re:
            logger.error(f"خطأ في إيقاف النشر الدوري عند الخروج: {_re}")

        settings_file = os.path.join(SESSIONS_DIR, f"{user_id}.json")
        if os.path.exists(settings_file):
            try:
                settings = load_settings(user_id)
                settings.update({
                    'phone': '',
                    'authenticated': False,
                    'connected': False,
                    'rotating_persistent': False,
                    'monitoring_persistent': False
                })
                save_settings(user_id, settings)
                logger.info(f"Settings cleared for {user_id}")
            except Exception as e:
                logger.error(f"خطأ في مسح الإعدادات: {e}")

        socketio.emit('log_update', {
            "message": "🚪 تم تسجيل الخروج وإنهاء جلسة التليجرام"
        }, to=user_id)

        socketio.emit('connection_status', {
            "status": "disconnected"
        }, to=user_id)

        socketio.emit('login_status', {
            "logged_in": False,
            "connected": False,
            "awaiting_code": False,
            "awaiting_password": False,
            "is_running": False
        }, to=user_id)

        logger.info(f"User {user_id} logged out successfully")

        return jsonify({
            "success": True,
            "message": "✅ تم تسجيل الخروج وإنهاء جلسة التليجرام بنجاح"
        })

    except Exception as e:
        logger.error(f"خطأ في تسجيل الخروج للمستخدم {user_id}: {str(e)}")
        return jsonify({
            "success": False,
            "message": f"❌ خطأ في تسجيل الخروج: {str(e)}"
        })

@app.route("/api/get_account_info", methods=["GET"])
def api_get_account_info():
    user_id = session.get('user_id', 'user_1')
    try:
        try:
            telegram_manager.ensure_client_active(user_id)
        except Exception as e:
            logger.debug(f"ensure_client_active in get_account_info: {e}")
        with USERS_LOCK:
            udata = USERS.get(user_id, {})
            cached = {
                "account_name": udata.get('account_name'),
                "account_username": udata.get('account_username'),
                "account_phone": udata.get('account_phone'),
                "account_avatar": udata.get('account_avatar'),
                "authenticated": udata.get('authenticated', False)
            }
        if not cached["account_name"] and cached["authenticated"]:
            try:
                cached["account_name"] = telegram_manager._fetch_account_name(user_id)
                with USERS_LOCK:
                    cached["account_username"] = USERS.get(user_id, {}).get('account_username')
                    cached["account_phone"] = USERS.get(user_id, {}).get('account_phone')
                    cached["account_avatar"] = USERS.get(user_id, {}).get('account_avatar')
            except Exception as e:
                logger.error(f"get_account_info refresh failed: {e}")
        if not cached.get("account_avatar"):
            avatar_file = os.path.join(SESSIONS_DIR, 'avatars', f"{user_id}.jpg")
            if os.path.exists(avatar_file) and os.path.getsize(avatar_file) > 0:
                cached["account_avatar"] = f"/api/account_avatar/{user_id}"
        return jsonify({
            "success": True,
            "user_id": user_id,
            "predefined_name": PREDEFINED_USERS.get(user_id, {}).get('name'),
            **cached
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/account_avatar/<uid>", methods=["GET"])
def api_account_avatar(uid):
    try:
        from flask import send_file, abort
        avatar_file = os.path.join(SESSIONS_DIR, 'avatars', f"{uid}.jpg")
        if os.path.exists(avatar_file) and os.path.getsize(avatar_file) > 0:
            return send_file(avatar_file, mimetype='image/jpeg', max_age=60)
        return ('', 404)
    except Exception as e:
        logger.error(f"Avatar serving error for {uid}: {e}")
        return ('', 404)

@app.route("/api/switch_user", methods=["POST"])
def api_switch_user():
    try:
        data = request.get_json()
        new_user_id = data.get('user_id')

        if not new_user_id or new_user_id not in PREDEFINED_USERS:
            return jsonify({
                "success": False,
                "message": "❌ مستخدم غير صحيح"
            })

        old_user_id = session.get('user_id', 'user_1')

        if old_user_id in USERS:
            current_settings = USERS[old_user_id].get('settings', {})
            if current_settings:
                save_settings(old_user_id, current_settings)
                logger.info(f"✅ Settings saved for user {old_user_id} - Operations continue running")

        with USERS_LOCK:
            if new_user_id not in USERS:
                saved_settings = load_settings(new_user_id)

                USERS[new_user_id] = {
                    'client_manager': None,
                    'settings': saved_settings,
                    'thread': None,
                    'is_running': False,
                    'stats': {"sent": 0, "errors": 0},
                    'connected': False,
                    'authenticated': False,
                    'awaiting_code': False,
                    'awaiting_password': False,
                    'phone_code_hash': None,
                    'monitoring_active': False,
                    'event_handlers_registered': False,
                    'sent_batches': (saved_settings or {}).get('sent_batches', []) or []
                }

                session_file = os.path.join(SESSIONS_DIR, f"{new_user_id}_session.session")
                if os.path.exists(session_file) and saved_settings.get('phone'):
                    USERS[new_user_id]['connected'] = True
                    USERS[new_user_id]['authenticated'] = True
                    logger.info(f"Found existing session for user {new_user_id}")
            else:
                saved_settings = load_settings(new_user_id)
                USERS[new_user_id]['settings'].update(saved_settings)

        session['user_id'] = new_user_id
        session.permanent = True

        logger.info(f"✅ User switched from {old_user_id} to {new_user_id} - All operations remain active")

        active_operations_summary = get_all_users_operations_status()

        socketio.emit('user_settings', USERS[new_user_id]['settings'], to=new_user_id)

        account_name = None
        account_avatar = None
        try:
            telegram_manager.ensure_client_active(new_user_id)

            with USERS_LOCK:
                account_name = USERS[new_user_id].get('account_name')
                account_avatar = USERS[new_user_id].get('account_avatar')
            if not account_name and USERS[new_user_id].get('authenticated'):
                account_name = telegram_manager._fetch_account_name(new_user_id)
                with USERS_LOCK:
                    account_avatar = USERS[new_user_id].get('account_avatar')
        except Exception as e:
            logger.error(f"Could not load account name on switch: {e}")

        if not account_avatar:
            avatar_file = os.path.join(SESSIONS_DIR, 'avatars', f"{new_user_id}.jpg")
            if os.path.exists(avatar_file) and os.path.getsize(avatar_file) > 0:
                account_avatar = f"/api/account_avatar/{new_user_id}"

        return jsonify({
            "success": True,
            "message": f"✅ تم التبديل إلى {PREDEFINED_USERS[new_user_id]['name']}" + (f" — حساب تليجرام: {account_name}" if account_name else ""),
            "switched": old_user_id != new_user_id,
            "previous_user_id": old_user_id,
            "user": {
                "id": new_user_id,
                "name": PREDEFINED_USERS[new_user_id]['name'],
                "icon": PREDEFINED_USERS[new_user_id]['icon'],
                "color": PREDEFINED_USERS[new_user_id]['color'],
                "account_name": account_name,
                "account_avatar": account_avatar,
                "authenticated": USERS[new_user_id].get('authenticated', False)
            },
            "account_name": account_name,
            "account_avatar": account_avatar,
            "settings": USERS[new_user_id]['settings'],
            "active_operations": active_operations_summary
        })

    except Exception as e:
        logger.error(f"Error in user switching API: {str(e)}")
        return jsonify({
            "success": False,
            "message": f"❌ خطأ في التبديل: {str(e)}"
        })

@app.route("/api/start_monitoring", methods=["POST"])
def api_start_monitoring():
    if 'user_id' not in session:
        return jsonify({
            "success": False, 
            "message": "❌ الجلسة غير صالحة، يرجى إعادة تحميل الصفحة"
        })

    user_id = session['user_id']

    with USERS_LOCK:
        if user_id not in USERS:
            return jsonify({
                "success": False, 
                "message": "❌ لم يتم إعداد الحساب"
            })

        if not USERS[user_id].get('authenticated'):
            return jsonify({
                "success": False, 
                "message": "❌ يجب تسجيل الدخول أولاً"
            })

        if USERS[user_id]['is_running']:
            return jsonify({
                "success": False, 
                "message": "✅ النظام يعمل بالفعل"
            })

        USERS[user_id]['is_running'] = True

    try:
        _settings = load_settings(user_id)
        _settings['monitoring_persistent'] = True
        save_settings(user_id, _settings)
    except Exception as _e:
        logger.error(f"Failed to persist monitoring flag for {user_id}: {_e}")

    log_user_event(user_id, 'INFO', "🚀 بدء تشغيل نظام المراقبة...")
    socketio.emit('log_update', {
        "message": "🚀 بدء تشغيل نظام المراقبة المحسن مع Event Handlers..."
    }, to=user_id)

    try:
        monitoring_thread = _OSThread(
            target=monitoring_worker, 
            args=(user_id,), 
            daemon=True
        )
        monitoring_thread.start()

        with USERS_LOCK:
            USERS[user_id]['thread'] = monitoring_thread

        socketio.emit('monitoring_status', {
            "monitoring_active": True,
            "status": "running",
            "is_running": True
        }, to=user_id)

        socketio.emit('update_monitoring_buttons', {
            "is_running": True
        }, to=user_id)

        return jsonify({
            "success": True, 
            "message": "🚀 بدأت المراقبة المحسنة مع Event Handlers"
        })

    except Exception as e:
        logger.error(f"Failed to start monitoring for {user_id}: {str(e)}")

        with USERS_LOCK:
            USERS[user_id]['is_running'] = False

        return jsonify({
            "success": False, 
            "message": f"❌ فشل في بدء المراقبة: {str(e)}"
        })

@app.route("/api/stop_monitoring", methods=["POST"])
def api_stop_monitoring():
    if 'user_id' not in session:
        return jsonify({
            "success": False, 
            "message": "❌ الجلسة غير صالحة، يرجى إعادة تحميل الصفحة"
        })

    user_id = session['user_id']

    try:
        _settings = load_settings(user_id)
        _settings['monitoring_persistent'] = False
        save_settings(user_id, _settings)
    except Exception as _e:
        logger.error(f"Failed to clear monitoring flag for {user_id}: {_e}")

    with USERS_LOCK:
        if user_id in USERS and USERS[user_id]['is_running']:
            USERS[user_id]['is_running'] = False
            log_user_event(user_id, 'INFO', "⏹ تم إيقاف نظام المراقبة")
            socketio.emit('log_update', {
                "message": "⏹ إيقاف نظام المراقبة..."
            }, to=user_id)

            socketio.emit('monitoring_status', {
                "monitoring_active": False,
                "status": "stopped",
                "is_running": False
            }, to=user_id)

            socketio.emit('update_monitoring_buttons', {
                "is_running": False
            }, to=user_id)

            return jsonify({
                "success": True, 
                "message": "⏹ تم إيقاف المراقبة"
            })

    return jsonify({
        "success": False, 
        "message": "❌ النظام غير مشغل"
    })

@app.route("/api/send_now", methods=["POST"])
def api_send_now():
    if 'user_id' not in session:
        return jsonify({
            "success": False, 
            "message": "❌ الجلسة غير صالحة، يرجى إعادة تحميل الصفحة"
        })

    user_id = session['user_id']

    with USERS_LOCK:
        if user_id not in USERS:
            return jsonify({
                "success": False, 
                "message": "❌ لم يتم إعداد الحساب"
            })

        if not USERS[user_id].get('authenticated'):
            return jsonify({
                "success": False, 
                "message": "❌ يجب تسجيل الدخول أولاً"
            })

    data = request.get_json()
    if not data:
        return jsonify({
            "success": False, 
            "message": "❌ لا توجد بيانات مرسلة"
        })

    message = data.get('message', '').strip()
    groups = data.get('groups', '').strip()
    images = data.get('images', [])

    if not message and not images:
        return jsonify({
            "success": False, 
            "message": "❌ يجب كتابة رسالة أو رفع صورة للإرسال"
        })

    if not groups:
        return jsonify({
            "success": False, 
            "message": "❌ يجب تحديد المجموعات للإرسال إليها"
        })

    raw_groups = [g.strip() for g in groups.replace('\n', ',').split(',') if g.strip()]
    original_count = len(raw_groups)
    groups_list = dedupe_groups(raw_groups)
    duplicates_removed = original_count - len(groups_list)
    if duplicates_removed > 0:
        socketio.emit('log_update', {
            "message": f"♻️ تم تجاهل {duplicates_removed} رابط مكرر في قائمة الإرسال"
        }, to=user_id)

    if not groups_list:
        return jsonify({
            "success": False, 
            "message": "❌ يجب تحديد مجموعة واحدة على الأقل"
        })

    image_files = []
    if images:
        try:
            for img_data in images:
                raw_data = img_data.get('data', '')
                if ',' in raw_data:
                    base64_data = raw_data.split(',', 1)[1]
                else:
                    base64_data = raw_data
                image_bytes = base64.b64decode(base64_data)

                mime = img_data.get('type', 'image/jpeg')
                ext = mime.split('/')[-1].lower()
                if ext in ('jpeg', 'jpg'):
                    ext = 'jpg'
                elif ext not in ('png', 'gif', 'webp', 'bmp'):
                    ext = 'jpg'

                temp_file = tempfile.NamedTemporaryFile(
                    delete=False, suffix=f'.{ext}', mode='wb'
                )
                temp_file.write(image_bytes)
                temp_file.flush()
                temp_file.close()

                image_files.append({
                    'path': temp_file.name,
                    'name': img_data.get('name', f'image.{ext}'),
                    'type': mime
                })

            socketio.emit('log_update', {
                "message": f"📷 تم تحضير {len(image_files)} صورة للإرسال"
            }, to=user_id)

        except Exception as e:
            logger.error(f"Error processing images: {str(e)}")
            return jsonify({
                "success": False,
                "message": f"❌ خطأ في معالجة الصور: {str(e)}"
            })

    content_type = "رسالة"
    if images and message:
        content_type = f"رسالة مع {len(images)} صورة"
    elif images:
        content_type = f"{len(images)} صورة"

    socketio.emit('log_update', {
        "message": f"🚀 بدء الإرسال الفوري: {content_type} إلى {len(groups_list)} مجموعة"
    }, to=user_id)

    def send_messages_with_images():
        try:
            successful = 0
            failed = 0
            batch_id = str(uuid.uuid4())
            batch_entries = []

            for i, group in enumerate(groups_list, 1):
                try:
                    if images and message:
                        result = telegram_manager.send_message_with_media_async(
                            user_id, group, message, image_files
                        )
                    elif images:
                        result = telegram_manager.send_media_async(
                            user_id, group, image_files
                        )
                    else:
                        result = telegram_manager.send_message_async(user_id, group, message)

                    if isinstance(result, dict) and result.get('skipped'):
                        socketio.emit('log_update', {
                            "message": f"⏭️ [{i}/{len(groups_list)}] تم تخطي المجموعة المحمية: {group}"
                        }, to=user_id)
                    else:
                        socketio.emit('log_update', {
                            "message": f"✅ [{i}/{len(groups_list)}] نجح إلى: {group}"
                        }, to=user_id)
                        successful += 1
                        # حفظ معرف الرسالة لدفعة "رسائلي"
                        msg_id = None
                        if isinstance(result, dict):
                            msg_id = result.get('message_id') or (result.get('message_ids') or [None])[0]
                        if msg_id:
                            batch_entries.append({"group": group, "msg_id": msg_id})
                        with USERS_LOCK:
                            if user_id in USERS:
                                USERS[user_id]['stats']['sent'] += 1
                        with USERS_LOCK:
                            if user_id in USERS:
                                socketio.emit('stats_update', USERS[user_id]['stats'], to=user_id)

                    if i < len(groups_list):
                        time.sleep(3)

                except Exception as e:
                    error_msg = str(e)
                    if "banned" in error_msg.lower() or "ban" in error_msg.lower():
                        error_type = "محظور"
                    elif "flood" in error_msg.lower():
                        error_type = "تجاوز حد الإرسال — انتظر"
                    elif "private" in error_msg.lower():
                        error_type = "خاص/محدود"
                    elif "can't write" in error_msg.lower() or "write" in error_msg.lower():
                        error_type = "غير مسموح بالإرسال"
                    elif "not found" in error_msg.lower() or "invalid" in error_msg.lower():
                        error_type = "مجموعة غير موجودة"
                    else:
                        error_type = error_msg[:60]
                    log_user_event(user_id, 'ERROR', f"❌ فشل الإرسال إلى {group}: {error_type}")

                    logger.error(f"Send error to {group}: {error_msg}")
                    socketio.emit('log_update', {
                        "message": f"❌ [{i}/{len(groups_list)}] فشل إلى {group}: {error_type}"
                    }, to=user_id)

                    failed += 1
                    with USERS_LOCK:
                        if user_id in USERS:
                            USERS[user_id]['stats']['errors'] += 1
                            socketio.emit('stats_update', USERS[user_id]['stats'], to=user_id)

            socketio.emit('log_update', {
                "message": f"📊 انتهى الإرسال: ✅ {successful} نجح | ❌ {failed} فشل"
            }, to=user_id)

            # ── حفظ الدفعة في "رسائلي" ──
            if batch_entries:
                batch_record = {
                    "id": batch_id,
                    "text": message or "",
                    "has_media": bool(images),
                    "sent_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    "sent_count": successful,
                    "entries": batch_entries
                }
                with USERS_LOCK:
                    ud = USERS.get(user_id)
                    if ud is not None:
                        if 'sent_batches' not in ud or ud['sent_batches'] is None:
                            ud['sent_batches'] = []
                        ud['sent_batches'].append(batch_record)
                        if len(ud['sent_batches']) > 100:
                            ud['sent_batches'] = ud['sent_batches'][-100:]
                        # حفظ الدفعات في الإعدادات لضمان الاستمرارية بعد إعادة التشغيل
                        try:
                            settings = load_settings(user_id)
                            settings['sent_batches'] = ud['sent_batches']
                            save_settings(user_id, settings)
                        except Exception:
                            pass
                socketio.emit('batch_saved', batch_record, to=user_id)

        except Exception as e:
            logger.error(f"Send thread error: {str(e)}")
        finally:
            for img_file in image_files:
                try:
                    if os.path.exists(img_file['path']):
                        os.unlink(img_file['path'])
                        logger.info(f"Cleaned up temp file: {img_file['name']}")
                except Exception as e:
                    logger.error(f"Error cleaning temp file {img_file.get('name', 'unknown')}: {str(e)}")

    _OSThread(target=send_messages_with_images, daemon=True).start()

    return jsonify({
        "success": True, 
        "message": f"🚀 بدأ إرسال {content_type} لـ {len(groups_list)} مجموعة"
    })

@app.route("/api/scan_groups_protection", methods=["POST"])
def api_scan_groups_protection():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "غير مسجّل"}), 401
    try:
        with USERS_LOCK:
            client_manager = USERS.get(user_id, {}).get('client_manager')
        if not client_manager:
            return jsonify({"error": "العميل غير متصل"}), 400

        data = request.get_json(force=True, silent=True) or {}
        raw_groups = data.get('groups', '')
        group_list = [g.strip() for g in re.split(r'[\n,]+', raw_groups) if g.strip()]
        if not group_list:
            return jsonify({"error": "لا توجد مجموعات للفحص"}), 400

        results = []
        for g in group_list[:50]:
            try:
                try:
                    entity_obj = client_manager.run_coroutine(
                        client_manager.client.get_entity(g)
                    )
                except Exception:
                    g2 = ('@' + g) if not g.startswith('@') and not g.startswith('https://') else g
                    entity_obj = client_manager.run_coroutine(
                        client_manager.client.get_entity(g2)
                    )
                is_prot, reason = client_manager.run_coroutine(
                    client_manager.is_group_protected(entity_obj)
                )
                title = getattr(entity_obj, 'title', g)
                results.append({
                    "group": g,
                    "title": title,
                    "protected": is_prot,
                    "reason": reason or ('غير محمية ✅' if not is_prot else '')
                })
            except Exception as e:
                results.append({"group": g, "title": g, "protected": False, "reason": f"خطأ: {str(e)[:60]}"})

        protected_count = sum(1 for r in results if r['protected'])
        return jsonify({"success": True, "results": results, "protected_count": protected_count, "total": len(results)})
    except Exception as e:
        logger.error(f"Scan groups error: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/get_stats", methods=["GET"])
def api_get_stats():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"sent": 0, "errors": 0})

    with USERS_LOCK:
        if user_id in USERS:
            return jsonify(USERS[user_id]['stats'])

    return jsonify({"sent": 0, "errors": 0})

@app.route("/api/get_login_status", methods=["GET"])
def api_get_login_status():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"logged_in": False, "connected": False})

    with USERS_LOCK:
        if user_id in USERS:
            user_data = USERS[user_id]
            client_manager = user_data.get('client_manager')
            authenticated = user_data.get('authenticated', False)
            connected = user_data.get('connected', False)

            if not authenticated and 'settings' in user_data and 'phone' in user_data['settings']:
                session_file = os.path.join(SESSIONS_DIR, f"{user_id}_session.session")
                if os.path.exists(session_file):
                    authenticated = True
                    connected = True
                    USERS[user_id]['authenticated'] = True
                    USERS[user_id]['connected'] = True

            return jsonify({
                "logged_in": authenticated, 
                "connected": connected,
                "is_running": user_data.get('is_running', False)
            })

    return jsonify({"logged_in": False, "connected": False, "is_running": False})

@app.route("/api/get_user_info", methods=["GET"])
def api_get_user_info():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"success": False, "message": "غير مسجل دخول"})

    with USERS_LOCK:
        if user_id in USERS and 'settings' in USERS[user_id]:
            settings = USERS[user_id]['settings']
            return jsonify({
                "success": True,
                "phone": settings.get('phone', ''),
                "name": settings.get('name', ''),
                "user_id": user_id[:8] + "..."
            })

    return jsonify({"success": False, "message": "لم يتم العثور على معلومات المستخدم"})

@app.route("/api/resend_code", methods=["POST"])
def api_resend_code():
    try:
        if 'user_id' not in session:
            return jsonify({"success": False, "message": "❌ الجلسة غير صالحة"})
        user_id = session['user_id']
        data = request.json or {}
        force_sms = bool(data.get('force_sms', False))

        with USERS_LOCK:
            if user_id not in USERS:
                return jsonify({"success": False, "message": "❌ يرجى البدء بإدخال رقم الهاتف أولاً"})
            client_manager = USERS[user_id].get('client_manager')
            settings = USERS[user_id].get('settings', {})
            phone = settings.get('phone')

        if not client_manager or not client_manager.client or not phone:
            return jsonify({"success": False, "message": "❌ لم يتم إعداد العميل"})

        sent = client_manager.run_coroutine(
            client_manager.client.send_code_request(phone, force_sms=force_sms)
        )
        with USERS_LOCK:
            if user_id in USERS:
                USERS[user_id]['awaiting_code'] = True
                USERS[user_id]['phone_code_hash'] = sent.phone_code_hash

        msg = "📱 تم إعادة الإرسال عبر SMS" if force_sms else "📱 تم إعادة إرسال الكود"
        socketio.emit('log_update', {"message": msg}, to=user_id)
        return jsonify({"success": True, "message": msg})
    except Exception as e:
        logger.error(f"Resend code error: {str(e)}")
        return jsonify({"success": False, "message": f"❌ {str(e)}"})

@app.route("/api/reset_login", methods=["POST"])
def api_reset_login():
    user_id = session.get('user_id', 'user_1')

    if user_id not in PREDEFINED_USERS:
        return jsonify({
            "success": False,
            "message": "❌ مستخدم غير صحيح"
        })

    try:
        logger.info(f"Resetting login for user {user_id}")

        with USERS_LOCK:
            if user_id in USERS:
                if USERS[user_id].get('is_running', False):
                    USERS[user_id]['is_running'] = False

                client_manager = USERS[user_id].get('client_manager')
                if client_manager:
                    try:
                        if hasattr(client_manager, 'stop'):
                            client_manager.stop()
                        if hasattr(client_manager, 'client') and client_manager.client:
                            client_manager.client.disconnect()
                        logger.info(f"Client stopped and disconnected for user {user_id}")
                    except Exception as e:
                        logger.error(f"Error stopping client for {user_id}: {e}")

                del USERS[user_id]
                logger.info(f"User data removed from memory for {user_id}")

        session_file = os.path.join(SESSIONS_DIR, f"{user_id}_session.session")
        if os.path.exists(session_file):
            try:
                os.remove(session_file)
                logger.info(f"Session file removed for {user_id}")
            except Exception as e:
                logger.error(f"Failed to remove session file for {user_id}: {str(e)}")

        socketio.emit('log_update', {
            "message": f"🔄 تم إعادة تعيين جلسة تسجيل الدخول لـ {PREDEFINED_USERS[user_id]['name']}"
        }, to=user_id)

        socketio.emit('connection_status', {
            "status": "disconnected"
        }, to=user_id)

        socketio.emit('login_status', {
            "logged_in": False,
            "connected": False,
            "awaiting_code": False,
            "awaiting_password": False,
            "is_running": False
        }, to=user_id)

        logger.info(f"Login reset completed for user {user_id}")

        return jsonify({
            "success": True, 
            "message": f"✅ تم إعادة تعيين جلسة {PREDEFINED_USERS[user_id]['name']} بنجاح"
        })

    except Exception as e:
        logger.error(f"Error resetting login for {user_id}: {str(e)}")
        return jsonify({
            "success": False,
            "message": f"❌ خطأ في إعادة التعيين: {str(e)}"
        })

@app.route("/api/system_health", methods=["GET"])
def api_system_health():
    try:
        import psutil

        memory = psutil.virtual_memory()
        disk = psutil.disk_usage('/')
        cpu_percent = psutil.cpu_percent(interval=1)
        network = psutil.net_io_counters()

        health_info = {
            'memory': {
                'total': memory.total,
                'available': memory.available,
                'percent': memory.percent,
                'used': memory.used
            },
            'disk': {
                'total': disk.total,
                'used': disk.used,
                'free': disk.free,
                'percent': (disk.used / disk.total) * 100
            },
            'cpu': {
                'percent': cpu_percent,
                'count': psutil.cpu_count()
            },
            'network': {
                'bytes_sent': network.bytes_sent,
                'bytes_recv': network.bytes_recv
            },
            'timestamp': time.time()
        }

        return jsonify({
            "success": True,
            "health": health_info
        })

    except Exception as e:
        return jsonify({
            "success": False,
            "message": f"خطأ: {str(e)}"
        })

def extract_telegram_links(text):
    if not text:
        return []

    patterns = [
        r'https?://t\.me/([a-zA-Z0-9_]+)(?:/\d+)?',
        r'https?://telegram\.me/([a-zA-Z0-9_]+)(?:/\d+)?',
        r'https?://t\.me/\+([a-zA-Z0-9_\-]+)',
        r'https?://telegram\.me/\+([a-zA-Z0-9_\-]+)',
        r't\.me/([a-zA-Z0-9_]+)',
        r't\.me/\+([a-zA-Z0-9_\-]+)',
        r'telegram\.me/([a-zA-Z0-9_]+)',
        r'@([a-zA-Z0-9_]{5,})',
    ]

    found_links = set()

    for pattern in patterns:
        matches = re.findall(pattern, text, re.IGNORECASE)
        for match in matches:
            clean_match = match if isinstance(match, str) else match[0] if match else ''

            if pattern.startswith(r'@'):
                clean_link = f"https://t.me/{clean_match}"
            elif '+' in clean_match or pattern.find(r'\+') != -1:
                clean_link = f"https://t.me/+{clean_match.replace('+', '')}"
            elif clean_match and not clean_match.startswith('http'):
                clean_link = f"https://t.me/{clean_match}"
            elif clean_match.startswith('http'):
                clean_link = f"https://t.me/{clean_match.split('/')[-1]}"
            else:
                clean_link = clean_match

            if clean_link and len(clean_link) > 15:
                clean_link = clean_link.split('?')[0].split('#')[0]
                found_links.add(clean_link)

    links_list = sorted(list(found_links))
    result_links = []
    for link in links_list:
        username = link.split('/')[-1].replace('@', '')
        result_links.append({
            'url': link,
            'username': username,
            'type': 'invite' if '+' in link else 'channel'
        })

    return result_links

async def join_telegram_group(client, group_link, user_id=None, client_manager=None):
    try:
        if group_link.startswith('https://t.me/'):
            group_identifier = group_link.replace('https://t.me/', '')
        elif group_link.startswith('https://telegram.me/'):
            group_identifier = group_link.replace('https://telegram.me/', '')
        elif group_link.startswith('@'):
            group_identifier = group_link[1:]
        else:
            group_identifier = group_link

        try:
            entity = await client.get_entity(group_identifier)
            if hasattr(entity, 'megagroup') or hasattr(entity, 'broadcast'):
                result = await client(functions.channels.JoinChannelRequest(entity))
            else:
                raise Exception("مجموعة عادية - يجب استخدام رابط دعوة")

            return {
                "success": True,
                "already_joined": False,
                "message": "تم الانضمام بنجاح"
            }

        except UserAlreadyParticipantError:
            return {
                "success": True,
                "already_joined": True,
                "message": "منضم مسبقاً للمجموعة"
            }

        except FloodWaitError as e:
            return {
                "success": False,
                "message": f"يرجى الانتظار {e.seconds} ثانية"
            }

        except InviteHashExpiredError:
            return {
                "success": False,
                "message": "انتهت صلاحية رابط الدعوة"
            }

        except InviteHashInvalidError:
            return {
                "success": False,
                "message": "رابط الدعوة غير صحيح"
            }

        except Exception as group_error:
            error_str = str(group_error).lower()
            appeal_url = None
            appeal_note = ""

            if "cas" in error_str or "combot" in error_str:
                appeal_url = "https://cas.chat/appeal"
                appeal_note = "تم حظرك بواسطة CAS (Combot Anti-Spam). توجه إلى الرابط أعلاه لتقديم استئناف."
            elif "spamwatch" in error_str:
                appeal_url = "https://spamwat.ch/appeal"
                appeal_note = "تم حظرك بواسطة SpamWatch. استخدم الرابط أعلاه للاستئناف."
            elif "shieldy" in error_str:
                appeal_url = "https://t.me/Shieldy_Bot?start=appeal"
                appeal_note = "تم حظرك بواسطة Shieldy. افتح البوت في الخاص لطلب فك الحظر."
            elif "rose" in error_str or "missrose" in error_str:
                appeal_url = "https://t.me/MissRose_Bot?start=appeal"
                appeal_note = "تم حظرك بواسطة Rose. أرسل /start إلى البوت ثم اتبع التعليمات."
            elif "groupguard" in error_str:
                appeal_url = "https://t.me/GroupGuardBot?start=appeal"
                appeal_note = "تم حظرك بواسطة GroupGuard. اتصل بالبوت."
            elif "antispam" in error_str or "spam" in error_str:
                appeal_url = "https://t.me/SpamBot"
                appeal_note = "قد يكون حسابك مصنفاً كسبام. تواصل مع @SpamBot للتحقق."
            else:
                if "banned" in error_str or "blocked" in error_str or "forbidden" in error_str:
                    appeal_url = "https://t.me/SpamBot"
                    appeal_note = "حسابك ربما محظور من الانضمام. جرب التواصل مع @SpamBot أو مشرف المجموعة."

            if appeal_url and user_id and client_manager:
                message_text = f"""🚫 **فشل الانضمام إلى المجموعة** 🚫

**الرابط:** {group_link}
**السبب:** {error_str[:200]}

**إجراء مقترح للاستئناف:**
{appeal_note}
🔗 **رابط الاستئناف:** {appeal_url}

يرجى فتح الرابط ومتابعة التعليمات لرفع الحظر. بعد إلغاء الحظر، يمكنك إعادة المحاولة.
"""
                try:
                    await client_manager.send_to_saved_messages(message_text)
                except Exception as save_err:
                    logger.error(f"Could not send appeal to saved messages: {save_err}")

            try:
                if '/' in group_identifier:
                    result = await client(functions.messages.ImportChatInviteRequest(group_identifier.split('/')[-1]))
                    return {
                        "success": True,
                        "already_joined": False,
                        "message": "تم الانضمام عبر رابط الدعوة"
                    }
                else:
                    raise group_error
            except UserAlreadyParticipantError:
                return {
                    "success": True,
                    "already_joined": True,
                    "message": "منضم مسبقاً للمجموعة"
                }
            except Exception as final_error:
                return {
                    "success": False,
                    "message": f"فشل الانضمام: {str(final_error)}",
                    "appeal_url": appeal_url
                }

    except Exception as e:
        return {
            "success": False,
            "message": f"خطأ: {str(e)}"
        }

@app.route("/api/extract_group_links", methods=["POST"])
def api_extract_group_links():
    try:
        data = request.json
        if not data or not data.get('text'):
            return jsonify({
                "success": False,
                "message": "❌ لم يتم إرسال النص"
            })

        text = data.get('text', '')
        links = extract_telegram_links(text)

        return jsonify({
            "success": True,
            "links": links,
            "count": len(links),
            "message": f"✅ تم استخراج {len(links)} رابط"
        })

    except Exception as e:
        logger.error(f"Error extracting links: {str(e)}")
        return jsonify({
            "success": False,
            "message": f"❌ خطأ: {str(e)}"
        })

@app.route("/api/join_group", methods=["POST"])
def api_join_group():
    try:
        user_id = session.get('user_id', 'user_1')

        if user_id not in PREDEFINED_USERS:
            return jsonify({
                "success": False,
                "message": "❌ مستخدم غير صحيح"
            })

        data = request.json

        if not data or not data.get('group_link'):
            return jsonify({
                "success": False,
                "message": "❌ لم يتم إرسال رابط المجموعة"
            })

        group_link_raw = data.get('group_link', '')
        if isinstance(group_link_raw, dict):
            group_link = group_link_raw.get('url', '') or group_link_raw.get('link', '') or str(group_link_raw)
        else:
            group_link = str(group_link_raw)

        group_link = group_link.strip()

        with USERS_LOCK:
            if user_id not in USERS:
                return jsonify({
                    "success": False,
                    "message": f"❌ المستخدم {PREDEFINED_USERS[user_id]['name']} غير مسجل"
                })

            client_manager = USERS[user_id].get('client_manager')
            if not client_manager or not client_manager.client:
                return jsonify({
                    "success": False,
                    "message": "❌ يرجى تسجيل الدخول أولاً"
                })

        result = client_manager.run_coroutine(
            join_telegram_group(client_manager.client, group_link, user_id, client_manager)
        )

        socketio.emit('log_update', {
            "message": f"{'✅' if result['success'] else '❌'} {group_link}: {result['message']}"
        }, to=user_id)

        return jsonify(result)

    except Exception as e:
        logger.error(f"Error joining group: {str(e)}")
        return jsonify({
            "success": False,
            "message": f"❌ خطأ: {str(e)}"
        })

@app.route("/api/start_auto_join", methods=["POST"])
def api_start_auto_join():
    try:
        user_id = session.get('user_id', 'user_1')

        if user_id not in PREDEFINED_USERS:
            return jsonify({
                "success": False,
                "message": "❌ مستخدم غير صحيح"
            })

        data = request.json
        if not data or not data.get('links'):
            return jsonify({
                "success": False,
                "message": "❌ لم يتم إرسال روابط المجموعات"
            })

        links = data.get('links', [])
        delay = data.get('delay', 3)

        if not links:
            return jsonify({
                "success": False,
                "message": "❌ لا توجد روابط للانضمام إليها"
            })

        with USERS_LOCK:
            if user_id not in USERS:
                return jsonify({
                    "success": False,
                    "message": f"❌ المستخدم {PREDEFINED_USERS[user_id]['name']} غير مسجل"
                })

            client_manager = USERS[user_id].get('client_manager')
            if not client_manager or not client_manager.client:
                return jsonify({
                    "success": False,
                    "message": "❌ يرجى تسجيل الدخول أولاً"
                })

        import threading

        def auto_join_worker():
            success_count = 0
            fail_count = 0
            already_joined_count = 0

            socketio.emit('log_update', {
                "message": f"🚀 بدء الانضمام التلقائي لـ {len(links)} مجموعة..."
            }, to=user_id)

            for i, link_obj in enumerate(links):
                try:
                    if isinstance(link_obj, dict):
                        group_link = link_obj.get('url', '') or link_obj.get('link', '') or str(link_obj)
                    else:
                        group_link = str(link_obj)

                    group_link = group_link.strip()

                    socketio.emit('join_progress', {
                        'current': i + 1,
                        'total': len(links),
                        'link': group_link
                    }, to=user_id)

                    result = client_manager.run_coroutine(
                        join_telegram_group(client_manager.client, group_link, user_id, client_manager)
                    )

                    if result['success']:
                        if result.get('already_joined', False):
                            already_joined_count += 1
                            socketio.emit('log_update', {
                                "message": f"ℹ️ منضم مسبقاً: {group_link}"
                            }, to=user_id)
                        else:
                            success_count += 1
                            socketio.emit('log_update', {
                                "message": f"✅ تم الانضمام: {group_link}"
                            }, to=user_id)
                    else:
                        fail_count += 1
                        socketio.emit('log_update', {
                            "message": f"❌ فشل: {group_link} - {result['message']}"
                        }, to=user_id)

                    socketio.emit('join_stats', {
                        'success': success_count,
                        'fail': fail_count,
                        'already_joined': already_joined_count
                    }, to=user_id)

                    if i < len(links) - 1:
                        time.sleep(delay)

                except Exception as e:
                    fail_count += 1
                    socketio.emit('log_update', {
                        "message": f"❌ خطأ في {group_link}: {str(e)}"
                    }, to=user_id)

            socketio.emit('auto_join_completed', {
                'success': success_count,
                'fail': fail_count,
                'already_joined': already_joined_count,
                'total': len(links)
            }, to=user_id)

            socketio.emit('log_update', {
                "message": f"🎉 انتهى الانضمام التلقائي! النجح: {success_count}, فشل: {fail_count}, منضم مسبقاً: {already_joined_count}"
            }, to=user_id)

        thread = _OSThread(target=auto_join_worker, daemon=True)
        thread.start()

        return jsonify({
            "success": True,
            "message": f"✅ تم بدء الانضمام التلقائي لـ {len(links)} مجموعة",
            "total_links": len(links)
        })

    except Exception as e:
        logger.error(f"Error starting auto join: {str(e)}")
        return jsonify({
            "success": False,
            "message": f"❌ خطأ في بدء الانضمام التلقائي: {str(e)}"
        })

# ==========================
# الإضافات الجديدة (الإرسال المتسلسل، الانضمام المتقدم، البوت التعليمي، الردود التلقائية، البحث)
# ==========================

class RotatingSendManager:
    def __init__(self):
        self.threads = {}
        self.stop_events = {}
        self.next_send_at = {}
        self.interval_seconds = {}

    def start(self, user_id, groups, messages, interval_minutes, callback=None):
        if user_id in self.threads and self.threads[user_id] and self.threads[user_id].is_alive():
            self.stop(user_id)

        stop_event = threading.Event()
        self.stop_events[user_id] = stop_event

        thread = _OSThread(target=self._worker, args=(user_id, groups, messages, interval_minutes, stop_event, callback), daemon=True)
        self.threads[user_id] = thread
        thread.start()
        return True

    def stop(self, user_id):
        if user_id in self.stop_events:
            self.stop_events[user_id].set()
        if user_id in self.threads and self.threads[user_id]:
            self.threads[user_id].join(timeout=2)
        self.next_send_at.pop(user_id, None)
        self.interval_seconds.pop(user_id, None)
        return True

    def _worker(self, user_id, groups, messages, interval_minutes, stop_event, callback):
        messages = [m.strip() for m in messages if m and m.strip()]
        if not messages:
            return

        index = 0
        sleep_seconds = max(60, int(interval_minutes * 60))
        self.interval_seconds[user_id] = sleep_seconds

        while not stop_event.is_set():
            try:
                current_msg = messages[index % len(messages)]
                for group in groups:
                    if stop_event.is_set():
                        break
                    try:
                        telegram_manager.send_message_async(user_id, group, current_msg)
                        if callback:
                            callback(user_id, 'success', group, current_msg)
                    except Exception as e:
                        if callback:
                            callback(user_id, 'error', group, str(e))
                    time.sleep(2)
                index += 1
                self.next_send_at[user_id] = time.time() + sleep_seconds
                for _ in range(sleep_seconds):
                    if stop_event.is_set():
                        break
                    time.sleep(1)
            except Exception as e:
                logger.error(f"Rotating send error for {user_id}: {str(e)}")
                time.sleep(10)

rotating_manager = RotatingSendManager()

@app.route("/api/rotating/save", methods=["POST"])
def api_rotating_save():
    try:
        user_id = session.get('user_id', 'user_1')
        data = request.json
        messages = data.get('messages', [''] * 5)
        groups = data.get('groups', [])
        interval = int(data.get('interval', 5))

        settings = load_settings(user_id)
        settings['rotating_messages'] = messages
        settings['rotating_groups'] = dedupe_groups(groups)
        settings['rotating_interval'] = interval
        save_settings(user_id, settings)

        return jsonify({"success": True, "message": "تم حفظ إعدادات الإرسال المتسلسل"})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/rotating/start", methods=["POST"])
def api_rotating_start():
    try:
        user_id = session.get('user_id', 'user_1')
        settings = load_settings(user_id)
        messages = settings.get('rotating_messages', [])
        groups = dedupe_groups(settings.get('rotating_groups', []))
        interval = settings.get('rotating_interval', 5)

        if not groups:
            return jsonify({"success": False, "message": "لا توجد مجموعات محددة"})
        valid_messages = [m for m in messages if m and m.strip()]
        if not valid_messages:
            return jsonify({"success": False, "message": "لا توجد رسائل صالحة"})

        def callback(uid, status, group, info):
            if status == 'success':
                socketio.emit('log_update', {"message": f"🔄 [متسلسل] أرسل إلى {group}"}, to=uid)
            else:
                socketio.emit('log_update', {"message": f"❌ [متسلسل] فشل إلى {group}: {info}"}, to=uid)

        rotating_manager.start(user_id, groups, valid_messages, interval, callback)

        try:
            settings['rotating_persistent'] = True
            save_settings(user_id, settings)
        except Exception as _e:
            logger.error(f"Failed to persist rotating flag for {user_id}: {_e}")

        socketio.emit('log_update', {"message": f"🔄 بدأ الإرسال المتسلسل ({len(valid_messages)} رسائل) كل {interval} دقيقة"}, to=user_id)
        return jsonify({"success": True, "message": "تم بدء الإرسال المتسلسل"})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/rotating/stop", methods=["POST"])
def api_rotating_stop():
    try:
        user_id = session.get('user_id', 'user_1')
        rotating_manager.stop(user_id)

        try:
            _settings = load_settings(user_id)
            _settings['rotating_persistent'] = False
            save_settings(user_id, _settings)
        except Exception as _e:
            logger.error(f"Failed to clear rotating flag for {user_id}: {_e}")

        socketio.emit('log_update', {"message": "⏹ تم إيقاف الإرسال المتسلسل"}, to=user_id)
        return jsonify({"success": True, "message": "تم إيقاف الإرسال المتسلسل"})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/rotating/status", methods=["GET"])
def api_rotating_status():
    try:
        user_id = session.get('user_id', 'user_1')
        settings = load_settings(user_id)
        is_active = user_id in rotating_manager.threads and rotating_manager.threads[user_id] and rotating_manager.threads[user_id].is_alive()
        next_send_in = None
        next_send_at = rotating_manager.next_send_at.get(user_id)
        if is_active and next_send_at:
            remaining = int(next_send_at - time.time())
            next_send_in = max(0, remaining)
        return jsonify({
            "success": True,
            "active": is_active,
            "messages": settings.get('rotating_messages', []),
            "groups": settings.get('rotating_groups', []),
            "interval": settings.get('rotating_interval', 5),
            "next_send_in": next_send_in,
            "interval_seconds": rotating_manager.interval_seconds.get(user_id)
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

def _classify_join_error(msg):
    if not msg:
        return ("خطأ غير معروف", "❓")
    s = str(msg).lower()
    if "anti-spam" in s or "antispam" in s or "spam" in s or "spambot" in s:
        return ("الحساب موسوم كسبام (يحتاج استئناف عبر @SpamBot)", "🚫")
    if "banned" in s or "blocked" in s or "forbidden" in s:
        return ("الحساب محظور من المجموعة أو من الانضمام", "⛔")
    if "expired" in s or "انتهت" in s:
        return ("انتهت صلاحية رابط الدعوة", "⏰")
    if "invalid" in s or "غير صحيح" in s or "غير صالح" in s:
        return ("رابط غير صالح", "🔗")
    if "flood" in s or "wait" in s or "ثانية" in s:
        return ("حد التليجرام مؤقت — يجب الانتظار قبل المحاولة مجدداً", "⏳")
    if "channel_private" in s or "private" in s or "خاص" in s:
        return ("القناة/المجموعة خاصة وتحتاج رابط دعوة", "🔒")
    if "not found" in s or "no user" in s or "could not find" in s or "غير موجود" in s:
        return ("المجموعة غير موجودة أو الرابط خاطئ", "🔍")
    if "too many channels" in s or "channels_too_much" in s:
        return ("الحساب وصل الحد الأقصى من القنوات (500)", "📛")
    if "user_deactivated" in s or "deactivated" in s:
        return ("الحساب معطل من تيليجرام", "🛑")
    if "captcha" in s or "verification" in s:
        return ("المجموعة تتطلب تحقق يدوي (كابتشا)", "🤖")
    if "admin" in s and "approval" in s:
        return ("الانضمام بحاجة موافقة المشرف", "👮")
    if "request" in s and ("send" in s or "join" in s):
        return ("تم إرسال طلب انضمام — بانتظار الموافقة", "📨")
    short = str(msg).strip()
    if len(short) > 120:
        short = short[:120] + "…"
    return (short, "❌")

@app.route("/api/auto_join/advanced", methods=["POST"])
def api_auto_join_advanced():
    try:
        user_id = session.get('user_id', 'user_1')
        data = request.json
        raw_links = data.get('links', [])
        delay = max(1, int(data.get('delay', 3)))
        max_retries = max(1, int(data.get('max_retries', 1)))

        norm_links = []
        for link in raw_links:
            url = link.get('url', link) if isinstance(link, dict) else link
            if url and isinstance(url, str) and url.strip():
                norm_links.append(url.strip())
        norm_links = dedupe_groups(norm_links)

        if not norm_links:
            return jsonify({"success": False, "message": "لا توجد روابط صالحة"})

        with USERS_LOCK:
            if user_id not in USERS:
                return jsonify({"success": False, "message": "المستخدم غير موجود"})
            client_manager = USERS[user_id].get('client_manager')
            if not client_manager or not client_manager.client:
                return jsonify({"success": False, "message": "العميل غير متصل، يرجى تسجيل الدخول"})

        # احفظ الروابط والتأخير لهذا الحساب
        try:
            _s = load_settings(user_id)
            _s['auto_join_links'] = norm_links
            _s['auto_join_delay'] = delay
            save_settings(user_id, _s)
        except Exception:
            pass

        # أنشئ حدث إيقاف لهذا الحساب
        import threading as _threading
        stop_event = _threading.Event()
        with USERS_LOCK:
            if user_id in USERS:
                USERS[user_id]['auto_join_stop'] = stop_event

        def _save_join_state(state_dict):
            """يحفظ حالة الانضمام في USERS لاستعادتها عند العودة للحساب"""
            with USERS_LOCK:
                if user_id in USERS:
                    USERS[user_id]['auto_join_state'] = state_dict

        def advanced_join_worker():
            total = len(norm_links)
            results = {"success": 0, "fail": 0, "already": 0, "total": total, "items": []}

            _save_join_state({
                'running': True, 'total': total, 'done': 0,
                'success': 0, 'already': 0, 'fail': 0,
                'items': [], 'links': norm_links, 'delay': delay
            })

            socketio.emit('auto_join_started', {"total": total}, to=user_id)
            socketio.emit('log_update', {
                "message": f"🚀 بدء الانضمام التلقائي لـ {total} مجموعة (بعد إزالة المكررات)"
            }, to=user_id)

            for idx, url in enumerate(norm_links, 1):
                # تحقق من طلب الإيقاف
                if stop_event.is_set():
                    socketio.emit('log_update', {
                        "message": f"⏹ تم إيقاف الانضمام بعد {idx - 1} مجموعة"
                    }, to=user_id)
                    socketio.emit('auto_join_stopped', {
                        "stopped_at": idx - 1, "total": total,
                        "success": results['success'], "fail": results['fail'],
                        "already": results['already']
                    }, to=user_id)
                    break

                item = {
                    "idx": idx, "total": total, "url": url,
                    "status": "processing", "reason": "", "icon": "⏳",
                    "group_title": ""
                }
                socketio.emit('auto_join_progress', dict(item), to=user_id)

                last_error = None
                final = None
                for attempt in range(max_retries):
                    try:
                        result = client_manager.run_coroutine(
                            join_telegram_group(client_manager.client, url, user_id, client_manager)
                        )
                        if result.get('success'):
                            final = result
                            break
                        else:
                            last_error = result.get('message') or 'فشل غير محدد'
                    except Exception as e:
                        last_error = str(e)
                    if attempt < max_retries - 1:
                        time.sleep(delay)

                if final and final.get('success'):
                    if final.get('already_joined'):
                        item['status'] = 'already'
                        item['icon'] = '📌'
                        item['reason'] = 'منضم مسبقاً'
                        results['already'] += 1
                    else:
                        item['status'] = 'success'
                        item['icon'] = '✅'
                        item['reason'] = final.get('message', 'تم الانضمام بنجاح')
                        results['success'] += 1
                    sleep_time = delay          # انتظار كامل بعد النجاح
                else:
                    reason_text, icon = _classify_join_error(last_error)
                    item['status'] = 'failed'
                    item['icon'] = icon
                    item['reason'] = reason_text
                    item['raw_error'] = (str(last_error)[:200] if last_error else '')
                    results['fail'] += 1
                    sleep_time = delay / 2      # نصف الانتظار بعد الفشل

                results['items'].append(item)

                counts_now = {
                    "success": results['success'],
                    "already": results['already'],
                    "fail": results['fail'],
                    "done": idx,
                    "total": total
                }
                # حفظ الحالة الجارية حتى يمكن استعادتها عند العودة للحساب
                _save_join_state({
                    'running': True, 'total': total, 'done': idx,
                    'success': results['success'], 'already': results['already'],
                    'fail': results['fail'],
                    'items': results['items'][-100:],
                    'links': norm_links, 'delay': delay
                })

                socketio.emit('auto_join_progress', {
                    **item,
                    "counts": counts_now
                }, to=user_id)

                socketio.emit('log_update', {
                    "message": f"{item['icon']} [{idx}/{total}] {url} — {item['reason']}"
                }, to=user_id)

                if idx < total and not stop_event.is_set():
                    socketio.emit('log_update', {
                        "message": f"⏳ انتظار {sleep_time:.0f}ث قبل الرابط التالي…"
                    }, to=user_id)
                    # نوم بخطوات صغيرة حتى نستجيب للإيقاف فوراً
                    elapsed = 0.0
                    step = 0.5
                    while elapsed < sleep_time and not stop_event.is_set():
                        time.sleep(min(step, sleep_time - elapsed))
                        elapsed += step
            else:
                # اكتملت الحلقة دون إيقاف — أرسل التقرير الختامي
                fail_breakdown = {}
                for it in results['items']:
                    if it.get('status') == 'failed':
                        reason = it.get('reason') or 'سبب غير معروف'
                        icon = it.get('icon') or '❌'
                        key = f"{icon} {reason}"
                        if key not in fail_breakdown:
                            fail_breakdown[key] = {"count": 0, "reason": reason, "icon": icon, "links": []}
                        fail_breakdown[key]["count"] += 1
                        fail_breakdown[key]["links"].append(it.get('url', ''))
                results['fail_breakdown'] = sorted(
                    fail_breakdown.values(),
                    key=lambda x: x['count'],
                    reverse=True
                )

                socketio.emit('auto_join_completed', results, to=user_id)

                socketio.emit('log_update', {
                    "message": (
                        f"🎉 انتهى الانضمام: ✅ {results['success']} نجح | "
                        f"📌 {results['already']} منضم مسبقاً | "
                        f"❌ {results['fail']} فشل (المجموع: {total})"
                    )
                }, to=user_id)
                if results.get('fail_breakdown'):
                    lines = ["📊 ملخّص أسباب الفشل:"]
                    for entry in results['fail_breakdown']:
                        lines.append(f"  {entry['icon']} {entry['reason']} — {entry['count']} مجموعة")
                    socketio.emit('log_update', {"message": "\n".join(lines)}, to=user_id)

            # تنظيف حدث الإيقاف وتحديث الحالة لغير جارٍ
            with USERS_LOCK:
                if user_id in USERS:
                    USERS[user_id].pop('auto_join_stop', None)
                    st = USERS[user_id].get('auto_join_state', {})
                    st['running'] = False
                    USERS[user_id]['auto_join_state'] = st

        _OSThread(target=advanced_join_worker, daemon=True).start()
        return jsonify({
            "success": True,
            "total": len(norm_links),
            "message": f"بدأ الانضمام إلى {len(norm_links)} مجموعة — تابع التقدم في الأسفل"
        })
    except Exception as e:
        logger.error(f"auto_join_advanced error: {e}")
        return jsonify({"success": False, "message": str(e)})

class LearningBotManager:
    def __init__(self):
        self.bots = {}
        self.user_settings = {}

    def get_bot(self, user_id):
        if user_id not in self.bots:
            self.bots[user_id] = LearningBot(user_id)
        return self.bots[user_id]

    def is_active(self, user_id, chat_type='private'):
        """التحقق من تفعيل البوت لنوع معين من المحادثات"""
        # تحميل من الذاكرة أولاً، وإذا لم يوجد نحمّل من الملف
        if user_id not in self.user_settings:
            saved = load_settings(user_id)
            self.user_settings[user_id] = {
                'active_private': saved.get('learning_active_private', False),
                'active_group': saved.get('learning_active_group', False),
            }
        settings = self.user_settings.get(user_id, {})
        if chat_type == 'private':
            return settings.get('active_private', False)
        elif chat_type == 'group':
            return settings.get('active_group', False)
        return settings.get('active_private', False) or settings.get('active_group', False)

    def set_active(self, user_id, active, chat_type='private'):
        """تفعيل/إلغاء تفعيل البوت لنوع معين من المحادثات"""
        if user_id not in self.user_settings:
            self.user_settings[user_id] = {}
        if chat_type == 'private':
            self.user_settings[user_id]['active_private'] = active
        elif chat_type == 'group':
            self.user_settings[user_id]['active_group'] = active
        settings = load_settings(user_id)
        settings['learning_active_private'] = self.user_settings[user_id].get('active_private', False)
        settings['learning_active_group'] = self.user_settings[user_id].get('active_group', False)
        save_settings(user_id, settings)

    def get_settings(self, user_id):
        """الحصول على إعدادات البوت للمستخدم"""
        saved = load_settings(user_id)
        if user_id not in self.user_settings:
            self.user_settings[user_id] = {
                'active_private': saved.get('learning_active_private', False),
                'active_group': saved.get('learning_active_group', False),
            }
        return self.user_settings.get(user_id, {
            'active_private': False,
            'active_group': False
        })


class LearningBot:
    def __init__(self, user_id):
        self.user_id = user_id
        self.knowledge = self.load_knowledge()
        self.unknown_requests = []
        self.conversations = {}
        self.groq_client = None
        try:
            from groq import Groq as _Groq
            api_key = os.environ.get('GROQ_API_KEY', '').strip()
            if api_key:
                self.groq_client = _Groq(api_key=api_key)
        except Exception as e:
            logger.warning(f"Failed to initialize Groq client for learning bot: {e}")

    def load_knowledge(self):
        path = os.path.join(get_user_session_dir(self.user_id), "knowledge.json")
        if os.path.exists(path):
            try:
                with open(path, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except:
                pass
        return {
            "حل واجب": {"description": "حل الواجبات والمسائل الدراسية", "keywords": ["حل", "واجب", "مسألة", "تمارين"], "price_range": "50-200 ريال", "time_range": "2-24 ساعة"},
            "بحث": {"description": "إعداد البحوث الأكاديمية", "keywords": ["بحث", "تقرير", "موضوع"], "price_range": "100-500 ريال", "time_range": "1-5 أيام"},
            "تلخيص": {"description": "تلخيص الكتب والمحاضرات", "keywords": ["تلخيص", "ملخص", "اختصار"], "price_range": "30-150 ريال", "time_range": "2-12 ساعة"},
            "ترجمة": {"description": "ترجمة النصوص", "keywords": ["ترجمة", "ترجم", "ترجمة نص"], "price_range": "20-100 ريال لكل صفحة", "time_range": "1-24 ساعة"},
            "تحليل بيانات": {"description": "تحليل إحصائي وبيانات", "keywords": ["تحليل", "بيانات", "إحصاء", "SPSS", "Excel"], "price_range": "100-400 ريال", "time_range": "1-3 أيام"},
            "تصميم": {"description": "تصميم عروض وبوسترات", "keywords": ["تصميم", "بوستر", "عرض", "PowerPoint", "PPT"], "price_range": "50-250 ريال", "time_range": "2-24 ساعة"}
        }

    def save_knowledge(self):
        path = os.path.join(get_user_session_dir(self.user_id), "knowledge.json")
        with open(path, 'w', encoding='utf-8') as f:
            json.dump(self.knowledge, f, ensure_ascii=False, indent=2)

    def detect_service(self, text):
        """كشف الخدمة المطلوبة من النص"""
        text_low = text.lower()
        best_match = None
        best_score = 0
        for service, data in self.knowledge.items():
            for kw in data.get('keywords', []):
                if kw in text_low:
                    score = len(kw)
                    if score > best_score:
                        best_score = score
                        best_match = service
        return best_match

    def is_service_request(self, text: str) -> tuple:
        """تحليل النص لتحديد إذا كان طلب خدمة حقيقي أم إعلان/رسالة عادية"""
        if not self.groq_client:
            return self._fallback_detection(text)
        try:
            response = self.groq_client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[
                    {"role": "system", "content": """أنت مساعد لتحليل الرسائل. حدد نوع الرسالة:
- 'service': إذا كان المستخدم يطلب خدمة أكاديمية (حل واجب، بحث، ترجمة، تلخيص، تحليل بيانات، تصميم)
- 'promo': إذا كانت الرسالة إعلانية أو ترويجية (يوجد روابط، أرقام تواصل، عروض)
- 'normal': إذا كانت رسالة عادية أو سؤال عام

أجب فقط بكلمة واحدة: service أو promo أو normal"""},
                    {"role": "user", "content": text[:500]}
                ],
                max_tokens=10,
                temperature=0.1
            )
            result = response.choices[0].message.content.strip().lower()
            if result == 'service':
                return True, "service"
            elif result == 'promo':
                return False, "promo"
            else:
                return False, "normal"
        except Exception as e:
            logger.error(f"AI classification error: {e}")
            return self._fallback_detection(text)

    def _fallback_detection(self, text: str) -> tuple:
        """كشف تقليدي احتياطي"""
        text_low = text.lower()
        service_keywords = ['حل', 'واجب', 'بحث', 'تقرير', 'تلخيص', 'ترجمة', 'تحليل', 'تصميم', 'مساعدة']
        promo_keywords = ['للتواصل', 'واتساب', 'تليجرام', 'إعلان', 'عرض', 'خصم', 'كود', 'كاش', 'رابط']
        is_service = any(kw in text_low for kw in service_keywords)
        is_promo = any(kw in text_low for kw in promo_keywords)
        if is_service and not is_promo:
            return True, "service"
        elif is_promo:
            return False, "promo"
        else:
            return False, "normal"

    async def generate_intelligent_response(self, sender_name: str, text: str, detected_service: str = None, conversation_context: str = None) -> str:
        """توليد رد ذكي باستخدام الذكاء الاصطناعي"""
        if not self.groq_client:
            return self._generate_fallback_response(detected_service)
        services_info = "\n".join([
            f"- {s}: {d['description']} (السعر: {d.get('price_range', 'حسب الطلب')}, الوقت: {d.get('time_range', 'حسب الطلب')})"
            for s, d in self.knowledge.items()
        ])
        context = f"\nسياق المحادثة السابقة: {conversation_context}" if conversation_context else ""
        system_prompt = f"""أنت مساعد ذكي في مركز سرعة انجاز للخدمات الأكاديمية. لديك المعرفة التالية:
{services_info}

قواعد الرد:
1. إذا كان المستخدم يطلب خدمة محددة، اسأل عن التفاصيل: نوع المهمة، عدد الصفحات/الأسئلة، الموعد النهائي
2. إذا سأل عن السعر، أعطه نطاق سعري حسب الخدمة
3. إذا سأل عن الوقت، أعطه الوقت المتوقع
4. ردودك قصيرة ومفيدة (جملتين إلى ثلاث)
5. استخدم لغة عربية بسيطة وودية
6. اطلب من المستخدم إرسال الملف أو المهمة إذا لزم الأمر
{context}

الخدمة المكتشفة: {detected_service or 'غير محددة'}
"""
        try:
            response = self.groq_client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": f"رسالة من {sender_name}: {text}"}
                ],
                max_tokens=300,
                temperature=0.7,
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            logger.error(f"Groq response error: {e}")
            return self._generate_fallback_response(detected_service)

    def _generate_fallback_response(self, detected_service: str = None) -> str:
        """توليد رد احتياطي في حالة فشل الذكاء الاصطناعي"""
        if detected_service and detected_service in self.knowledge:
            data = self.knowledge[detected_service]
            return (f"📚 {detected_service}: {data['description']}\n"
                    f"💰 السعر: {data.get('price_range', 'حسب الطلب')}\n"
                    f"⏰ الوقت المتوقع: {data.get('time_range', 'حسب الطلب')}\n\n"
                    f"يرجى إرسال التفاصيل (عدد الصفحات/الأسئلة، الموعد النهائي) وسنرد عليك قريباً.")
        return "📚 مرحباً! نحن في مركز سرعة انجاز نقدم خدمات أكاديمية متنوعة. أرسل لي تفاصيل ما تحتاجه (حل واجب، بحث، ترجمة، تلخيص) وسأرد عليك."

    def add_service(self, name, description, keywords):
        if name and description:
            self.knowledge[name] = {
                "description": description,
                "keywords": [k.strip() for k in keywords if k.strip()] or [name]
            }
            self.save_knowledge()
            return True
        return False

    def delete_service(self, name):
        if name in self.knowledge:
            del self.knowledge[name]
            self.save_knowledge()
            return True
        return False

    def get_unknown_requests(self):
        return self.unknown_requests

    def clear_unknown(self):
        self.unknown_requests = []

    async def handle_incoming_message(self, event, client_manager):
        try:
            user_id = self.user_id
            is_private = event.is_private
            is_group = event.is_group or event.is_channel

            if is_private and not learning_manager.is_active(user_id, 'private'):
                return
            if is_group and not learning_manager.is_active(user_id, 'group'):
                return

            message = event.message
            if not message.text:
                return
            if getattr(message, 'out', False):
                return

            text = message.text
            sender = await event.get_sender()
            sender_name = getattr(sender, 'first_name', '') or getattr(sender, 'username', '') or 'مستخدم'
            sender_id = str(getattr(sender, 'id', ''))

            is_service_req, msg_type = self.is_service_request(text)
            if msg_type == 'promo':
                logger.info(f"Ignoring promo message from {sender_name}")
                return

            conversation_key = f"{sender_id}_{event.chat_id}"
            context = self.conversations.get(conversation_key, {}).get('last_response', '')
            detected_service = self.detect_service(text)

            response = await self.generate_intelligent_response(sender_name, text, detected_service, context)

            self.conversations[conversation_key] = {
                'last_request': text,
                'last_response': response,
                'timestamp': time.time(),
                'detected_service': detected_service
            }

            now = time.time()
            to_delete = [k for k, v in self.conversations.items() if now - v.get('timestamp', 0) > 3600]
            for k in to_delete:
                del self.conversations[k]

            await event.reply(response)
            socketio.emit('log_update', {
                "message": f"🤖 رد تعليمي لـ {sender_name}: {response[:100]}..."
            }, to=user_id)

            if is_private and is_service_req and not detected_service:
                self.unknown_requests.append({
                    "text": text[:200],
                    "sender": sender_name,
                    "sender_id": sender_id,
                    "time": datetime.now().strftime("%Y-%m-%d %H:%M"),
                    "chat_id": event.chat_id
                })
                socketio.emit('new_unknown_request', self.unknown_requests[-1], to=user_id)

        except Exception as e:
            logger.error(f"Learning bot error for {self.user_id}: {str(e)}")

learning_manager = LearningBotManager()

@app.route("/api/learning/status", methods=["GET"])
def api_learning_status():
    user_id = session.get('user_id', 'user_1')
    settings = learning_manager.get_settings(user_id)
    return jsonify({
        "success": True,
        "active_private": settings.get('active_private', False),
        "active_group": settings.get('active_group', False),
        "reply_in_groups": settings.get('active_group', False)
    })

@app.route("/api/learning/toggle", methods=["POST"])
def api_learning_toggle():
    user_id = session.get('user_id', 'user_1')
    data = request.json
    chat_type = data.get('chat_type', 'private')
    current_settings = learning_manager.get_settings(user_id)
    if chat_type == 'private':
        current = current_settings.get('active_private', False)
    else:
        current = current_settings.get('active_group', False)
    new_active = not current if data.get('active') is None else bool(data.get('active'))
    learning_manager.set_active(user_id, new_active, chat_type)
    return jsonify({"success": True, "active": new_active, "chat_type": chat_type})

@app.route("/api/learning/toggle_all", methods=["POST"])
def api_learning_toggle_all():
    """تفعيل/إلغاء تفعيل كلا النوعين معاً"""
    user_id = session.get('user_id', 'user_1')
    data = request.json
    active_private = data.get('active_private', False)
    active_group = data.get('active_group', False)
    learning_manager.set_active(user_id, active_private, 'private')
    learning_manager.set_active(user_id, active_group, 'group')
    return jsonify({"success": True, "active_private": active_private, "active_group": active_group})

@app.route("/api/learning/services", methods=["GET"])
def api_learning_services():
    user_id = session.get('user_id', 'user_1')
    bot = learning_manager.get_bot(user_id)
    return jsonify({"success": True, "services": bot.knowledge})

@app.route("/api/learning/add_service", methods=["POST"])
def api_learning_add_service():
    user_id = session.get('user_id', 'user_1')
    data = request.json
    name = data.get('name', '').strip()
    description = data.get('description', '').strip()
    keywords = data.get('keywords', [])
    if not name or not description:
        return jsonify({"success": False, "message": "الاسم والوصف مطلوبان"})
    bot = learning_manager.get_bot(user_id)
    if bot.add_service(name, description, keywords):
        return jsonify({"success": True, "message": f"تم إضافة الخدمة {name}"})
    return jsonify({"success": False, "message": "فشل في الإضافة"})

@app.route("/api/learning/delete_service", methods=["POST"])
def api_learning_delete_service():
    user_id = session.get('user_id', 'user_1')
    data = request.json
    name = data.get('name', '')
    bot = learning_manager.get_bot(user_id)
    if bot.delete_service(name):
        return jsonify({"success": True, "message": f"تم حذف الخدمة {name}"})
    return jsonify({"success": False, "message": "الخدمة غير موجودة"})

@app.route("/api/learning/unknown_requests", methods=["GET"])
def api_learning_unknown():
    user_id = session.get('user_id', 'user_1')
    bot = learning_manager.get_bot(user_id)
    return jsonify({"success": True, "requests": bot.get_unknown_requests()})

@app.route("/api/learning/clear_unknown", methods=["POST"])
def api_learning_clear_unknown():
    user_id = session.get('user_id', 'user_1')
    bot = learning_manager.get_bot(user_id)
    bot.clear_unknown()
    return jsonify({"success": True, "message": "تم مسح الطلبات"})

def _normalize_auto_reply(rule):
    if not isinstance(rule, dict):
        return None
    keyword = (rule.get('keyword') or rule.get('trigger') or '').strip()
    reply = (rule.get('reply') or '').strip()
    if not keyword or not reply:
        return None
    scope = (rule.get('scope') or 'all').lower()
    if scope not in ('all', 'private', 'groups'):
        scope = 'all'
    match = (rule.get('match') or 'contains').lower()
    if match not in ('contains', 'exact', 'regex'):
        match = 'contains'
    return {
        'keyword': keyword,
        'reply': reply,
        'scope': scope,
        'match': match,
        'used_count': int(rule.get('used_count') or 0),
        'last_used': rule.get('last_used') or '',
    }

@app.route("/api/auto_replies", methods=["GET"])
@app.route("/api/get_auto_replies", methods=["GET"])
def api_get_auto_replies():
    user_id = session.get('user_id', 'user_1')
    settings = load_settings(user_id)
    return jsonify({
        "success": True,
        "enabled": settings.get('auto_reply_enabled', True),
        "auto_replies": settings.get('auto_replies', []) or []
    })

@app.route("/api/add_auto_reply", methods=["POST"])
def api_add_auto_reply():
    user_id = session.get('user_id', 'user_1')
    data = request.json or {}
    rule = _normalize_auto_reply({
        'keyword': data.get('keyword') or data.get('trigger') or '',
        'reply': data.get('reply') or '',
        'scope': data.get('scope') or 'all',
        'match': data.get('match') or 'contains',
    })
    if not rule:
        return jsonify({"success": False, "message": "❌ الكلمة المفتاحية ونص الرد مطلوبان"})

    settings = load_settings(user_id)
    rules = settings.get('auto_replies', []) or []
    rules.append(rule)
    settings['auto_replies'] = rules
    if save_settings(user_id, settings):
        return jsonify({"success": True, "message": "✅ تم إضافة الرد التلقائي", "auto_replies": rules})
    return jsonify({"success": False, "message": "❌ فشل حفظ القاعدة"})

@app.route("/api/update_auto_reply", methods=["POST"])
def api_update_auto_reply():
    user_id = session.get('user_id', 'user_1')
    data = request.json or {}
    try:
        index = int(data.get('index', -1))
    except (TypeError, ValueError):
        index = -1
    settings = load_settings(user_id)
    rules = settings.get('auto_replies', []) or []
    if not (0 <= index < len(rules)):
        return jsonify({"success": False, "message": "❌ الفهرس غير صحيح"})
    new_rule = _normalize_auto_reply({
        'keyword': data.get('keyword'),
        'reply': data.get('reply'),
        'scope': data.get('scope'),
        'match': data.get('match'),
        'used_count': rules[index].get('used_count'),
        'last_used': rules[index].get('last_used'),
    })
    if not new_rule:
        return jsonify({"success": False, "message": "❌ بيانات غير صالحة"})
    rules[index] = new_rule
    settings['auto_replies'] = rules
    save_settings(user_id, settings)
    return jsonify({"success": True, "message": "✅ تم تحديث القاعدة", "auto_replies": rules})

@app.route("/api/delete_auto_reply", methods=["POST"])
def api_delete_auto_reply():
    user_id = session.get('user_id', 'user_1')
    data = request.json or {}
    try:
        index = int(data.get('index', -1))
    except (TypeError, ValueError):
        index = -1
    settings = load_settings(user_id)
    rules = settings.get('auto_replies', []) or []
    if 0 <= index < len(rules):
        removed = rules.pop(index)
        settings['auto_replies'] = rules
        save_settings(user_id, settings)
        return jsonify({"success": True, "message": f"🗑️ تم حذف الرد '{removed.get('keyword','')[:30]}'",
                        "auto_replies": rules})
    return jsonify({"success": False, "message": "❌ فهرس غير صحيح"})

@app.route("/api/save_auto_replies", methods=["POST"])
def api_save_auto_replies():
    user_id = session.get('user_id', 'user_1')
    data = request.json or {}
    raw = data.get('auto_replies', []) or []
    cleaned = []
    for r in raw:
        nr = _normalize_auto_reply(r)
        if nr:
            cleaned.append(nr)
    settings = load_settings(user_id)
    settings['auto_replies'] = cleaned
    save_settings(user_id, settings)
    return jsonify({"success": True, "message": f"✅ تم حفظ {len(cleaned)} قاعدة رد", "auto_replies": cleaned})

@app.route("/api/toggle_auto_reply", methods=["POST"])
def api_toggle_auto_reply():
    user_id = session.get('user_id', 'user_1')
    data = request.json or {}
    enabled = bool(data.get('enabled', True))
    settings = load_settings(user_id)
    settings['auto_reply_enabled'] = enabled
    save_settings(user_id, settings)
    return jsonify({
        "success": True,
        "enabled": enabled,
        "message": "✅ تم تفعيل الردود التلقائية" if enabled else "⏸️ تم تعطيل الردود التلقائية"
    })

async def resolve_link_group_name(client, url):
    """يحاول الحصول على اسم المجموعة/القناة من رابط تيليجرام"""
    try:
        username = url.split('/')[-1].replace('@', '')
        if username.startswith('+'):
            return None  # روابط الدعوة الخاصة لا يمكن حلها
        entity = await client.get_entity(username)
        return getattr(entity, 'title', None) or getattr(entity, 'first_name', None)
    except Exception:
        return None

async def search_links_in_chats(client, since_date, resolve_names=True):
    from datetime import timezone as _tz
    # تحويل since_date إلى timezone-aware للمقارنة الصحيحة
    if since_date.tzinfo is None:
        since_dt = since_date.replace(tzinfo=_tz.utc)
    else:
        since_dt = since_date

    found_links = []
    seen_urls   = set()
    try:
        async for dialog in client.iter_dialogs():
            try:
                if not dialog.entity:
                    continue
                chat_title = dialog.title or "محادثة غير معروفة"
                async for message in client.iter_messages(dialog, limit=500):
                    if not message.date:
                        continue
                    msg_date = message.date if message.date.tzinfo else message.date.replace(tzinfo=_tz.utc)
                    if msg_date < since_dt:
                        break
                    if not message.text:
                        continue
                    links = extract_telegram_links(message.text)
                    sender_id = getattr(message, 'sender_id', None)
                    for link in links:
                        url = link['url']
                        if url and url not in seen_urls:
                            seen_urls.add(url)
                            found_links.append({
                                'url': url,
                                'username': link['username'],
                                'group_name': '',
                                'date': message.date.strftime('%Y-%m-%d %H:%M'),
                                'chat_title': chat_title,
                                'sender_id': str(sender_id) if sender_id else '',
                                'link_type': link.get('type', 'channel'),
                            })
                if len(found_links) >= 1000:
                    break
            except Exception as e:
                logger.warning(f"تخطي محادثة بسبب خطأ: {str(e)}")
                continue
    except Exception as e:
        logger.error(f"خطأ في البحث عن الروابط: {str(e)}")

    found_links.sort(key=lambda x: x['date'], reverse=True)

    # ── استخراج اسم المجموعة لكل رابط عام ──
    if resolve_names and found_links:
        resolve_count = 0
        MAX_RESOLVE = 80
        for item in found_links:
            if resolve_count >= MAX_RESOLVE:
                break
            if '/+' in item['url']:
                continue  # روابط الدعوة الخاصة لا يمكن حلها
            try:
                uname = item['url'].rstrip('/').split('/')[-1]
                if uname and len(uname) >= 3:
                    entity = await client.get_entity(uname)
                    title = (getattr(entity, 'title', None)
                             or getattr(entity, 'first_name', None))
                    if title:
                        item['group_name'] = title
                    resolve_count += 1
                    await asyncio.sleep(0.3)
            except Exception:
                pass

    return found_links

async def search_public_telegram(client, query, limit=50):
    results = []
    try:
        global_search = await client(functions.messages.SearchGlobalRequest(
            q=query,
            offset_date=None,
            offset_peer=None,
            offset_id=0,
            limit=limit
        ))
        for message in global_search.messages:
            if hasattr(message, 'peer_id') and hasattr(message.peer_id, 'channel_id'):
                channel_id = message.peer_id.channel_id
                for chat in global_search.chats:
                    if hasattr(chat, 'id') and chat.id == channel_id:
                        if isinstance(chat, Channel):
                            username = chat.username if hasattr(chat, 'username') else None
                            result_item = {
                                'id': str(chat.id),
                                'title': chat.title,
                                'username': username,
                                'participants_count': getattr(chat, 'participants_count', 0),
                                'megagroup': getattr(chat, 'megagroup', False),
                                'verified': getattr(chat, 'verified', False),
                                'scam': getattr(chat, 'scam', False)
                            }
                            if not any(r['id'] == result_item['id'] for r in results):
                                results.append(result_item)
        if len(results) < 10:
            try:
                if not query.startswith('@'):
                    potential_username = '@' + query.replace(' ', '').replace('@', '')
                    try:
                        entity = await client.get_entity(potential_username)
                        if isinstance(entity, (Channel, Chat)):
                            result_item = {
                                'id': str(entity.id),
                                'title': entity.title,
                                'username': getattr(entity, 'username', None),
                                'participants_count': getattr(entity, 'participants_count', 0),
                                'megagroup': getattr(entity, 'megagroup', False),
                                'verified': getattr(entity, 'verified', False),
                                'scam': getattr(entity, 'scam', False)
                            }
                            if not any(r['id'] == result_item['id'] for r in results):
                                results.append(result_item)
                    except Exception:
                        pass
            except Exception:
                pass
    except Exception as e:
        logger.warning(f"خطأ في البحث العام: {str(e)}")
    results.sort(key=lambda x: x.get('participants_count', 0), reverse=True)
    return results[:limit]

@app.route("/api/search_my_links", methods=["POST"])
def api_search_my_links():
    """بحث كلاسيكي — يُرجع كل النتائج مرة واحدة (للتوافق مع CSV)"""
    try:
        if 'user_id' not in session:
            return jsonify({"success": False, "message": "❌ يرجى تسجيل الدخول أولاً"})
        user_id = session['user_id']
        data = request.json
        days = data.get('days', 60)
        if days <= 0 or days > 365:
            days = 60
        with USERS_LOCK:
            if user_id not in USERS:
                return jsonify({"success": False, "message": "❌ المستخدم غير مسجل"})
            client_manager = USERS[user_id].get('client_manager')
            if not client_manager or not client_manager.client:
                return jsonify({"success": False, "message": "❌ يرجى تسجيل الدخول أولاً"})
        since_date = datetime.now() - timedelta(days=days)
        result = client_manager.run_coroutine(
            search_links_in_chats(client_manager.client, since_date)
        )
        return jsonify({"success": True, "links": result, "message": f"تم العثور على {len(result)} رابط"})
    except Exception as e:
        logger.error(f"خطأ في البحث عن الروابط: {str(e)}")
        return jsonify({"success": False, "message": f"❌ خطأ في البحث: {str(e)}"})

@app.route("/api/search_my_links/start", methods=["POST"])
def api_search_my_links_start():
    """بحث متدفق — يبدأ البحث في الخلفية ويُرسل النتائج فور اكتشافها عبر socket.io"""
    try:
        if 'user_id' not in session:
            return jsonify({"success": False, "message": "❌ يرجى تسجيل الدخول أولاً"})
        user_id = session['user_id']
        data = request.json or {}
        days = int(data.get('days', 60))
        if days <= 0 or days > 365:
            days = 60
        with USERS_LOCK:
            if user_id not in USERS:
                return jsonify({"success": False, "message": "❌ المستخدم غير مسجل"})
            client_manager = USERS[user_id].get('client_manager')
            if not client_manager or not client_manager.client:
                return jsonify({"success": False, "message": "❌ يرجى تسجيل الدخول أولاً"})

        async def _stream_search():
            from datetime import timezone as _tz
            since_dt = datetime.now(_tz.utc) - timedelta(days=days)
            seen_urls   = set()
            total_found = 0
            resolve_count = 0
            MAX_RESOLVE   = 80  # حد أقصى لتجنب الحظر من Telegram

            try:
                async for dialog in client_manager.client.iter_dialogs():
                    try:
                        if not dialog.entity:
                            continue
                        chat_title = dialog.title or "محادثة غير معروفة"
                        batch = []

                        # ── جلب رسائل المحادثة (الأحدث أولاً) ──
                        async for message in client_manager.client.iter_messages(dialog, limit=500):
                            if not message.date:
                                continue
                            # التوقف عند الوصول لرسائل أقدم من الفترة المحددة
                            msg_date = message.date if message.date.tzinfo else message.date.replace(tzinfo=_tz.utc)
                            if msg_date < since_dt:
                                break
                            if not message.text:
                                continue
                            links = extract_telegram_links(message.text)
                            for lnk in links:
                                url = lnk.get('url', '')
                                if url and url not in seen_urls:
                                    seen_urls.add(url)
                                    batch.append({
                                        'url': url,
                                        'username': lnk.get('username', ''),
                                        'group_name': '',
                                        'date': message.date.strftime('%Y-%m-%d %H:%M'),
                                        'chat_title': chat_title,
                                        'link_type': lnk.get('type', 'channel'),
                                    })
                                    total_found += 1

                        # ── استخراج اسم المجموعة لكل رابط عام في الدفعة ──
                        if batch and resolve_count < MAX_RESOLVE:
                            for item in batch:
                                if resolve_count >= MAX_RESOLVE:
                                    break
                                if '/+' in item['url']:
                                    continue  # روابط الدعوة الخاصة لا يمكن حلها
                                try:
                                    uname = item['url'].rstrip('/').split('/')[-1]
                                    if uname and len(uname) >= 3:
                                        entity = await client_manager.client.get_entity(uname)
                                        title = (getattr(entity, 'title', None)
                                                 or getattr(entity, 'first_name', None))
                                        if title:
                                            item['group_name'] = title
                                        resolve_count += 1
                                        await asyncio.sleep(0.25)
                                except Exception:
                                    pass

                        if batch:
                            socketio.emit('search_link_batch', {
                                'items': batch,
                                'chat_title': chat_title,
                                'count': total_found
                            }, to=user_id)
                        if total_found >= 1000:
                            break
                    except Exception as ex:
                        logger.warning(f"تخطي محادثة {dialog.title}: {ex}")
                        continue
            except Exception as e:
                logger.error(f"خطأ في البحث المتدفق: {e}")
                socketio.emit('search_links_done', {'total': total_found, 'error': str(e)}, to=user_id)
                return
            socketio.emit('search_links_done', {'total': total_found}, to=user_id)

        def _run_stream():
            client_manager.run_coroutine(_stream_search())

        _OSThread(target=_run_stream, daemon=True).start()
        return jsonify({"success": True, "message": "بدأ البحث — ستظهر النتائج فور اكتشافها"})
    except Exception as e:
        logger.error(f"خطأ في بدء البحث المتدفق: {e}")
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/search_my_links/csv", methods=["POST"])
def api_search_my_links_csv():
    """تحميل نتائج البحث كملف CSV"""
    try:
        if 'user_id' not in session:
            return jsonify({"success": False, "message": "❌ يرجى تسجيل الدخول"})
        user_id = session['user_id']
        data = request.json or {}
        days = int(data.get('days', 60))
        links = data.get('links', [])
        if not links:
            with USERS_LOCK:
                client_manager = USERS.get(user_id, {}).get('client_manager')
            if not client_manager or not client_manager.client:
                return jsonify({"success": False, "message": "❌ العميل غير متصل"})
            since_date = datetime.now() - timedelta(days=days)
            links = client_manager.run_coroutine(
                search_links_in_chats(client_manager.client, since_date, resolve_names=False)
            )
        output = io.StringIO()
        output.write('\ufeff')  # BOM للعربية في Excel
        output.write('الرابط,اسم المجموعة,وجد في,التاريخ,نوع الرابط\n')
        for item in links:
            url = str(item.get('url', '')).replace('"', '""')
            group_name = str(item.get('group_name', '') or item.get('username', '')).replace('"', '""')
            chat_title = str(item.get('chat_title', '')).replace('"', '""')
            date = str(item.get('date', ''))
            link_type = 'دعوة خاصة' if '+' in url else 'قناة/مجموعة'
            output.write(f'"{url}","{group_name}","{chat_title}","{date}","{link_type}"\n')
        csv_content = output.getvalue()
        response = make_response(csv_content)
        response.headers['Content-Type'] = 'text/csv; charset=utf-8-sig'
        response.headers['Content-Disposition'] = f'attachment; filename=telegram_links_{datetime.now().strftime("%Y%m%d_%H%M")}.csv'
        return response
    except Exception as e:
        logger.error(f"CSV export error: {e}")
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/app_logs", methods=["GET"])
def api_app_logs():
    """[مُصلح] إرجاع سجلات التطبيق — مخصصة لكل مستخدم"""
    try:
        level = request.args.get('level', 'ALL')
        user_id = request.args.get('user_id') or session.get('user_id')

        if user_id and user_id in PREDEFINED_USERS:
            records = _get_user_logs(user_id, None if level == 'ALL' else level)
        else:
            records = _mem_log_handler.get_records(None if level == 'ALL' else level)

        # إذا لم تكن هناك سجلات بعد — أضف سجلات حالة فورية
        if not records:
            from datetime import datetime as _dt
            _now = _dt.now().strftime('%H:%M:%S')
            records = [
                {'time': _now, 'level': 'INFO', 'msg': '🚀 النظام يعمل — في انتظار العمليات...', 'name': 'system'},
                {'time': _now, 'level': 'INFO', 'msg': '📡 Flask + SocketIO + Telethon جاهزون', 'name': 'system'},
            ]
            if user_id and user_id in PREDEFINED_USERS:
                with USERS_LOCK:
                    ud = USERS.get(user_id, {})
                st = ud.get('stats', {})
                records.append({
                    'time': _now, 'level': 'INFO',
                    'msg': (f"👤 {PREDEFINED_USERS[user_id]['name']}: "
                            f"{'✅ متصل' if ud.get('connected') else '⭕ غير متصل'} | "
                            f"مرسل: {st.get('sent', 0)} | أخطاء: {st.get('errors', 0)}"),
                    'name': user_id,
                })

        return jsonify({"success": True, "logs": records, "count": len(records)})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/diagnose_logs", methods=["POST"])
def api_diagnose_logs():
    """[مُصلح] تشخيص عمليات وأخطاء مستخدم محدد باستخدام الذكاء الاصطناعي"""
    try:
        data = request.json or {}
        logs_text = data.get('logs', '')
        user_id = data.get('user_id') or session.get('user_id')
        user_name = PREDEFINED_USERS.get(user_id, {}).get('name', user_id or 'غير محدد') if user_id else 'غير محدد'

        # جمع سجلات المستخدم إذا لم يُرسل نص
        if not logs_text:
            if user_id and user_id in PREDEFINED_USERS:
                recs = _get_user_logs(user_id, None)
            else:
                recs = _mem_log_handler.get_records('WARNING')
            logs_text = '\n'.join(f"[{r['time']}] {r['level']}: {r['msg']}" for r in recs[-80:])

        # بناء سياق المستخدم الحالي
        user_context = ""
        if user_id and user_id in PREDEFINED_USERS:
            with USERS_LOCK:
                ud = USERS.get(user_id, {})
            st = ud.get('stats', {})
            settings = ud.get('settings', {})
            user_context = (
                f"\n\nحالة {user_name} الحالية:\n"
                f"- الاتصال: {'✅ متصل' if ud.get('connected') else '❌ غير متصل'}\n"
                f"- المصادقة: {'✅ نعم' if ud.get('authenticated') else '❌ لا'}\n"
                f"- المراقبة: {'🟢 تعمل' if ud.get('is_running') else '⭕ متوقفة'}\n"
                f"- رسائل مرسلة: {st.get('sent', 0)} | أخطاء: {st.get('errors', 0)}\n"
                f"- الهاتف: {settings.get('phone', 'غير مسجل')}"
            )

        if not logs_text.strip():
            return jsonify({"success": True, "diagnosis": f"✅ لا توجد أخطاء مسجلة لـ {user_name} حتى الآن.{user_context}"})

        try:
            from groq import Groq as _Groq
            _g = _Groq(api_key=GROQ_API_KEY)
            system_prompt = (
                f"أنت خبير في تشخيص مشاكل تطبيقات Python/Flask/Telegram.\n"
                f"أنت تحلل سجلات المستخدم: {user_name} (ID: {user_id}).{user_context}\n\n"
                "مهمتك:\n"
                "1. تحديد الأخطاء الرئيسية وأسبابها لهذا المستخدم تحديداً\n"
                "2. تحليل العمليات الناجحة والفاشلة\n"
                "3. تفسير كل خطأ بلغة بسيطة\n"
                "4. اقتراح الحل المناسب لكل مشكلة\n"
                "5. ملخص عام لحالة هذا المستخدم\n"
                "الإجابة باللغة العربية فقط، منظمة ومختصرة."
            )
            resp = _g.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": f"سجلات {user_name}:\n{logs_text[:4000]}"}
                ],
                max_tokens=1000,
                temperature=0.3
            )
            diagnosis = resp.choices[0].message.content
        except Exception as ai_err:
            diagnosis = f"⚠️ تعذر الاتصال بالذكاء الاصطناعي: {ai_err}\n\nالسجلات لـ {user_name}:\n{logs_text[:500]}"
        return jsonify({"success": True, "diagnosis": diagnosis})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})


@app.route("/api/ai/analyze-error", methods=["POST"])
def api_ai_analyze_error():
    """تحليل خطأ محدد من السجلات باستخدام الذكاء الاصطناعي"""
    try:
        data = request.json or {}
        message  = data.get('message', '').strip()
        details  = data.get('details', {})
        timestamp = data.get('timestamp', '')
        user_id  = data.get('user_id') or session.get('user_id')
        user_name = PREDEFINED_USERS.get(user_id, {}).get('name', user_id or 'النظام') if user_id else 'النظام'

        if not message:
            return jsonify({"success": False, "error": "لا توجد رسالة خطأ"})

        user_context = ""
        if user_id and user_id in PREDEFINED_USERS:
            with USERS_LOCK:
                ud = USERS.get(user_id, {})
            st = ud.get('stats', {})
            user_context = (
                f"\nحالة المستخدم {user_name}:\n"
                f"- الاتصال: {'✅ متصل' if ud.get('connected') else '❌ غير متصل'}\n"
                f"- المراقبة: {'🟢 تعمل' if ud.get('is_running') else '⭕ متوقفة'}\n"
                f"- رسائل مرسلة: {st.get('sent', 0)} | أخطاء: {st.get('errors', 0)}"
            )

        prompt = (
            f"⛔ الخطأ:\n{message}\n\n"
            f"📋 التفاصيل:\n{str(details) if details else 'لا توجد'}\n"
            f"⏰ الوقت: {timestamp}"
            f"{user_context}\n\n"
            "المطلوب:\n"
            "1. تحديد نوع الخطأ بالضبط\n"
            "2. شرح سبب حدوثه\n"
            "3. خطوات الإصلاح العملية\n"
            "4. كيف نمنع تكراره"
        )

        from groq import Groq as _Groq
        _g = _Groq(api_key=GROQ_API_KEY)
        resp = _g.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {"role": "system", "content": (
                    f"أنت خبير تقني في تحليل أخطاء Python/Flask/Telegram. "
                    f"تساعد مستخدم '{user_name}' في برنامج 'مركز سرعة انجاز'. "
                    "أجب بالعربية بشكل منظم ومختصر."
                )},
                {"role": "user", "content": prompt}
            ],
            max_tokens=800,
            temperature=0.3
        )
        return jsonify({"success": True, "analysis": resp.choices[0].message.content})
    except Exception as e:
        return jsonify({"success": False, "error": f"فشل تحليل الخطأ: {str(e)}"})


@app.route("/api/ai/chat", methods=["POST"])
def api_ai_chat():
    """دردشة عامة مع الذكاء الاصطناعي حول البرنامج والأخطاء"""
    try:
        data     = request.json or {}
        message  = data.get('message', '').strip()
        history  = data.get('history', [])
        user_id  = data.get('user_id') or session.get('user_id')

        if not message:
            return jsonify({"success": False, "error": "الرجاء إدخال رسالة"})

        user_name = PREDEFINED_USERS.get(user_id, {}).get('name', user_id or 'المستخدم') if user_id else 'المستخدم'
        user_context = ""
        if user_id and user_id in PREDEFINED_USERS:
            with USERS_LOCK:
                ud = USERS.get(user_id, {})
            st = ud.get('stats', {})
            settings = ud.get('settings', {})
            user_context = (
                f"\nحالة {user_name}:\n"
                f"- الاتصال: {'✅ متصل' if ud.get('connected') else '❌ غير متصل'}\n"
                f"- المصادقة: {'✅ نعم' if ud.get('authenticated') else '❌ لا'}\n"
                f"- المراقبة: {'🟢 تعمل' if ud.get('is_running') else '⭕ متوقفة'}\n"
                f"- رسائل مرسلة: {st.get('sent', 0)} | أخطاء: {st.get('errors', 0)}"
            )

        system_prompt = (
            f"أنت مساعد برنامج 'مركز سرعة انجاز' للخدمات الطلابية والأكاديمية.\n"
            f"تساعد المستخدم '{user_name}' في تحليل الأخطاء وحل المشاكل.\n"
            f"البرنامج: Python/Flask + Telegram API + gevent.{user_context}\n\n"
            "كن دقيقاً ومباشراً، أجب بالعربية."
        )

        messages = [{"role": "system", "content": system_prompt}]
        for h in history[-10:]:
            if h.get('role') in ('user', 'assistant') and h.get('content'):
                messages.append({"role": h['role'], "content": h['content']})
        messages.append({"role": "user", "content": message})

        from groq import Groq as _Groq
        _g = _Groq(api_key=GROQ_API_KEY)
        resp = _g.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=messages,
            max_tokens=800,
            temperature=0.7
        )
        return jsonify({"success": True, "response": resp.choices[0].message.content})
    except Exception as e:
        return jsonify({"success": False, "error": f"فشل الدردشة: {str(e)}"})


@app.route("/api/search_public_channels", methods=["POST"])
def api_search_public_channels():
    try:
        if 'user_id' not in session:
            return jsonify({
                "success": False,
                "message": "❌ يرجى تسجيل الدخول أولاً"
            })

        user_id = session['user_id']
        data = request.json
        query = data.get('query', '').strip()
        if not query:
            return jsonify({
                "success": False,
                "message": "❌ يرجى كتابة نص للبحث"
            })
        limit = min(data.get('limit', 50), 100)

        with USERS_LOCK:
            if user_id not in USERS:
                return jsonify({
                    "success": False,
                    "message": "❌ المستخدم غير مسجل"
                })
            client_manager = USERS[user_id].get('client_manager')
            if not client_manager or not client_manager.client:
                return jsonify({
                    "success": False,
                    "message": "❌ يرجى تسجيل الدخول أولاً"
                })

        logger.info(f"🌐 بدء البحث العام للمستخدم {user_id} عن: {query}")
        result = client_manager.run_coroutine(
            search_public_telegram(client_manager.client, query, limit)
        )
        logger.info(f"✅ تم العثور على {len(result)} قناة/مجموعة للمستخدم {user_id}")
        return jsonify({
            "success": True,
            "channels": result,
            "message": f"تم العثور على {len(result)} قناة/مجموعة"
        })
    except Exception as e:
        logger.error(f"خطأ في البحث العام: {str(e)}")
        return jsonify({
            "success": False,
            "message": f"❌ خطأ في البحث: {str(e)}"
        })

# ===========================
# 📊 الأدوات الأكاديمية
# ===========================

@app.route("/tools/analyze_stats", methods=["POST"])
def api_academic_analyze_stats():
    try:
        data = request.get_json(force=True, silent=True) or {}
        raw = data.get('data', '')
        nums = [float(x) for x in re.findall(r'[-+]?\d*\.?\d+', str(raw)) if x]
        if len(nums) < 2:
            return jsonify({"error": "أدخل على الأقل رقمين للتحليل"}), 400

        arr = np.array(nums)
        mode_val = float(stats.mode(arr, keepdims=True).mode[0])
        stats_result = {
            "count":    int(len(arr)),
            "sum":      round(float(np.sum(arr)), 4),
            "mean":     round(float(np.mean(arr)), 4),
            "median":   round(float(np.median(arr)), 4),
            "mode":     round(mode_val, 4),
            "std":      round(float(np.std(arr)), 4),
            "variance": round(float(np.var(arr)), 4),
            "min":      round(float(np.min(arr)), 4),
            "max":      round(float(np.max(arr)), 4),
            "range":    round(float(np.max(arr) - np.min(arr)), 4),
            "q1":       round(float(np.percentile(arr, 25)), 4),
            "q3":       round(float(np.percentile(arr, 75)), 4),
            "iqr":      round(float(np.percentile(arr, 75) - np.percentile(arr, 25)), 4),
            "skewness": round(float(stats.skew(arr)), 4),
            "kurtosis": round(float(stats.kurtosis(arr)), 4),
        }

        fig, axes = plt.subplots(1, 2, figsize=(10, 4))
        axes[0].hist(arr, bins='auto', color='#4e73df', edgecolor='white', alpha=0.85)
        axes[0].axvline(stats_result['mean'],   color='red',   linestyle='--', label=f"المتوسط: {stats_result['mean']}")
        axes[0].axvline(stats_result['median'], color='green', linestyle='--', label=f"الوسيط: {stats_result['median']}")
        axes[0].set_title('توزيع البيانات'); axes[0].legend()
        axes[0].set_xlabel('القيمة'); axes[0].set_ylabel('التكرار')
        axes[1].boxplot(arr, vert=True, patch_artist=True,
                        boxprops=dict(facecolor='#4e73df', alpha=0.7))
        axes[1].set_title('المربع الجذري (Boxplot)'); axes[1].set_ylabel('القيمة')
        plt.tight_layout()
        buf = BytesIO()
        plt.savefig(buf, format='png', dpi=100, bbox_inches='tight')
        plt.close()
        buf.seek(0)
        chart_b64 = base64.b64encode(buf.read()).decode('utf-8')

        return jsonify({"success": True, "stats": stats_result, "chart": chart_b64})
    except Exception as e:
        logger.error(f"Analyze stats error: {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/tools/format_file", methods=["POST"])
def api_academic_format_file():
    try:
        if 'file' not in request.files:
            return jsonify({"error": "لم يتم رفع ملف"}), 400
        f = request.files['file']
        filename = f.filename.lower()
        content_parts = []

        file_bytes = f.read()
        if filename.endswith('.pdf'):
            if pdfplumber is None:
                return jsonify({"error": "مكتبة PDF غير متاحة، شغّل: pip install pdfplumber"}), 500
            import io as _io2
            with pdfplumber.open(_io2.BytesIO(file_bytes)) as pdf:
                for i, page in enumerate(pdf.pages[:50], 1):
                    text = page.extract_text() or ''
                    if text.strip():
                        content_parts.append(f"--- صفحة {i} ---\n{text.strip()}")
                    for tbl in page.extract_tables():
                        if tbl:
                            rows_txt = '\n'.join('\t'.join(str(c or '') for c in row) for row in tbl)
                            content_parts.append(f"[جدول]\n{rows_txt}")
        elif filename.endswith('.docx'):
            if docx is None:
                return jsonify({"error": "مكتبة Word غير متاحة، شغّل: pip install python-docx"}), 500
            import io as _io2
            doc_obj = docx.Document(_io2.BytesIO(file_bytes))
            for para in doc_obj.paragraphs:
                if para.text.strip():
                    style = para.style.name if para.style else ''
                    prefix = '# ' if 'Heading 1' in style else '## ' if 'Heading' in style else ''
                    content_parts.append(prefix + para.text)
            for tbl in doc_obj.tables:
                rows_txt = '\n'.join('\t'.join(c.text for c in row.cells) for row in tbl.rows)
                content_parts.append(f"[جدول]\n{rows_txt}")
        elif filename.endswith('.txt'):
            content_parts.append(file_bytes.decode('utf-8', errors='replace'))
        else:
            return jsonify({"error": "صيغة غير مدعومة. استخدم PDF أو DOCX أو TXT"}), 400

        full_text = '\n\n'.join(content_parts)

        # ── تحسين النص بالذكاء الاصطناعي (اختياري) ──
        ai_summary = None
        use_ai = request.form.get('use_ai', 'false').lower() == 'true'
        if use_ai and full_text.strip():
            try:
                from groq import Groq as _Groq
                _gc = _Groq(api_key=GROQ_API_KEY)
                ai_resp = _gc.chat.completions.create(
                    model='llama-3.3-70b-versatile',
                    messages=[
                        {"role": "system", "content":
                         "أنت مساعد أكاديمي. لخّص الوثيقة التالية بشكل منظّم: العنوان، الأقسام الرئيسية، الأهداف، النتائج، التوصيات."},
                        {"role": "user", "content": f"لخّص هذا النص:\n\n{full_text[:6000]}"}
                    ],
                    max_tokens=800,
                    temperature=0.3,
                )
                ai_summary = ai_resp.choices[0].message.content
            except Exception as _ae:
                ai_summary = f"[تعذّر التلخيص: {_ae}]"

        words = len(full_text.split())
        return jsonify({
            "success":    True,
            "text":       full_text[:15000],
            "words":      words,
            "chars":      len(full_text),
            "sections":   len(content_parts),
            "filename":   f.filename,
            "ai_summary": ai_summary
        })
    except Exception as e:
        logger.error(f"Format file error: {e}")
        return jsonify({"error": str(e)}), 500


# ════════════════════════════════════════════════════════════
#  تحويل HTML إلى Word — محسّن كلياً
#  يدعم: فواصل الصفحات | جداول كاملة | صور base64 | ألوان | RTL
# ════════════════════════════════════════════════════════════

# ألوان CSS المسمّاة
_CSS_NAMED_COLORS = {
    'black':(0,0,0),'white':(255,255,255),'red':(255,0,0),'green':(0,128,0),
    'blue':(0,0,255),'yellow':(255,255,0),'orange':(255,165,0),'purple':(128,0,128),
    'pink':(255,192,203),'gray':(128,128,128),'grey':(128,128,128),'brown':(165,42,42),
    'cyan':(0,255,255),'magenta':(255,0,255),'navy':(0,0,128),'teal':(0,128,128),
    'lime':(0,255,0),'maroon':(128,0,0),'olive':(128,128,0),'silver':(192,192,192),
    'gold':(255,215,0),'coral':(255,127,80),'salmon':(250,128,114),'turquoise':(64,224,208),
    'indigo':(75,0,130),'violet':(238,130,238),'darkblue':(0,0,139),'darkgreen':(0,100,0),
    'darkred':(139,0,0),'darkgray':(169,169,169),'lightblue':(173,216,230),
    'lightgreen':(144,238,144),'lightyellow':(255,255,224),'lightgray':(211,211,211),
    'crimson':(220,20,60),'deepskyblue':(0,191,255),'forestgreen':(34,139,34),
    'hotpink':(255,105,180),'limegreen':(50,205,50),'mediumblue':(0,0,205),
    'orangered':(255,69,0),'royalblue':(65,105,225),'seagreen':(46,139,87),
    'skyblue':(135,206,235),'slategray':(112,128,144),'steelblue':(70,130,180),
    'tomato':(255,99,71),'yellowgreen':(154,205,50),'beige':(245,245,220),
    'ivory':(255,255,240),'khaki':(240,230,140),'lavender':(230,230,250),
    'mintcream':(245,255,250),'snow':(255,250,250),'wheat':(245,222,179),
}

def _w2_css_color(color_str):
    """تحويل لون CSS إلى RGBColor — يدعم hex, rgb(), rgba(), named colors"""
    try:
        from docx.shared import RGBColor
        s = (color_str or '').strip().lower()
        if not s or s == 'transparent' or s == 'inherit' or s == 'currentcolor':
            return None
        # hex
        if s.startswith('#'):
            c = s.lstrip('#')
            if len(c) == 3:
                c = c[0]*2 + c[1]*2 + c[2]*2
            if len(c) >= 6:
                return RGBColor(int(c[0:2],16), int(c[2:4],16), int(c[4:6],16))
        # rgb / rgba
        if s.startswith('rgb'):
            nums = re.findall(r'[\d.]+', s)
            if len(nums) >= 3:
                r,g,b = int(float(nums[0])), int(float(nums[1])), int(float(nums[2]))
                return RGBColor(min(r,255), min(g,255), min(b,255))
        # named
        if s in _CSS_NAMED_COLORS:
            t = _CSS_NAMED_COLORS[s]
            if t:
                return RGBColor(t[0], t[1], t[2])
    except Exception:
        pass
    return None

def _w2_parse_style(style_str):
    """تحليل CSS style مضمن → dict (يحافظ على القيمة الأخيرة عند التكرار)"""
    props = {}
    for part in (style_str or '').split(';'):
        part = part.strip()
        if ':' in part:
            k, v = part.split(':', 1)
            props[k.strip().lower()] = v.strip()
    return props

def _w2_align(style_props, default='LEFT'):
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    m = {'center': WD_ALIGN_PARAGRAPH.CENTER,
         'left':   WD_ALIGN_PARAGRAPH.LEFT,
         'right':  WD_ALIGN_PARAGRAPH.RIGHT,
         'justify':WD_ALIGN_PARAGRAPH.JUSTIFY}
    ta = style_props.get('text-align','').lower()
    return m.get(ta, getattr(WD_ALIGN_PARAGRAPH, default))

def _w2_parse_font_size_pt(fs):
    """تحويل قيمة font-size CSS إلى نقاط (Pt) أو None"""
    if not fs:
        return None
    fs = fs.strip().lower()
    try:
        if fs.endswith('pt'):
            return float(fs[:-2].strip())
        if fs.endswith('px'):
            return round(float(fs[:-2].strip()) * 0.75, 1)
        if fs.endswith('em'):
            return round(float(fs[:-2].strip()) * 12.0, 1)
        if fs.endswith('rem'):
            return round(float(fs[:-3].strip()) * 12.0, 1)
        if fs.endswith('%'):
            return round(float(fs[:-1].strip()) * 12.0 / 100.0, 1)
    except Exception:
        pass
    return None

def _w2_apply_styles_to_run(run, styles):
    """تطبيق dict من CSS styles على run — يدعم وراثة الأنماط"""
    try:
        from docx.shared import Pt, RGBColor
        fw = styles.get('font-weight', '')
        if fw in ('bold','700','800','900','bolder'): run.bold = True
        fi = styles.get('font-style', '')
        if fi == 'italic' or fi == 'oblique': run.italic = True
        td = styles.get('text-decoration', '')
        if 'underline' in td: run.underline = True
        if 'line-through' in td: run.font.strike = True
        c = _w2_css_color(styles.get('color', ''))
        if c: run.font.color.rgb = c
        ff = styles.get('font-family', '')
        if ff:
            fname = ff.split(',')[0].strip().strip('"\'')
            if fname: run.font.name = fname
        pt = _w2_parse_font_size_pt(styles.get('font-size', ''))
        if pt and 6 <= pt <= 72: run.font.size = Pt(pt)
    except Exception:
        pass

def _w2_apply_run(run, node, extra_tag=''):
    """تطبيق تنسيق inline على run (للتوافق مع الكود القديم)"""
    try:
        tag = (node.name or '').lower() if hasattr(node,'name') else extra_tag
        styles = {}
        if tag in ('b','strong'): styles['font-weight'] = 'bold'
        if tag in ('i','em'):     styles['font-style'] = 'italic'
        if tag == 'u':            styles['text-decoration'] = 'underline'
        if tag in ('s','strike','del'): styles['text-decoration'] = 'line-through'
        if tag == 'code':         styles['font-family'] = 'Courier New'
        if hasattr(node,'get'):
            styles.update(_w2_parse_style(node.get('style','')))
            if tag == 'font':
                fc = node.get('color',''); ff = node.get('face','')
                if fc: styles['color'] = fc
                if ff: styles['font-family'] = ff
        _w2_apply_styles_to_run(run, styles)
    except Exception:
        pass

def _w2_inline(node, para, inherited_styles=None):
    """
    معالجة المحتوى inline مع وراثة الأنماط الكاملة من الوالدين.
    inherited_styles: dict من CSS properties موروثة من العنصر الأب
    """
    from bs4 import NavigableString, Tag
    if inherited_styles is None:
        inherited_styles = {}

    INLINE_TAGS = {
        'b','strong','i','em','u','s','strike','del','span','a','font',
        'mark','sup','sub','small','big','code','kbd','abbr','cite','q',
        'bdi','bdo','time','var','ins','samp','dfn',
    }

    for child in node.children:
        if isinstance(child, NavigableString):
            txt = str(child)
            if txt:
                run = para.add_run(txt)
                _w2_apply_styles_to_run(run, inherited_styles)
        elif isinstance(child, Tag):
            ctag = (child.name or '').lower()
            if ctag == 'br':
                para.add_run('\n')
                continue
            if ctag == 'img':
                continue  # تعالجها _w2_node

            # بناء أنماط الابن = أنماط الوالد + أنماط الابن الخاصة
            child_styles = dict(inherited_styles)

            # أنماط مبنية على اسم الوسم
            if ctag in ('b','strong'):   child_styles['font-weight'] = 'bold'
            if ctag in ('i','em'):       child_styles['font-style']  = 'italic'
            if ctag == 'u':              child_styles['text-decoration'] = 'underline'
            if ctag in ('s','strike','del','ins'): child_styles['text-decoration'] = 'line-through'
            if ctag == 'mark':           child_styles['background-color'] = '#ffff00'
            if ctag in ('small',):
                old_pt = _w2_parse_font_size_pt(child_styles.get('font-size','12pt')) or 12
                child_styles['font-size'] = f'{old_pt * 0.85:.1f}pt'
            if ctag == 'code':           child_styles['font-family'] = 'Courier New'
            if ctag == 'font':
                fc = child.get('color',''); ff = child.get('face','')
                if fc: child_styles['color'] = fc
                if ff: child_styles['font-family'] = ff

            # الأنماط المضمّنة الخاصة بالعنصر (تتجاوز الموروثة)
            own_sp = _w2_parse_style(child.get('style','') if hasattr(child,'get') else '')
            child_styles.update(own_sp)

            if ctag in INLINE_TAGS:
                # تحقق هل يحتوي على عناصر block
                has_block = any(
                    isinstance(c, Tag) and (c.name or '').lower() in
                    ('div','p','table','ul','ol','h1','h2','h3','h4','h5','h6')
                    for c in child.children
                )
                if has_block:
                    _w2_inline(child, para, child_styles)
                else:
                    run = para.add_run(child.get_text())
                    _w2_apply_styles_to_run(run, child_styles)
            else:
                # عنصر block داخل inline — استخرج النص مع الأنماط الحالية
                run = para.add_run(child.get_text())
                _w2_apply_styles_to_run(run, child_styles)

def _w2_add_page_break(doc):
    """إضافة فاصل صفحة حقيقي في Word"""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    para = doc.add_paragraph()
    run = para.add_run()
    br = OxmlElement('w:br')
    br.set(qn('w:type'), 'page')
    run._r.append(br)
    return para

def _w2_shade_cell(cell, hex_color_str):
    """تلوين خلفية خلية جدول بلون hex"""
    try:
        from docx.oxml import parse_xml
        xml = ('<w:shd xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"'
               ' w:val="clear" w:color="auto" w:fill="' + hex_color_str.upper() + '"/>')
        cell._tc.get_or_add_tcPr().append(parse_xml(xml))
    except Exception:
        pass

def _w2_shade_para(para, hex_color_str):
    """تلوين خلفية فقرة بلون hex"""
    try:
        from docx.oxml import parse_xml
        pPr = para._p.get_or_add_pPr()
        shd_xml = (f'<w:shd xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"'
                   f' w:val="clear" w:color="auto" w:fill="{hex_color_str.upper()}"/>')
        pPr.append(parse_xml(shd_xml))
    except Exception:
        pass

def _w2_set_para_border(para, color_hex='E6B422', side='left', sz=24):
    """إضافة حد ملوّن على جانب فقرة"""
    try:
        from docx.oxml import parse_xml
        pPr = para._p.get_or_add_pPr()
        bdr_xml = (
            f'<w:pBdr xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            f'<w:{side} w:val="single" w:sz="{sz}" w:space="4" w:color="{color_hex.upper()}"/>'
            f'</w:pBdr>'
        )
        pPr.append(parse_xml(bdr_xml))
    except Exception:
        pass

def _w2_set_rtl_para(para):
    """تعيين اتجاه الفقرة RTL للنص العربي"""
    try:
        from docx.oxml import parse_xml
        pPr = para._p.get_or_add_pPr()
        rtl_xml = '<w:bidi xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"/>'
        pPr.append(parse_xml(rtl_xml))
    except Exception:
        pass

def _w2_is_rtl(text):
    """كشف إذا كان النص عربياً/RTL"""
    arabic = sum(1 for c in text if '\u0600' <= c <= '\u06FF' or '\u0750' <= c <= '\u077F')
    return arabic > len(text) * 0.3

def _w2_table(node, doc):
    """
    تحويل <table> HTML إلى جدول Word كامل مع:
    - ألوان رأس الجدول من CSS
    - colspan + rowspan
    - خلفيات الخلايا
    - وراثة أنماط النص
    """
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    # جمع كل الصفوف من thead + tbody + tfoot
    all_rows = []
    thead_rows = set()
    for section in node.find_all(['thead','tbody','tfoot'], recursive=False) or [node]:
        is_head = section.name == 'thead'
        for tr in section.find_all('tr', recursive=True if section == node else False):
            if is_head:
                thead_rows.add(id(tr))
            all_rows.append(tr)
    if not all_rows:
        all_rows = node.find_all('tr')
    if not all_rows:
        return

    # حساب أقصى عدد أعمدة مع مراعاة colspan
    def row_width(tr):
        w = 0
        for c in tr.find_all(['td','th'], recursive=False):
            w += max(1, int(c.get('colspan', 1) or 1))
        return w

    max_cols = max((row_width(r) for r in all_rows), default=1)
    if max_cols == 0:
        return

    num_rows = len(all_rows)
    table = doc.add_table(rows=num_rows, cols=max_cols)
    table.style = 'Table Grid'

    # شبكة لتتبع rowspan
    grid = [[False]*max_cols for _ in range(num_rows)]

    # قراءة لون border الجدول إن وجد
    tbl_sp = _w2_parse_style(node.get('style',''))
    tbl_border_color = tbl_sp.get('border-color', '').strip()

    for r_idx, tr in enumerate(all_rows):
        is_header_row = (id(tr) in thead_rows) or (r_idx == 0 and not thead_rows)
        # فحص إذا كان الصف هو tr داخل thead
        if tr.parent and tr.parent.name == 'thead':
            is_header_row = True

        cells_el = tr.find_all(['td','th'], recursive=False)
        col_cursor = 0

        for cell_el in cells_el:
            # تجاوز الخلايا المحجوزة بـ rowspan
            while col_cursor < max_cols and grid[r_idx][col_cursor]:
                col_cursor += 1
            if col_cursor >= max_cols:
                break

            colspan = max(1, int(cell_el.get('colspan', 1) or 1))
            rowspan = max(1, int(cell_el.get('rowspan', 1) or 1))
            is_th   = (cell_el.name == 'th') or is_header_row

            # تحديد الخلية في Word
            try:
                cell = table.cell(r_idx, col_cursor)
            except Exception:
                col_cursor += colspan
                continue

            # دمج colspan
            if colspan > 1:
                end_col = min(col_cursor + colspan - 1, max_cols - 1)
                if end_col > col_cursor:
                    try:
                        cell = cell.merge(table.cell(r_idx, end_col))
                    except Exception:
                        pass

            # دمج rowspan
            if rowspan > 1:
                end_row = min(r_idx + rowspan - 1, num_rows - 1)
                if end_row > r_idx:
                    try:
                        cell = cell.merge(table.cell(end_row, col_cursor))
                        # تحديد الشبكة
                        for rr in range(r_idx, end_row + 1):
                            for cc in range(col_cursor, min(col_cursor + colspan, max_cols)):
                                grid[rr][cc] = True
                    except Exception:
                        pass

            # مسح النص الافتراضي
            cell.text = ''
            para = cell.paragraphs[0] if cell.paragraphs else cell.add_paragraph()

            # تحليل أنماط الخلية
            cell_sp = _w2_parse_style(cell_el.get('style',''))
            ta = cell_sp.get('text-align','').lower()
            if not ta:
                ta = 'center' if is_th else 'left'
            align_map = {
                'center': WD_ALIGN_PARAGRAPH.CENTER,
                'left':   WD_ALIGN_PARAGRAPH.LEFT,
                'right':  WD_ALIGN_PARAGRAPH.RIGHT,
                'justify':WD_ALIGN_PARAGRAPH.JUSTIFY,
            }
            para.alignment = align_map.get(ta, WD_ALIGN_PARAGRAPH.LEFT)

            # تلوين خلفية الخلية
            bg = (cell_sp.get('background-color','') or cell_sp.get('background','')).strip()
            if not bg and is_th:
                bg = '#1E4A6E'  # أزرق غامق لرأس الجدول
            if bg:
                rgb = _w2_css_color(bg)
                if rgb:
                    hex_c = '{:02X}{:02X}{:02X}'.format(int(rgb[0]), int(rgb[1]), int(rgb[2]))
                    _w2_shade_cell(cell, hex_c)

            # الأنماط الموروثة للـ runs داخل الخلية
            cell_text_styles = {}
            txt_c = _w2_css_color(cell_sp.get('color',''))
            if is_th:
                cell_text_styles['font-weight'] = 'bold'
                if not txt_c:
                    txt_c = _w2_css_color('#FFFFFF')  # نص أبيض على رأس أزرق
            if txt_c:
                cell_text_styles['color'] = cell_sp.get('color','') or ('#FFFFFF' if is_th else '')

            # معالجة محتوى الخلية
            _w2_inline(cell_el, para, cell_text_styles)

            # تطبيق bold على رأس الجدول بعد _w2_inline
            if is_th:
                for run in para.runs:
                    run.bold = True
                    try:
                        if not run.font.color.rgb:
                            run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                    except Exception:
                        pass

            # تلوين خاص (matrix-high, matrix-mid, matrix-low)
            cell_classes = cell_el.get('class', [])
            if isinstance(cell_classes, str):
                cell_classes = cell_classes.split()
            if 'matrix-high' in cell_classes:
                _w2_shade_cell(cell, '2ECC71')
                for run in para.runs:
                    try: run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); run.bold=True
                    except Exception: pass
            elif 'matrix-mid' in cell_classes:
                _w2_shade_cell(cell, 'F39C12')
                for run in para.runs:
                    try: run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); run.bold=True
                    except Exception: pass
            elif 'matrix-low' in cell_classes:
                _w2_shade_cell(cell, 'E74C3C')
                for run in para.runs:
                    try: run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); run.bold=True
                    except Exception: pass

            col_cursor += colspan
            grid[r_idx][col_cursor-colspan] = True

    doc.add_paragraph()

def _w2_svg_to_png_bytes(svg_bytes, width_px=400, height_px=300):
    """محاولة تحويل SVG إلى PNG"""
    try:
        import cairosvg
        return cairosvg.svg2png(bytestring=svg_bytes, output_width=width_px, output_height=height_px)
    except Exception:
        pass
    try:
        from PIL import Image, ImageDraw
        img = Image.new('RGB', (width_px, height_px), color=(240, 245, 255))
        draw = ImageDraw.Draw(img)
        draw.rectangle([2,2,width_px-3,height_px-3], outline=(100,120,180), width=2)
        draw.text((10, height_px//2-10), "[ SVG Chart ]", fill=(100,120,180))
        from io import BytesIO as _BIO
        buf = _BIO()
        img.save(buf, 'PNG')
        return buf.getvalue()
    except Exception:
        return None

def _w2_embed_image(src, doc, width_inches=5.0, center=True, caption=None):
    """
    تضمين صورة في ملف Word — يدعم:
    - base64 PNG/JPEG/GIF/WebP/SVG
    - روابط HTTP/HTTPS
    - مسارات /static/ المحلية
    """
    try:
        from docx.shared import Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from io import BytesIO
        import base64 as _b64

        img_bytes = None
        if not src:
            return
        # تجاهل مسارات غير قابلة للوصول
        if src.startswith('/sdcard') or src.startswith('file://') or src == '#':
            return

        if 'data:image/svg' in src[:30]:
            parts = src.split(',', 1)
            raw = parts[1] if len(parts) == 2 else ''
            try:
                if '%3C' in raw or '%3c' in raw:
                    import urllib.parse
                    raw = urllib.parse.unquote(raw)
                svg_b = _b64.b64decode(raw) if (raw and not raw.startswith('<')) else raw.encode()
            except Exception:
                svg_b = raw.encode() if isinstance(raw, str) else raw
            img_bytes = _w2_svg_to_png_bytes(svg_b, 600, 400)
            if not img_bytes: return
            width_inches = min(width_inches, 5.5)

        elif src.startswith('data:image'):
            try:
                header, data = src.split(',', 1)
                img_bytes = _b64.b64decode(data + '==')
            except Exception:
                return

        elif src.startswith('http://') or src.startswith('https://'):
            try:
                resp = requests.get(src, timeout=12, headers={'User-Agent': 'Mozilla/5.0'})
                if resp.ok:
                    ct = resp.headers.get('Content-Type','')
                    img_bytes = _w2_svg_to_png_bytes(resp.content) if 'svg' in ct else resp.content
            except Exception:
                return

        elif src.startswith('/static/') or src.startswith('static/'):
            local_path = os.path.join(_PROJECT_ROOT, src.lstrip('/'))
            if os.path.exists(local_path):
                with open(local_path, 'rb') as _lf:
                    img_bytes = _lf.read()

        if not img_bytes:
            return

        # التحقق من صحة الصورة وتحويلها إذا لزم
        try:
            from PIL import Image as _PILImg
            pil_obj = _PILImg.open(BytesIO(img_bytes))
            if pil_obj.mode in ('RGBA','P','LA','CMYK'):
                pil_obj = pil_obj.convert('RGB')
                buf = BytesIO()
                pil_obj.save(buf, 'PNG')
                img_bytes = buf.getvalue()
            else:
                pil_obj.verify()  # تحقق فقط
        except Exception:
            try:
                from PIL import Image as _PILImg
                pil_obj = _PILImg.open(BytesIO(img_bytes)).convert('RGB')
                buf = BytesIO()
                pil_obj.save(buf, 'PNG')
                img_bytes = buf.getvalue()
            except Exception:
                return

        para = doc.add_paragraph()
        if center:
            para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = para.add_run()
        run.add_picture(BytesIO(img_bytes), width=Inches(min(float(width_inches), 6.0)))

        if caption:
            cap_p = doc.add_paragraph(caption)
            cap_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for r in cap_p.runs:
                r.italic = True
                try:
                    from docx.shared import Pt
                    r.font.size = Pt(10)
                except Exception: pass

    except Exception as e:
        logger.debug(f"Image embed skip ({str(src)[:60]}): {e}")

def _w2_node(node, doc, in_para=None):
    """
    المعالج الرئيسي: يحوّل عنصر HTML إلى عناصر DOCX مع دعم كامل لـ:
    - ألوان الخلفية على divs والفقرات
    - الجداول الكاملة مع rowspan/colspan
    - الصور بجميع أنواعها
    - القوائم المتداخلة
    - RTL للنص العربي
    - blockquote وpre وcode
    """
    from bs4 import NavigableString, Tag
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    if isinstance(node, NavigableString):
        txt = str(node)
        if in_para and txt.strip():
            in_para.add_run(txt)
        elif in_para and txt and not txt.strip():
            in_para.add_run(txt)
        return

    if not isinstance(node, Tag):
        return

    tag = (node.name or '').lower()
    classes = node.get('class', [])
    if isinstance(classes, str):
        classes = classes.split()
    classes_set = set(classes)
    sp = _w2_parse_style(node.get('style',''))

    # ── تجاهل العناصر غير المرئية ──
    if tag in ('script','style','noscript','meta','link','head'):
        return

    # ── فاصل صفحة ──
    is_page_break = (
        'page-break' in classes_set or 'pagebreak' in classes_set or
        sp.get('page-break-before','') == 'always' or
        sp.get('break-before','') in ('page','always') or
        sp.get('page-break-after','') == 'always'
    )
    if is_page_break and tag in ('div','section','p','hr','span'):
        _w2_add_page_break(doc)
        for child in node.children:
            _w2_node(child, doc)
        return

    # ── عناوين H1-H6 ──
    if tag in ('h1','h2','h3','h4','h5','h6'):
        lvl = int(tag[1])
        try:
            para = doc.add_heading(level=min(lvl, 4))
        except Exception:
            para = doc.add_paragraph()
        para.clear()
        sizes = {1:24, 2:20, 3:17, 4:15, 5:13, 6:12}

        # تطبيق خلفية العنوان إذا وُجدت
        bg = (sp.get('background-color','') or sp.get('background','')).strip()
        if bg:
            rgb_bg = _w2_css_color(bg)
            if rgb_bg:
                _w2_shade_para(para, '{:02X}{:02X}{:02X}'.format(int(rgb_bg[0]),int(rgb_bg[1]),int(rgb_bg[2])))

        # بناء أنماط الوراثة للعنوان
        heading_styles = {'font-weight': 'bold', 'font-size': f'{sizes.get(lvl,14)}pt'}
        c_str = sp.get('color','')
        if c_str: heading_styles['color'] = c_str

        _w2_inline(node, para, heading_styles)

        # تطبيق الحجم والـ bold على جميع runs بعد الإضافة
        for run in para.runs:
            run.bold = True
            run.font.size = Pt(sizes.get(lvl, 14))
            if c_str:
                c = _w2_css_color(c_str)
                if c: run.font.color.rgb = c

        para.alignment = _w2_align(sp, 'LEFT')

        # RTL للعناوين العربية
        txt_check = node.get_text()
        if _w2_is_rtl(txt_check):
            _w2_set_rtl_para(para)
        return

    # ── فقرات <p> ──
    if tag == 'p':
        para = doc.add_paragraph()
        para.alignment = _w2_align(sp, 'LEFT')

        # خلفية الفقرة
        bg = (sp.get('background-color','') or sp.get('background','')).strip()
        if bg:
            rgb_bg = _w2_css_color(bg)
            if rgb_bg:
                _w2_shade_para(para, '{:02X}{:02X}{:02X}'.format(int(rgb_bg[0]),int(rgb_bg[1]),int(rgb_bg[2])))

        # أنماط الفقرة
        p_styles = {}
        c_str = sp.get('color','')
        if c_str: p_styles['color'] = c_str
        fw = sp.get('font-weight','')
        if fw: p_styles['font-weight'] = fw
        fs = sp.get('font-size','')
        if fs: p_styles['font-size'] = fs
        ff = sp.get('font-family','')
        if ff: p_styles['font-family'] = ff

        _w2_inline(node, para, p_styles)

        txt_check = node.get_text()
        if _w2_is_rtl(txt_check):
            _w2_set_rtl_para(para)
        return

    # ── جدول ──
    if tag == 'table':
        _w2_table(node, doc)
        return

    # ── صورة ──
    if tag == 'img':
        src = node.get('src','')
        alt = node.get('alt','')
        w_attr = node.get('width','')
        try:
            w_in = min(float(str(w_attr).replace('px','').strip()) / 96.0, 6.0) if w_attr else 5.0
        except Exception:
            w_in = 5.0
        if src:
            _w2_embed_image(src, doc, width_inches=w_in)
        return

    # ── canvas ──
    if tag == 'canvas':
        # محاولة استخراج data-url إن وُجدت
        du = node.get('data-url','') or node.get('data-image','')
        if du:
            _w2_embed_image(du, doc, width_inches=5.5)
        return

    # ── figure ──
    if tag == 'figure':
        img_el = node.find('img')
        if img_el:
            src = img_el.get('src','')
            if src:
                _w2_embed_image(src, doc, width_inches=5.0)
        cap = node.find('figcaption')
        if cap:
            cp = doc.add_paragraph(cap.get_text(strip=True))
            cp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for r in cp.runs:
                r.italic = True
                try: r.font.size = Pt(10)
                except Exception: pass
        return

    # ── قوائم ──
    if tag == 'ul':
        for li in node.find_all('li', recursive=False):
            para = doc.add_paragraph(style='List Bullet')
            para.alignment = WD_ALIGN_PARAGRAPH.LEFT
            _w2_inline(li, para)
            if _w2_is_rtl(li.get_text()):
                _w2_set_rtl_para(para)
        return
    if tag == 'ol':
        for li in node.find_all('li', recursive=False):
            para = doc.add_paragraph(style='List Number')
            para.alignment = WD_ALIGN_PARAGRAPH.LEFT
            _w2_inline(li, para)
            if _w2_is_rtl(li.get_text()):
                _w2_set_rtl_para(para)
        return

    # ── hr ──
    if tag == 'hr':
        p = doc.add_paragraph()
        try:
            from docx.oxml import parse_xml
            pPr = p._p.get_or_add_pPr()
            bdr = ('<w:pBdr xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
                   '<w:bottom w:val="single" w:sz="6" w:space="1" w:color="AAAAAA"/>'
                   '</w:pBdr>')
            pPr.append(parse_xml(bdr))
        except Exception:
            p.add_run('─' * 60)
        return

    # ── br مستقل ──
    if tag == 'br':
        if in_para:
            in_para.add_run('\n')
        else:
            doc.add_paragraph()
        return

    # ── blockquote ──
    if tag == 'blockquote':
        para = doc.add_paragraph()
        _w2_inline(node, para)
        try:
            from docx.shared import Inches as _IN
            para.paragraph_format.left_indent  = _IN(0.4)
            para.paragraph_format.right_indent = _IN(0.4)
        except Exception: pass
        _w2_set_para_border(para, 'AAAAAA', 'left', 18)
        _w2_shade_para(para, 'F5F5F5')
        return

    # ── pre / code block ──
    if tag == 'pre':
        code_el = node.find('code')
        code_text = (code_el or node).get_text()
        for line in code_text.split('\n'):
            p = doc.add_paragraph(line)
            p.paragraph_format.left_indent = Pt(18)
            _w2_shade_para(p, 'F4F4F4')
            for r in p.runs:
                r.font.name = 'Courier New'
                try: r.font.size = Pt(10)
                except Exception: pass
        return

    # ── عناصر خاصة بالتصميم ──
    if 'insight-box' in classes_set:
        para = doc.add_paragraph()
        para.alignment = WD_ALIGN_PARAGRAPH.LEFT
        _w2_inline(node, para)
        _w2_shade_para(para, 'FEF8E7')
        _w2_set_para_border(para, 'E6B422', 'left', 24)
        return

    if 'toc-row' in classes_set or 'list-row' in classes_set:
        spans = node.find_all(['span','div'], recursive=False)
        texts = [s.get_text(strip=True) for s in spans] if spans else [node.get_text(strip=True)]
        line = '  '.join(filter(None, texts))
        if line:
            para = doc.add_paragraph(line)
            para.alignment = WD_ALIGN_PARAGRAPH.LEFT
        return

    if 'ref-item' in classes_set:
        para = doc.add_paragraph()
        para.alignment = WD_ALIGN_PARAGRAPH.LEFT
        _w2_inline(node, para)
        try:
            para.paragraph_format.left_indent = Pt(18)
            para.paragraph_format.space_after  = Pt(6)
        except Exception: pass
        return

    if 'figure-caption' in classes_set:
        para = doc.add_paragraph(node.get_text(strip=True))
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r in para.runs:
            r.italic = True
            try: r.font.size = Pt(10)
            except Exception: pass
        return

    if 'page-number' in classes_set:
        return

    if 'header-logo' in classes_set:
        img_el = node.find('img')
        if img_el:
            src = img_el.get('src','')
            if src and 'svg' not in src[:30]:
                _w2_embed_image(src, doc, width_inches=1.5)
        return

    # ── div / section مع خلفية ملوّنة → صندوق نصي ──
    CONTAINER_TAGS = {
        'div','section','article','main','aside','header','footer','nav',
        'body','html','form','fieldset','details','summary','address',
        'span','a','label','li',
    }
    if tag in CONTAINER_TAGS:
        bg = (sp.get('background-color','') or sp.get('background','')).strip()
        border_color = sp.get('border-left-color','') or sp.get('border-color','')

        if in_para and tag in ('span','a','label'):
            own_sp = {}
            own_sp.update(sp)
            _w2_inline(node, in_para, own_sp)
            return

        if bg and bg not in ('transparent','inherit','none',''):
            # معالجة الأبناء أولاً في فقرات عادية، ثم تلوين الفقرة الأولى
            start_idx = len(doc.paragraphs)
            for child in node.children:
                _w2_node(child, doc, in_para=None)
            end_idx = len(doc.paragraphs)
            rgb_bg = _w2_css_color(bg)
            if rgb_bg:
                hex_bg = '{:02X}{:02X}{:02X}'.format(int(rgb_bg[0]),int(rgb_bg[1]),int(rgb_bg[2]))
                for i in range(start_idx, end_idx):
                    try:
                        _w2_shade_para(doc.paragraphs[i], hex_bg)
                        if border_color:
                            bc = _w2_css_color(border_color)
                            if bc:
                                _w2_set_para_border(doc.paragraphs[i],
                                    '{:02X}{:02X}{:02X}'.format(int(bc[0]),int(bc[1]),int(bc[2])),
                                    'left', 18)
                    except Exception: pass
        else:
            for child in node.children:
                _w2_node(child, doc, in_para=None)
        return

    # ── أي عنصر آخر ──
    for child in node.children:
        _w2_node(child, doc, in_para=in_para)


@app.route("/formatter/")
@app.route("/formatter")
def formatter_page():
    """منسق الملفات"""
    resp = make_response(render_template('formatter.html'))
    resp.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate, max-age=0'
    return resp


@app.route("/formatter/<path:filename>")
def formatter_static(filename):
    """ملفات ثابتة لمنسق الملفات"""
    from flask import send_from_directory as _sfd
    return _sfd(os.path.join(os.path.dirname(__file__), 'static', 'wf_build'), filename)


# ════════════════════════════════════════════════════════════
#  PDF → Word مع صور كاملة وجداول (PyMuPDF)
# ════════════════════════════════════════════════════════════

def _smart_draw_charts(doc, text_content, add_charts=True):
    """نظام الرسم الذكي: يحلل المحتوى ويضيف رسوماً بيانية مناسبة"""
    if not add_charts:
        return
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import matplotlib.font_manager as fm
        import json as _json
        from io import BytesIO
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Inches, Pt

        # طلب AI لتحليل البيانات واقتراح الرسم
        from groq import Groq as _Groq
        client = _Groq(api_key=GROQ_API_KEY)
        prompt = f"""Analyze the following text and determine if it contains numerical data suitable for a chart.
Reply ONLY with valid JSON (no explanation):
{{
  "has_chart": true,
  "chart_type": "bar",
  "title": "Chart Title",
  "direction": "vertical",
  "colors": ["#1e4a2f","#2d6b41","#4f9a66","#7cb56e","#a8d08d"],
  "labels": ["Label1","Label2","Label3"],
  "values": [10, 25, 15],
  "x_label": "X axis",
  "y_label": "Y axis"
}}
If no chart is appropriate, set has_chart to false.
Text (first 2000 chars):
{text_content[:2000]}"""

        resp = client.chat.completions.create(
            model='llama-3.3-70b-versatile',
            messages=[{"role": "user", "content": prompt}],
            max_tokens=400, temperature=0.1
        )
        raw = resp.choices[0].message.content.strip()
        # استخراج JSON من الرد
        import re as _re
        m = _re.search(r'\{.*\}', raw, _re.DOTALL)
        if not m:
            return
        chart_data = _json.loads(m.group())
        if not chart_data.get('has_chart', False):
            return

        labels  = chart_data.get('labels', [])
        values  = chart_data.get('values', [])
        colors  = chart_data.get('colors', ['#1e4a2f','#2d6b41','#4f9a66','#7cb56e'])
        title   = chart_data.get('title', 'مخطط البيانات')
        ctype   = chart_data.get('chart_type', 'bar')
        direc   = chart_data.get('direction', 'vertical')
        xlabel  = chart_data.get('x_label', '')
        ylabel  = chart_data.get('y_label', '')

        if not labels or not values:
            return
        n = min(len(labels), len(values))
        labels, values = labels[:n], [float(v) for v in values[:n]]
        while len(colors) < n:
            colors += colors

        fig, ax = plt.subplots(figsize=(8, 4.5))
        fig.patch.set_facecolor('#f8fdf4')
        ax.set_facecolor('#f0f7ea')

        if ctype == 'pie':
            wedges, texts, autotexts = ax.pie(
                values, labels=labels, colors=colors[:n],
                autopct='%1.1f%%', startangle=90,
                pctdistance=0.82, wedgeprops=dict(width=0.6))
            for at in autotexts:
                at.set_fontsize(9); at.set_fontweight('bold')
        elif ctype == 'line':
            ax.plot(labels, values, marker='o', color=colors[0],
                    linewidth=2.5, markersize=7, markerfacecolor='white',
                    markeredgewidth=2)
            ax.fill_between(range(n), values, alpha=0.12, color=colors[0])
            ax.set_xticks(range(n)); ax.set_xticklabels(labels, rotation=20, ha='right')
            if xlabel: ax.set_xlabel(xlabel)
            if ylabel: ax.set_ylabel(ylabel)
            ax.grid(axis='y', linestyle='--', alpha=0.5)
        elif ctype == 'scatter':
            ax.scatter(range(n), values, c=colors[:n], s=120, zorder=5)
            ax.set_xticks(range(n)); ax.set_xticklabels(labels, rotation=20, ha='right')
        else:  # bar (default)
            if direc == 'horizontal':
                bars = ax.barh(labels, values, color=colors[:n], edgecolor='white', linewidth=0.5)
                if xlabel: ax.set_xlabel(xlabel)
                if ylabel: ax.set_ylabel(ylabel)
                for bar, val in zip(bars, values):
                    ax.text(bar.get_width()+max(values)*0.01, bar.get_y()+bar.get_height()/2,
                            f'{val:g}', va='center', fontsize=9)
            else:
                bars = ax.bar(labels, values, color=colors[:n], edgecolor='white',
                              linewidth=0.5, width=0.65)
                ax.set_xticks(range(n)); ax.set_xticklabels(labels, rotation=20, ha='right')
                if xlabel: ax.set_xlabel(xlabel)
                if ylabel: ax.set_ylabel(ylabel)
                for bar, val in zip(bars, values):
                    ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+max(values)*0.01,
                            f'{val:g}', ha='center', fontsize=9)
            ax.grid(axis='y' if direc != 'horizontal' else 'x', linestyle='--', alpha=0.4)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)

        ax.set_title(title, fontsize=13, fontweight='bold', pad=10)
        plt.tight_layout()

        buf = BytesIO()
        plt.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                    facecolor=fig.get_facecolor())
        plt.close(fig)
        buf.seek(0)

        # فاصل + الرسم + تسمية
        doc.add_paragraph('─' * 50)
        para = doc.add_paragraph()
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = para.add_run()
        run.add_picture(buf, width=Inches(5.5))
        cap = doc.add_paragraph(f"📊 {title}")
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r in cap.runs:
            r.italic = True
            try: r.font.size = Pt(10)
            except Exception: pass
        doc.add_paragraph()

    except Exception as e:
        logger.debug(f"Smart draw error: {e}")


@app.route("/tools/pdf_to_word", methods=["POST"])
def api_pdf_to_word():
    """
    تحويل PDF → DOCX مع:
    - الحفاظ الكامل على ألوان الخطوط والخلفيات
    - تجنب تكرار نص الجداول
    - جداول منسّقة مع تلوين رأس الجدول
    - صور مضمّنة كاملة
    - دعم RTL للنص العربي
    - تطبيق إعدادات المستخدم (خط / هوامش / حجم)
    """
    try:
        if 'file' not in request.files:
            return jsonify({"error": "لم يتم رفع ملف PDF"}), 400

        f = request.files['file']
        if not f.filename.lower().endswith('.pdf'):
            return jsonify({"error": "يرجى رفع ملف PDF فقط"}), 400

        # إعدادات المستخدم من الواجهة
        add_smart_draw  = request.form.get('smart_draw', 'false').lower() == 'true'
        user_font       = request.form.get('font_family', 'Times New Roman').strip() or 'Times New Roman'
        user_font_size  = float(request.form.get('font_size', '12') or '12')
        margin_in_str   = request.form.get('margin', '1.0')
        try:
            margin_in = float(margin_in_str)
        except Exception:
            margin_in = 1.0

        file_bytes = f.read()
        safe_name  = re.sub(r'[^\w\u0600-\u06FF\-_]', '_',
                            f.filename.rsplit('.', 1)[0]) or 'document'

        import io as _io
        from io import BytesIO
        import docx as _docx
        from docx.shared import Pt, Inches, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = _docx.Document()
        for sec in doc.sections:
            sec.top_margin    = Inches(margin_in)
            sec.bottom_margin = Inches(margin_in)
            sec.left_margin   = Inches(max(margin_in, 1.0))
            sec.right_margin  = Inches(max(margin_in, 1.0))
        doc.styles['Normal'].font.name = user_font
        doc.styles['Normal'].font.size = Pt(user_font_size)

        all_text = []
        total    = 0

        # ── دالة مساعدة: استخراج لون فيتز (int) → (r,g,b) أو None ──
        def _fitz_color(color_val):
            try:
                if color_val is None:
                    return None
                if isinstance(color_val, (list, tuple)) and len(color_val) >= 3:
                    return (int(color_val[0]*255), int(color_val[1]*255), int(color_val[2]*255))
                ci = int(color_val)
                r = (ci >> 16) & 0xFF
                g = (ci >> 8)  & 0xFF
                b =  ci        & 0xFF
                if r == g == b == 0:
                    return None  # أسود = افتراضي، تجاهل
                return (r, g, b)
            except Exception:
                return None

        # ── دالة: هل يتداخل بلوك نصي مع منطقة جدول ──
        def _in_table_region(block_bbox, table_bboxes, threshold=0.5):
            bx0,by0,bx1,by1 = block_bbox
            for tx0,ty0,tx1,ty1 in table_bboxes:
                ix0 = max(bx0,tx0); iy0 = max(by0,ty0)
                ix1 = min(bx1,tx1); iy1 = min(by1,ty1)
                if ix1 > ix0 and iy1 > iy0:
                    inter = (ix1-ix0)*(iy1-iy0)
                    area  = max((bx1-bx0)*(by1-by0), 1)
                    if inter/area > threshold:
                        return True
            return False

        try:
            import fitz as _fitz

            pdf_doc = _fitz.open(stream=file_bytes, filetype="pdf")
            total   = len(pdf_doc)

            # ── جمع بيانات الجداول من pdfplumber (bbox + data) ──
            plumb_pages = {}   # page_num -> {'bboxes': [...], 'tables': [...]}
            try:
                if pdfplumber:
                    with pdfplumber.open(_io.BytesIO(file_bytes)) as _plumb:
                        for pi, pp in enumerate(_plumb.pages[:total]):
                            found = pp.find_tables()
                            bboxes = [t.bbox for t in found]
                            data   = [t.extract() for t in found]
                            if bboxes:
                                plumb_pages[pi] = {'bboxes': bboxes, 'tables': data}
            except Exception:
                pass

            seen_img_xrefs = set()

            for page_num, page in enumerate(pdf_doc):
                if page_num > 0:
                    _w2_add_page_break(doc)

                page_table_bboxes = plumb_pages.get(page_num, {}).get('bboxes', [])

                # ── استخراج النص مع ألوان كاملة ──
                page_dict = page.get_text("dict", sort=True)

                for block in page_dict.get('blocks', []):
                    if block.get('type') != 0:
                        continue  # تجاهل بلوكات الصور — ستُعالَج لاحقاً

                    block_bbox = block.get('bbox', (0,0,0,0))

                    # تخطّ النص الواقع داخل منطقة جدول
                    if _in_table_region(block_bbox, page_table_bboxes):
                        continue

                    for line in block.get('lines', []):
                        # تجميع النص والتنسيق لكل span في السطر
                        spans_data = []
                        max_size = 0
                        line_is_bold = False
                        line_has_color = False

                        for span in line.get('spans', []):
                            t   = span.get('text', '')
                            if not t.strip():
                                continue
                            sz  = span.get('size', user_font_size)
                            flg = span.get('flags', 0)
                            col = span.get('color', 0)
                            rgb = _fitz_color(col)
                            is_bold = bool(flg & 16)
                            spans_data.append({
                                'text': t,
                                'size': sz,
                                'bold': is_bold,
                                'color': rgb,
                                'font': span.get('font',''),
                            })
                            if sz > max_size: max_size = sz
                            if is_bold: line_is_bold = True
                            if rgb: line_has_color = True

                        if not spans_data:
                            continue

                        line_txt = ' '.join(s['text'] for s in spans_data).strip()
                        all_text.append(line_txt)

                        # تحديد نوع الفقرة بناءً على الحجم
                        is_h1 = max_size >= 18
                        is_h2 = 14 <= max_size < 18
                        is_h3 = 12.5 <= max_size < 14 and line_is_bold

                        if is_h1:
                            para = doc.add_heading(level=1)
                            para.clear()
                            base_size = min(max_size, 24)
                        elif is_h2:
                            para = doc.add_heading(level=2)
                            para.clear()
                            base_size = min(max_size, 18)
                        elif is_h3:
                            para = doc.add_heading(level=3)
                            para.clear()
                            base_size = 14
                        else:
                            para = doc.add_paragraph()
                            base_size = user_font_size

                        # إذا كان السطر يحتوي على span واحد أو أقل اختلافاً في التنسيق
                        if len(spans_data) == 1 or (not line_has_color and not any(s['bold'] != line_is_bold for s in spans_data)):
                            run = para.add_run(line_txt)
                            run.bold = line_is_bold or is_h1 or is_h2 or is_h3
                            run.font.size = Pt(base_size)
                            run.font.name = user_font
                            if spans_data[0]['color']:
                                r,g,b = spans_data[0]['color']
                                try: run.font.color.rgb = RGBColor(r,g,b)
                                except Exception: pass
                        else:
                            # spans متعددة مع تنسيق مختلف
                            for span_d in spans_data:
                                run = para.add_run(span_d['text'] + ' ')
                                run.bold = span_d['bold'] or is_h1 or is_h2
                                run.font.size = Pt(min(span_d['size'], 36))
                                run.font.name = user_font
                                if span_d['color']:
                                    r,g,b = span_d['color']
                                    try: run.font.color.rgb = RGBColor(r,g,b)
                                    except Exception: pass

                        # RTL للنص العربي
                        if _w2_is_rtl(line_txt):
                            _w2_set_rtl_para(para)
                            para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                        elif is_h1 or is_h2:
                            para.alignment = WD_ALIGN_PARAGRAPH.CENTER

                # ── إدراج الجداول من pdfplumber ──
                page_tables = plumb_pages.get(page_num, {}).get('tables', [])
                for tbl_data in page_tables:
                    if not tbl_data:
                        continue
                    clean = [[str(c or '').strip() for c in row] for row in tbl_data]
                    clean = [r for r in clean if any(c for c in r)]
                    if not clean:
                        continue
                    max_cols = max(len(r) for r in clean)
                    if max_cols == 0:
                        continue
                    clean = [r + [''] * (max_cols - len(r)) for r in clean]

                    doc.add_paragraph()
                    tbl_word = doc.add_table(rows=len(clean), cols=max_cols)
                    tbl_word.style = 'Table Grid'

                    # تعيين عرض الأعمدة بالتساوي
                    try:
                        from docx.oxml.ns import qn as _qn
                        from docx.oxml import OxmlElement as _OXE
                        tbl_w = tbl_word._tbl
                        tblPr = tbl_w.tblPr if tbl_w.tblPr is not None else _OXE('w:tblPr')
                        tblW  = _OXE('w:tblW')
                        tblW.set(_qn('w:w'), '5000')
                        tblW.set(_qn('w:type'), 'pct')
                        tblPr.append(tblW)
                    except Exception:
                        pass

                    for ri, row in enumerate(clean):
                        for ci, val in enumerate(row):
                            cell = tbl_word.cell(ri, ci)
                            cell.text = ''
                            p = cell.paragraphs[0]
                            p.alignment = WD_ALIGN_PARAGRAPH.CENTER

                            is_rtl_cell = _w2_is_rtl(val)
                            run = p.add_run(val)
                            run.font.name = user_font
                            run.font.size = Pt(user_font_size)

                            if ri == 0:
                                run.bold = True
                                try: run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                                except Exception: pass
                                _w2_shade_cell(cell, '1E4A6E')  # أزرق غامق
                            elif ri % 2 == 0:
                                _w2_shade_cell(cell, 'EBF5FB')  # أزرق فاتح للصفوف الزوجية

                            if is_rtl_cell:
                                _w2_set_rtl_para(p)
                    doc.add_paragraph()

                # ── استخراج الصور من الصفحة ──
                img_list = page.get_images(full=True)
                for img_info in img_list:
                    xref = img_info[0]
                    if xref in seen_img_xrefs:
                        continue
                    seen_img_xrefs.add(xref)
                    try:
                        base_img   = pdf_doc.extract_image(xref)
                        img_bytes  = base_img["image"]
                        w_px, h_px = base_img.get("width",1), base_img.get("height",1)
                        if w_px < 80 or h_px < 80:
                            continue  # تجاهل ديكور صغير

                        from PIL import Image as _PILImg
                        try:
                            pil_img = _PILImg.open(BytesIO(img_bytes))
                            if pil_img.mode in ('RGBA','P','LA','CMYK'):
                                pil_img = pil_img.convert('RGB')
                            buf = BytesIO()
                            pil_img.save(buf, 'PNG')
                            buf.seek(0)
                            img_bytes = buf.getvalue()
                        except Exception:
                            pass

                        aspect = w_px / max(h_px, 1)
                        w_in   = min(5.5, max(1.5, aspect * 4.0))

                        para = doc.add_paragraph()
                        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        run  = para.add_run()
                        run.add_picture(BytesIO(img_bytes), width=Inches(w_in))
                    except Exception as _ie:
                        logger.debug(f"PDF image xref={xref}: {_ie}")

            pdf_doc.close()

        except ImportError:
            # ── fallback: pdfplumber فقط بدون PyMuPDF ──
            if pdfplumber is None:
                return jsonify({"error": "مكتبة PDF غير متاحة"}), 500
            with pdfplumber.open(_io.BytesIO(file_bytes)) as pdf:
                total = len(pdf.pages)
                for i, page in enumerate(pdf.pages[:60], 1):
                    if i > 1:
                        _w2_add_page_break(doc)

                    # جمع bboxes الجداول لهذه الصفحة
                    page_tbls = page.find_tables()
                    tbl_bboxes = [t.bbox for t in page_tbls]
                    tbl_data   = [t.extract() for t in page_tbls]

                    text = page.extract_text(layout=True) or ''
                    for line in text.split('\n'):
                        lt = line.strip()
                        if not lt:
                            continue
                        all_text.append(lt)
                        para = doc.add_paragraph()
                        run = para.add_run(lt)
                        run.font.name = user_font
                        run.font.size = Pt(user_font_size)
                        if _w2_is_rtl(lt):
                            _w2_set_rtl_para(para)
                            para.alignment = WD_ALIGN_PARAGRAPH.RIGHT

                    for tbl in tbl_data:
                        if not tbl:
                            continue
                        clean = [[str(c or '').strip() for c in r] for r in tbl]
                        clean = [r for r in clean if any(c for c in r)]
                        if not clean:
                            continue
                        mc = max(len(r) for r in clean)
                        clean = [r + ['']*(mc-len(r)) for r in clean]
                        tw = doc.add_table(rows=len(clean), cols=mc)
                        tw.style = 'Table Grid'
                        for ri, row in enumerate(clean):
                            for ci, val in enumerate(row):
                                c = tw.cell(ri, ci)
                                c.text = ''
                                p = c.paragraphs[0]
                                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                                run = p.add_run(val)
                                run.font.name = user_font
                                run.font.size = Pt(user_font_size)
                                if ri == 0:
                                    run.bold = True
                                    try: run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                                    except Exception: pass
                                    _w2_shade_cell(c, '1E4A6E')
                                elif ri % 2 == 0:
                                    _w2_shade_cell(c, 'EBF5FB')
                        doc.add_paragraph()

        # ── الرسم الذكي (اختياري) ──
        if add_smart_draw and all_text:
            _smart_draw_charts(doc, '\n'.join(all_text[:400]))

        final_path = str(_PPTX_OUTPUTS_DIR / (safe_name + '_converted.docx'))
        doc.save(final_path)

        return jsonify({
            "success": True,
            "download_url": f"/tools/word/download/{safe_name}_converted.docx",
            "filename": f"{safe_name}_converted.docx",
            "pages": total
        })

    except Exception as e:
        logger.error(f"PDF to Word error: {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500


@app.route("/tools/html_to_word", methods=["POST"])
def api_html_to_word():
    """
    تحويل HTML → DOCX مع دعم كامل لـ:
    - ألوان CSS كاملة مع وراثة الأنماط
    - جداول كاملة (colspan + rowspan + خلفيات)
    - صور base64 + HTTP + مسارات محلية
    - RTL للنص العربي تلقائياً
    - blockquote / pre / code / insight-box
    - تطبيق إعدادات المستخدم (خط / هوامش / حجم)
    """
    try:
        data = request.get_json() or {}
        html_content   = data.get('html', '')
        filename       = data.get('filename', 'document')
        add_smart_draw = data.get('smart_draw', False)

        # إعدادات المستخدم
        user_font      = (data.get('font_family','') or 'Times New Roman').strip()
        user_font_size = float(data.get('font_size', 12) or 12)
        margin_val     = float(data.get('margin', 1.0) or 1.0)
        rtl_doc        = data.get('rtl', False)

        if not html_content:
            return jsonify({"error": "محتوى HTML فارغ"}), 400

        from bs4 import BeautifulSoup
        import docx as _docx
        from docx.shared import Pt, Inches

        soup = BeautifulSoup(html_content, 'html.parser')

        doc = _docx.Document()

        # هوامش الصفحة
        for sec in doc.sections:
            sec.top_margin    = Inches(margin_val)
            sec.bottom_margin = Inches(margin_val)
            sec.left_margin   = Inches(max(margin_val, 0.8))
            sec.right_margin  = Inches(max(margin_val, 0.8))

        # النمط الافتراضي
        normal = doc.styles['Normal']
        normal.font.name = user_font
        normal.font.size = Pt(user_font_size)

        # تطبيق النمط الافتراضي على الـ Heading styles أيضاً
        for h_level in range(1, 5):
            try:
                hstyle = doc.styles[f'Heading {h_level}']
                hstyle.font.name = user_font
            except Exception:
                pass

        # إزالة العناصر غير المرئية
        body = soup.find('body') or soup
        for el in body.find_all(['script','style','noscript','head'], recursive=True):
            el.decompose()

        # معالجة المحتوى
        for child in body.children:
            _w2_node(child, doc)

        # حذف الفقرة الفارغة الأولى إن وُجدت
        try:
            first = doc.paragraphs[0]
            if not first.text.strip() and len(doc.paragraphs) > 1:
                first._element.getparent().remove(first._element)
        except Exception:
            pass

        # تطبيق اتجاه RTL على مستوى القسم إذا طُلب
        if rtl_doc:
            try:
                from docx.oxml import parse_xml
                for sec in doc.sections:
                    sectPr = sec._sectPr
                    bidi_xml = '<w:bidi xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"/>'
                    sectPr.append(parse_xml(bidi_xml))
            except Exception:
                pass

        # الرسم الذكي (اختياري)
        if add_smart_draw:
            body_text = ' '.join(p.text for p in doc.paragraphs if p.text.strip())[:3000]
            if body_text:
                _smart_draw_charts(doc, body_text)

        # حفظ الملف
        safe_name  = re.sub(r'[^\w\u0600-\u06FF\-_]', '_', filename) or 'document'
        final_path = str(_PPTX_OUTPUTS_DIR / (safe_name + '.docx'))
        doc.save(final_path)

        return jsonify({
            "success": True,
            "download_url": f"/tools/word/download/{safe_name}.docx",
            "filename": f"{safe_name}.docx"
        })
    except Exception as e:
        logger.error(f"HTML to Word error: {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500


@app.route("/tools/word/download/<path:filename>")
def tools_word_download(filename):
    """تحميل ملف Word المُنشأ"""
    from flask import send_from_directory
    outputs_dir = os.path.join(os.path.dirname(__file__), 'pptx_app', 'outputs')
    return send_from_directory(outputs_dir, filename, as_attachment=True)


@app.route("/tools/html_to_excel", methods=["POST"])
def api_html_to_excel():
    """تحويل HTML → Excel (.xlsx) مع استخراج الجداول والنصوص"""
    try:
        data = request.get_json() or {}
        html_content = data.get('html', '')
        filename     = data.get('filename', 'document')

        if not html_content:
            return jsonify({"error": "محتوى HTML فارغ"}), 400

        from bs4 import BeautifulSoup
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter

        soup = BeautifulSoup(html_content, 'html.parser')
        wb   = openpyxl.Workbook()
        ws   = wb.active
        ws.title = "المحتوى"

        current_row = 1

        # نمط الترويسات
        header_font  = Font(bold=True, size=12, color="FFFFFF")
        header_fill  = PatternFill("solid", fgColor="1a7a3c")
        center_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
        wrap_align   = Alignment(wrap_text=True, vertical="top")
        thin_border  = Border(
            left=Side(style='thin'), right=Side(style='thin'),
            top=Side(style='thin'), bottom=Side(style='thin')
        )

        def write_text_row(text, bold=False, size=11):
            nonlocal current_row
            if not text.strip():
                return
            cell = ws.cell(row=current_row, column=1, value=text.strip())
            cell.font = Font(bold=bold, size=size)
            cell.alignment = wrap_align
            current_row += 1

        def write_table(table_el):
            nonlocal current_row
            rows = table_el.find_all('tr')
            if not rows:
                return
            start_row = current_row
            max_col = 0
            for ri, tr in enumerate(rows):
                cells = tr.find_all(['th', 'td'])
                col = 1
                for ci, cell_el in enumerate(cells):
                    while ws.cell(row=current_row, column=col).value is not None:
                        col += 1
                    text  = cell_el.get_text(separator=' ', strip=True)
                    is_th = cell_el.name == 'th' or ri == 0
                    c = ws.cell(row=current_row, column=col, value=text)
                    c.border    = thin_border
                    c.alignment = center_align if is_th else wrap_align
                    if is_th:
                        c.font = header_font
                        c.fill = header_fill
                    else:
                        c.font = Font(size=11)
                    colspan = int(cell_el.get('colspan', 1))
                    rowspan = int(cell_el.get('rowspan', 1))
                    if colspan > 1 or rowspan > 1:
                        ws.merge_cells(
                            start_row=current_row, start_column=col,
                            end_row=current_row + rowspan - 1, end_column=col + colspan - 1
                        )
                    col += colspan
                    if col - 1 > max_col:
                        max_col = col - 1
                current_row += 1
            # ضبط عرض الأعمدة تلقائياً
            for c in range(1, max_col + 1):
                ws.column_dimensions[get_column_letter(c)].width = 22
            current_row += 1  # سطر فارغ بعد الجدول

        body = soup.find('body') or soup
        for el in body.find_all(['script', 'style', 'noscript'], recursive=True):
            el.decompose()

        for node in body.children:
            if not hasattr(node, 'name') or not node.name:
                t = str(node).strip()
                if t:
                    write_text_row(t)
                continue
            tag = node.name.lower()
            if tag in ('h1', 'h2', 'h3', 'h4', 'h5', 'h6'):
                level = int(tag[1])
                write_text_row(node.get_text(strip=True), bold=True, size=max(11, 16 - level))
            elif tag == 'table':
                write_table(node)
            elif tag in ('p', 'div', 'li', 'span', 'blockquote', 'pre'):
                # ابحث عن جداول بداخله
                inner_tables = node.find_all('table')
                if inner_tables:
                    for tbl in inner_tables:
                        write_table(tbl)
                else:
                    txt = node.get_text(separator=' ', strip=True)
                    if txt:
                        write_text_row(txt)
            elif tag in ('ul', 'ol'):
                for li in node.find_all('li', recursive=False):
                    write_text_row('• ' + li.get_text(strip=True))

        # حفظ الملف
        safe_name  = re.sub(r'[^\w\u0600-\u06FF\-_]', '_', filename) or 'document'
        outputs_dir = str(_PPTX_OUTPUTS_DIR)
        os.makedirs(outputs_dir, exist_ok=True)
        final_path = os.path.join(outputs_dir, safe_name + '.xlsx')
        wb.save(final_path)

        return jsonify({
            "success": True,
            "download_url": f"/tools/excel/download/{safe_name}.xlsx",
            "filename": f"{safe_name}.xlsx"
        })
    except Exception as e:
        logger.error(f"HTML to Excel error: {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500


@app.route("/tools/excel/download/<path:filename>")
def tools_excel_download(filename):
    """تحميل ملف Excel المُنشأ"""
    from flask import send_from_directory
    outputs_dir = os.path.join(os.path.dirname(__file__), 'pptx_app', 'outputs')
    return send_from_directory(outputs_dir, filename, as_attachment=True)


# ════════════════════════════════════════════════════════════
#  مساعد الذكاء الاصطناعي (AI Assistant) - يستطيع تعديل الكود
# ════════════════════════════════════════════════════════════
_PROJECT_ROOT = os.path.dirname(os.path.abspath(__file__))

def _ai_read_file(rel_path):
    """قراءة ملف من المشروع"""
    try:
        full = os.path.join(_PROJECT_ROOT, rel_path)
        full = os.path.realpath(full)
        if not full.startswith(_PROJECT_ROOT):
            return None, "مسار غير مسموح به"
        with open(full, 'r', encoding='utf-8', errors='replace') as f:
            return f.read(), None
    except Exception as e:
        return None, str(e)

def _ai_write_file(rel_path, content):
    """كتابة ملف في المشروع"""
    try:
        full = os.path.join(_PROJECT_ROOT, rel_path)
        full = os.path.realpath(full)
        if not full.startswith(_PROJECT_ROOT):
            return False, "مسار غير مسموح به"
        os.makedirs(os.path.dirname(full), exist_ok=True)
        with open(full, 'w', encoding='utf-8') as f:
            f.write(content)
        return True, None
    except Exception as e:
        return False, str(e)

def _ai_list_files():
    """قائمة ملفات المشروع"""
    result = []
    skip_dirs = {'.git', '__pycache__', '.local', 'node_modules', 'pptx_app/outputs', 'sessions', '.cache', '_extract'}
    skip_exts = {'.pyc', '.pyo', '.session', '.session-journal', '.lock', '.jpg', '.png', '.ico', '.svg', '.webp'}
    for root, dirs, files in os.walk(_PROJECT_ROOT):
        dirs[:] = [d for d in dirs if d not in skip_dirs and not d.startswith('.')]
        rel_root = os.path.relpath(root, _PROJECT_ROOT)
        for f in files:
            ext = os.path.splitext(f)[1].lower()
            if ext in skip_exts:
                continue
            rel = os.path.join(rel_root, f) if rel_root != '.' else f
            result.append(rel)
    return result[:200]


def _ai_github_push(files_dict: dict, commit_msg: str = "🤖 تعديل تلقائي بواسطة المساعد الذكي") -> dict:
    """
    رفع ملفات إلى GitHub عبر REST API.
    files_dict: {rel_path: content_str}
    """
    import base64 as _b64
    headers = {
        'Authorization': f'token {GITHUB_TOKEN}',
        'Accept': 'application/vnd.github.v3+json',
        'Content-Type': 'application/json',
    }
    api = 'https://api.github.com'
    results = {}
    for rel_path, content in files_dict.items():
        try:
            # الحصول على SHA الحالي للملف (إن وُجد)
            get_url = f"{api}/repos/{GITHUB_REPO}/contents/{rel_path}"
            gr = requests.get(get_url, headers=headers, params={'ref': GITHUB_BRANCH}, timeout=15)
            sha = gr.json().get('sha') if gr.ok else None
            # رفع الملف
            body = {
                'message': commit_msg,
                'content': _b64.b64encode(content.encode('utf-8', errors='replace')).decode(),
                'branch': GITHUB_BRANCH,
            }
            if sha:
                body['sha'] = sha
            pr = requests.put(get_url, headers=headers, json=body, timeout=30)
            results[rel_path] = 'pushed' if pr.ok else f"error {pr.status_code}"
        except Exception as ex:
            results[rel_path] = f"exception: {ex}"
    return results


def _ai_extract_code_blocks(text: str) -> list[dict]:
    """
    استخراج كتل الكود من رد الذكاء الاصطناعي.
    يدعم:  ```python app.py  أو  ```# app.py  أو  ```app.py
    """
    import re as _re
    blocks = []
    pattern = _re.compile(
        r'```(?P<lang>\w+)?\s*(?P<path>[^\n`]{2,100}\.[a-zA-Z]{1,6})?\n(?P<code>.*?)```',
        _re.DOTALL
    )
    for m in pattern.finditer(text):
        code = m.group('code').strip()
        path = (m.group('path') or '').strip()
        lang = (m.group('lang') or '').strip()
        blocks.append({'lang': lang, 'path': path, 'code': code})
    return blocks


@app.route("/api/ai_assistant", methods=["POST"])
def api_ai_assistant():
    """مساعد الذكاء الاصطناعي الكامل — قراءة، تعديل، رفع GitHub"""
    try:
        data     = request.get_json() or {}
        messages = data.get('messages', [])
        action   = data.get('action', 'chat')      # chat | apply | github_push
        target_file = data.get('file', '')         # ملف مستهدف يمكن تحديده يدوياً

        from groq import Groq as _Groq
        client = _Groq(api_key=GROQ_API_KEY)

        # ── قائمة الملفات والملخص ──────────────────────────────
        files_list   = _ai_list_files()
        file_summaries = []
        for mf in ['app.py', 'templates/index.html', 'templates/academic.html',
                   'static/js/app.js', 'requirements.txt', 'render.yaml']:
            c, _ = _ai_read_file(mf)
            if c:
                file_summaries.append(f"- {mf} ({len(c.splitlines())} سطر)")

        system_prompt = f"""أنت مساعد ذكاء اصطناعي متخصص في تطوير مشروع "مركز سرعة انجاز".

المشروع: Flask + Python + Telethon + Groq AI
الملفات الرئيسية:
{chr(10).join(file_summaries)}

جميع ملفات المشروع ({len(files_list)} ملف):
{chr(10).join(files_list[:60])}

مستودع GitHub: https://github.com/{GITHUB_REPO}  (فرع {GITHUB_BRANCH})

═══ قدراتك الكاملة ═══
1. قراءة أي ملف: اذكر اسم الملف وسيُقرأ تلقائياً
2. تعديل الكود: قدّم الكود في كتلة ``` مع اسم الملف في أول سطرها
   مثال:
   ```python app.py
   # الكود هنا
   ```
3. كتابة ملفات جديدة: نفس الأسلوب باسم ملف جديد
4. رفع إلى GitHub: يتم تلقائياً عند الطلب

تعليمات التعديل التلقائي:
- عند action=apply: سيُطبَّق الكود الأول المُعثَر عليه في ردك على الملف المحدد
- عند action=github_push: بعد التطبيق، يُرفع إلى GitHub تلقائياً
- عند طلب تعديل، اكتب الكود الكامل للملف المعني، لا جزءاً منه فقط

أجب بالعربية دائماً. كن دقيقاً وعملياً."""

        history = [{'role': 'system', 'content': system_prompt}] + messages[-20:]

        # ── إضافة سياق الملفات المذكورة في الرسالة ─────────────
        last_msg = messages[-1]['content'] if messages else ''
        file_context = ""
        for fname in files_list:
            if fname in last_msg or os.path.basename(fname) in last_msg:
                c, _ = _ai_read_file(fname)
                if c:
                    preview = c[:4000] + ('\n... [مقطوع]' if len(c) > 4000 else '')
                    file_context += f"\n\n═══ محتوى {fname} ═══\n{preview}"
        if target_file and not any(target_file in fc for fc in [file_context]):
            c, _ = _ai_read_file(target_file)
            if c:
                file_context += f"\n\n═══ محتوى {target_file} (الملف المستهدف) ═══\n{c[:5000]}"
        if file_context:
            history.append({'role': 'system', 'content': f"محتوى الملفات المطلوبة:{file_context}"})

        # ── استدعاء GROQ ──────────────────────────────────────
        resp = client.chat.completions.create(
            model='llama-3.3-70b-versatile',
            messages=history,
            max_tokens=4000,
            temperature=0.5,
        )
        reply = resp.choices[0].message.content

        # ── تطبيق الكود تلقائياً إذا طُلب ─────────────────────
        apply_result  = None
        github_result = None

        if action in ('apply', 'github_push'):
            blocks = _ai_extract_code_blocks(reply)
            applied = {}
            for blk in blocks:
                # تحديد مسار الملف
                fpath = blk['path'] or target_file
                if not fpath:
                    continue
                # تنظيف المسار
                fpath = fpath.strip().strip('`').strip()
                ok, err = _ai_write_file(fpath, blk['code'])
                applied[fpath] = 'تم التطبيق ✅' if ok else f'فشل: {err}'

            apply_result = applied if applied else {"info": "لم يُعثر على كود قابل للتطبيق في الرد"}

            if action == 'github_push' and applied:
                files_to_push = {}
                for fpath, status in applied.items():
                    if '✅' in status:
                        c, _ = _ai_read_file(fpath)
                        if c:
                            files_to_push[fpath] = c
                if files_to_push:
                    github_result = _ai_github_push(
                        files_to_push,
                        commit_msg=f"🤖 تعديل بواسطة المساعد الذكي: {last_msg[:80]}"
                    )

        return jsonify({
            "success":       True,
            "reply":         reply,
            "apply_result":  apply_result,
            "github_result": github_result,
            "files_list":    files_list[:40]
        })
    except Exception as e:
        logger.error(f"AI Assistant error: {e}", exc_info=True)
        return jsonify({"error": str(e), "success": False}), 500


@app.route("/api/ai_github_push", methods=["POST"])
def api_ai_github_push():
    """رفع ملفات محددة إلى GitHub"""
    try:
        data       = request.get_json() or {}
        file_paths = data.get('files', [])   # قائمة مسارات نسبية
        commit_msg = data.get('message', '🤖 تحديث بواسطة المساعد الذكي')
        if not file_paths:
            # رفع جميع الملفات الرئيسية
            file_paths = ['app.py', 'templates/index.html', 'templates/academic.html',
                          'static/js/app.js', 'requirements.txt', 'render.yaml']
        files_to_push = {}
        for fp in file_paths:
            c, _ = _ai_read_file(fp)
            if c:
                files_to_push[fp] = c
        if not files_to_push:
            return jsonify({"error": "لم يُعثر على ملفات"}), 400
        result = _ai_github_push(files_to_push, commit_msg)
        return jsonify({"success": True, "result": result, "pushed": len(files_to_push)})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/ai_read_file", methods=["POST"])
def api_ai_read_file_endpoint():
    """قراءة ملف للمساعد الذكي"""
    try:
        data = request.get_json() or {}
        rel_path = data.get('path', '')
        if not rel_path:
            return jsonify({"error": "المسار مطلوب"}), 400
        content, err = _ai_read_file(rel_path)
        if err:
            return jsonify({"error": err}), 400
        return jsonify({"success": True, "content": content, "lines": len(content.splitlines())})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/ai_write_file", methods=["POST"])
def api_ai_write_file_endpoint():
    """كتابة/تعديل ملف عبر المساعد الذكي"""
    try:
        data = request.get_json() or {}
        rel_path = data.get('path', '')
        content  = data.get('content', '')
        if not rel_path:
            return jsonify({"error": "المسار مطلوب"}), 400
        ok, err = _ai_write_file(rel_path, content)
        if not ok:
            return jsonify({"error": err}), 400
        return jsonify({"success": True, "message": f"تم حفظ {rel_path} بنجاح"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ════════════════════════════════════════════════════════════
#  منشئ العروض التقديمية — مدموج بالكامل داخل app.py
# ════════════════════════════════════════════════════════════

# ── مجلد المخرجات ──────────────────────────────────────────
from pathlib import Path
_PPTX_OUTPUTS_DIR = Path(__file__).parent / "pptx_app" / "outputs"
_PPTX_OUTPUTS_DIR.mkdir(parents=True, exist_ok=True)
_PPTX_IMG_CACHE   = _PPTX_OUTPUTS_DIR / "img_cache"
_PPTX_IMG_CACHE.mkdir(parents=True, exist_ok=True)

# ── قوالب الألوان (templates.py) ───────────────────────────
_PPTX_THEMES = {
    "blue":   {"name":"أزرق احترافي",  "primary":(0x66,0x7E,0xEA),"secondary":(0x76,0x4B,0xA2),"accent":(0xFF,0xFF,0xFF),"text_dark":(0x1A,0x1A,0x2E),"text_light":(0xFF,0xFF,0xFF),"bg":(0xF5,0xF7,0xFF)},
    "green":  {"name":"أخضر طبيعي",    "primary":(0x11,0x99,0x55),"secondary":(0x00,0x7A,0x3D),"accent":(0xFF,0xFF,0xFF),"text_dark":(0x1A,0x2E,0x1A),"text_light":(0xFF,0xFF,0xFF),"bg":(0xF0,0xFB,0xF4)},
    "red":    {"name":"أحمر حيوي",     "primary":(0xE5,0x39,0x35),"secondary":(0xB7,0x1C,0x1C),"accent":(0xFF,0xCC,0x02),"text_dark":(0x2E,0x1A,0x1A),"text_light":(0xFF,0xFF,0xFF),"bg":(0xFD,0xF2,0xF2)},
    "purple": {"name":"بنفسجي ملكي",   "primary":(0x6A,0x1B,0x9A),"secondary":(0x4A,0x14,0x8C),"accent":(0xFF,0xD7,0x00),"text_dark":(0x1A,0x1A,0x2E),"text_light":(0xFF,0xFF,0xFF),"bg":(0xF8,0xF0,0xFF)},
}

def _pptx_get_theme(color: str) -> dict:
    return _PPTX_THEMES.get(color, _PPTX_THEMES["blue"])

# ── مساعدات pptx عامة ──────────────────────────────────────
try:
    from pptx import Presentation as _Prs
    from pptx.util import Inches as _Inches, Pt as _Pt, Emu as _Emu
    from pptx.dml.color import RGBColor as _RGBColor
    from pptx.enum.text import PP_ALIGN as _PP_ALIGN
    from pptx.oxml.ns import qn as _qn
    from lxml import etree as _etree
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

def _pptx_rgb(t):
    return _RGBColor(*t)

def _pptx_set_rtl(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1")

def _pptx_set_run_font(run, font_name: str):
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:cs", "a:ea"):
        el = rPr.find(_qn(tag))
        if el is None:
            el = _etree.SubElement(rPr, _qn(tag))
        el.set("typeface", font_name)

def _pptx_add_text_box(slide, left, top, width, height, text, size,
                       bold=False, color=(0,0,0), align=None, wrap=True,
                       font_name="Traditional Arabic"):
    if align is None:
        align = _PP_ALIGN.RIGHT
    box = slide.shapes.add_textbox(_Inches(left), _Inches(top), _Inches(width), _Inches(height))
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align; _pptx_set_rtl(p)
    run = p.add_run()
    run.text = text; run.font.size = _Pt(size); run.font.bold = bold
    run.font.color.rgb = _pptx_rgb(color); _pptx_set_run_font(run, font_name)
    return box

def _pptx_set_cell_bg(cell, color_tuple):
    tc = cell._tc; tcPr = tc.get_or_add_tcPr()
    sf = _etree.SubElement(tcPr, _qn("a:solidFill"))
    sc = _etree.SubElement(sf, _qn("a:srgbClr"))
    sc.set("val", "{:02X}{:02X}{:02X}".format(*color_tuple))

# ── جالب الصور (image_fetcher.py) ──────────────────────────
import hashlib as _hashlib

_PPTX_ARABIC_TO_EN = {
    "تعليم":"education","مدرسة":"school","جامعة":"university","طالب":"students",
    "تدريب":"training","تعلم":"learning","اختبار":"exam","تقنية":"technology",
    "ذكاء اصطناعي":"artificial intelligence","برمجة":"coding","حاسوب":"computer",
    "شبكة":"network","بيانات":"data analytics","سحابة":"cloud computing",
    "ابتكار":"innovation","ذكاء":"artificial intelligence","أعمال":"business",
    "شركة":"corporate office","تجارة":"commerce","سوق":"market","مبيعات":"sales",
    "تسويق":"marketing","إدارة":"management","قيادة":"leadership","فريق":"team",
    "اجتماع":"meeting","عمل":"workplace","موظف":"employees","استراتيجية":"strategy",
    "خطة":"planning","ميزانية":"finance","استثمار":"investment","نمو":"growth chart",
    "صحة":"healthcare","طب":"medicine","مستشفى":"hospital","رياضة":"sports",
    "بيئة":"environment","طاقة":"renewable energy","طبيعة":"nature landscape",
    "مجتمع":"community","أسرة":"family","شباب":"youth","سياحة":"tourism",
    "بناء":"architecture","مشروع":"construction project","تصميم":"design",
    "هندسة":"engineering","خلاصة":"success achievement","نتائج":"results achievement",
    "هدف":"goal target","مستقبل":"future vision","إنجاز":"achievement","نجاح":"success",
}
_PPTX_SLIDE_TYPE_DEFAULTS = {
    "title":"professional presentation","conclusion":"success achievement team",
    "table":"data analytics chart","chart":"business graph analytics","bullets":"professional business",
}

def _pptx_get_slide_image(slide_title, slide_bullets, slide_type="bullets", groq_client=None):
    keywords = ""
    if groq_client:
        try:
            content = slide_title + "\n" + "\n".join(slide_bullets[:3])
            resp = groq_client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role":"user","content":f"Extract 2-3 English keywords for image search from Arabic:\n{content}\nReturn ONLY keywords separated by commas."}],
                max_tokens=25, temperature=0.2,
            )
            kw = resp.choices[0].message.content.strip().replace("\n",",")
            keywords = ",".join([k.strip() for k in kw.split(",") if k.strip()])
        except Exception:
            pass
    if not keywords:
        full = slide_title + " " + " ".join(slide_bullets[:3])
        for ar, en in _PPTX_ARABIC_TO_EN.items():
            if ar in full:
                keywords = en.replace(" ",","); break
    if not keywords:
        keywords = _PPTX_SLIDE_TYPE_DEFAULTS.get(slide_type, "professional business")
    safe_kw   = keywords.strip()[:80]
    cache_key = _hashlib.md5(safe_kw.encode()).hexdigest()[:10]
    cache_path = _PPTX_IMG_CACHE / f"{cache_key}.jpg"
    if cache_path.exists() and cache_path.stat().st_size > 4000:
        return str(cache_path)
    try:
        import requests as _req
        resp = _req.get(f"https://loremflickr.com/900/550/{safe_kw}", timeout=8, allow_redirects=True)
        if resp.status_code == 200 and len(resp.content) > 4000:
            cache_path.write_bytes(resp.content); return str(cache_path)
    except Exception:
        pass
    return None

# ── توليد الرسوم البيانية (chart_generator.py) ─────────────
def _pptx_create_chart(chart_type, labels, values, title="", color="#667eea"):
    try:
        import matplotlib; matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import tempfile
        colors_list = ["#667eea","#764ba2","#11aa55","#e53935","#ff6b35","#00bcd4","#ffc107","#9c27b0"]
        fig, ax = plt.subplots(figsize=(8,5))
        fig.patch.set_facecolor("#f8f9ff"); ax.set_facecolor("#f8f9ff")
        if chart_type == "bar":
            bars = ax.bar(labels, values, color=colors_list[:len(labels)])
            for bar, val in zip(bars, values):
                ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+max(values)*0.01, f"{val:,.0f}", ha="center", va="bottom", fontsize=9, fontweight="bold")
            ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
        elif chart_type == "pie":
            ax.pie(values, labels=labels, autopct="%1.1f%%", colors=colors_list[:len(labels)], startangle=90, pctdistance=0.85)
        elif chart_type == "line":
            ax.plot(labels, values, color=color, linewidth=2.5, marker="o", markersize=7)
            ax.fill_between(range(len(labels)), values, alpha=0.15, color=color)
            ax.set_xticks(range(len(labels))); ax.set_xticklabels(labels)
            ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
        elif chart_type == "horizontal_bar":
            ax.barh(labels, values, color=colors_list[:len(labels)])
            ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
        if title:
            ax.set_title(title, fontsize=13, fontweight="bold", pad=12)
        plt.tight_layout()
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        plt.savefig(tmp.name, dpi=150, bbox_inches="tight", facecolor=fig.get_facecolor())
        plt.close(fig); return tmp.name
    except Exception as e:
        logger.warning(f"Chart error: {e}"); return None

# ── معالج الذكاء الاصطناعي (ai_processor.py) ───────────────
_PPTX_GROQ_MODEL = "llama-3.3-70b-versatile"

def _pptx_subtitle(ptype):
    return {"business":"عرض تجاري احترافي","educational":"مواد تعليمية متميزة","sales":"عرض تسويقي متكامل","general":"عرض تقديمي شامل"}.get(ptype,"عرض تقديمي")

def _pptx_section_titles(ptype):
    return {"business":["نظرة عامة","الأهداف الاستراتيجية","الخطة التنفيذية","الميزانية والموارد","مؤشرات النجاح"],"educational":["المقدمة","المفاهيم الأساسية","التطبيقات العملية","الأمثلة والتدريبات","التقييم"],"sales":["المشكلة","الحل المقترح","المميزات والفوائد","الأسعار والعروض","لماذا نحن؟"],"general":["المقدمة","المحتوى الرئيسي","التفاصيل","النتائج","التوصيات"]}.get(ptype,["المقدمة","المحتوى","التفاصيل","النتائج","الخلاصة"])

def _pptx_process_locally(text, num_slides, ptype, extracted_tables):
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    sentences = []
    for line in lines:
        parts = re.split(r'[.،,;؛]', line)
        sentences.extend([p.strip() for p in parts if len(p.strip()) > 5])
    title = sentences[0][:70] if sentences else "العرض التقديمي"
    slides = [{"title":title,"subtitle":_pptx_subtitle(ptype),"bullets":[],"slide_type":"title"}]
    remaining = sentences[1:] if len(sentences)>1 else ["محتوى العرض"]
    content_count = max(1, num_slides-2)
    chunk_size = max(1, len(remaining)//max(1,content_count))
    sec_titles = _pptx_section_titles(ptype)
    tables_used = 0
    for i in range(content_count):
        chunk = remaining[i*chunk_size:(i+1)*chunk_size]
        bullets = [b[:120] for b in chunk if b][:5]
        if not bullets:
            bullets = [f"النقطة الرئيسية {i+1}","التفاصيل والمعلومات","الخلاصة الجزئية"]
        slide = {"title":sec_titles[i%len(sec_titles)],"subtitle":"","bullets":bullets,"slide_type":"bullets"}
        if extracted_tables and tables_used<len(extracted_tables) and i==content_count//2:
            slide["slide_type"]="table"; slide["table_data"]=extracted_tables[tables_used]; tables_used+=1
        slides.append(slide)
    slides.append({"title":"الخلاصة والتوصيات","subtitle":"","bullets":["✅ "+(sentences[-1][:90] if sentences else "تم عرض أهم النقاط"),"📌 نرحب بأسئلتكم واستفساراتكم","🙏 شكراً لاهتمامكم"],"slide_type":"conclusion"})
    return slides[:num_slides]

class _AIProcessor:
    def __init__(self):
        self.groq_available = False; self.client = None
        try:
            from groq import Groq
            api_key = os.environ.get("GROQ_API_KEY","").strip()
            if api_key:
                self.client = Groq(api_key=api_key); self.groq_available = True
        except Exception as e:
            logger.warning(f"[Groq init] {e}")

    @property
    def is_ai_available(self): return self.groq_available

    def text_to_presentation_structure(self, text, num_slides=6, presentation_type="general",
                                       title_override="", include_tables=False, include_charts=False,
                                       extracted_tables=None):
        if self.groq_available:
            slides = self._process_with_groq(text, num_slides, presentation_type, include_tables, include_charts)
        else:
            slides = _pptx_process_locally(text, num_slides, presentation_type, extracted_tables or [])
        if title_override and slides:
            slides[0]["title"] = title_override
        return slides

    def _process_with_groq(self, text, num_slides, ptype, include_tables, include_charts):
        label = {"general":"عام","business":"تجاري","educational":"تعليمي","sales":"تسويقي"}.get(ptype,"عام")
        extras = ""
        if include_tables: extras += '\n- أضف شريحة جدول بنوع "table" مع مفتاح "table_data"'
        if include_charts: extras += '\n- أضف شريحة رسم بياني بنوع "chart" مع chart_type/chart_labels/chart_values'
        user_prompt = f"""حلّل النص وأنشئ هيكل عرض {label} من {num_slides} شرائح.{extras}

النص:
{text[:4000]}

القواعد:
- الأولى: نوع "title" مع عنوان رئيسي وعنوان فرعي
- الأخيرة: نوع "conclusion" مع 3 نقاط
- الوسطى: نوع "bullets" مع 3-5 نقاط
- كل المحتوى بالعربية
- أرجع JSON فقط: [{{"title":"...","subtitle":"...","bullets":["..."],"slide_type":"title|bullets|table|chart|conclusion",...}}]"""
        try:
            resp = self.client.chat.completions.create(
                model=_PPTX_GROQ_MODEL,
                messages=[{"role":"system","content":"أنت خبير عروض تقديمية. تُنتج JSON صحيحاً فقط."},{"role":"user","content":user_prompt}],
                max_tokens=3000, temperature=0.65,
            )
            content = resp.choices[0].message.content.strip()
            match = re.search(r'\[.*\]', content, re.DOTALL)
            if match:
                slides = json.loads(match.group())
                if isinstance(slides, list) and slides:
                    return slides
        except Exception as e:
            logger.warning(f"[Groq pptx] {e}")
        return _pptx_process_locally(text, num_slides, ptype, [])

# ── مولّد PPTX (presentation_generator.py) ─────────────────
class _PresentationGenerator:
    def __init__(self):
        self.font_name = "Traditional Arabic"
        self.body_font_size = 22

    def create_presentation(self, slides_data, theme_color="blue", cover_data=None,
                            extracted_images=None, font_name="Traditional Arabic",
                            body_font_size=22, ai_images=False, groq_client=None):
        from datetime import datetime as _dt
        self.font_name = font_name; self.body_font_size = body_font_size
        theme = _pptx_get_theme(theme_color)
        prs = _Prs()
        prs.slide_width = _Inches(13.33); prs.slide_height = _Inches(7.5)
        if cover_data:
            self._add_cover_slide(prs, cover_data, theme)
        for i, sd in enumerate(slides_data):
            stype = sd.get("slide_type","bullets")
            if stype == "title":        self._add_title_slide(prs, sd, theme)
            elif stype == "table":      self._add_table_slide(prs, sd, theme)
            elif stype == "chart":      self._add_chart_slide(prs, sd, theme)
            elif stype == "conclusion": self._add_conclusion_slide(prs, sd, theme)
            else:
                img_path = None
                if ai_images:
                    img_path = _pptx_get_slide_image(sd.get("title",""), sd.get("bullets",[]), stype, groq_client)
                elif extracted_images and i < len(extracted_images):
                    img_path = extracted_images[i]
                self._add_content_slide(prs, sd, theme, i, img_path)
        ts = _dt.now().strftime("%Y%m%d_%H%M%S")
        out = str(_PPTX_OUTPUTS_DIR / f"عرض_{ts}.pptx")
        prs.save(out); return out

    def _add_cover_slide(self, prs, cover_data, theme):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _pptx_rgb(theme["primary"])
        rect = slide.shapes.add_shape(1,_Inches(0),_Inches(0),_Inches(13.33),_Inches(2.2))
        rect.fill.solid(); rect.fill.fore_color.rgb = _pptx_rgb(theme["secondary"]); rect.line.fill.background()
        logo_box = slide.shapes.add_textbox(_Inches(0.5),_Inches(0.35),_Inches(2),_Inches(1.5))
        lp = logo_box.text_frame.paragraphs[0]; lp.alignment = _PP_ALIGN.CENTER
        lr = lp.add_run(); lr.text = cover_data.get("logo","📊"); lr.font.size = _Pt(52)
        org = cover_data.get("organization","")
        if org:
            ob = slide.shapes.add_textbox(_Inches(2.8),_Inches(0.5),_Inches(10),_Inches(1))
            op = ob.text_frame.paragraphs[0]; op.alignment = _PP_ALIGN.RIGHT; _pptx_set_rtl(op)
            or_ = op.add_run(); or_.text = org; or_.font.size = _Pt(20); or_.font.bold = True
            or_.font.color.rgb = _RGBColor(220,220,255); _pptx_set_run_font(or_, self.font_name)
        fn = self.font_name
        _pptx_add_text_box(slide,1,2.5,11.3,1.8,cover_data.get("title","العنوان الرئيسي"),42,bold=True,color=theme["text_light"],align=_PP_ALIGN.CENTER,font_name=fn)
        subtitle = cover_data.get("subtitle","")
        if subtitle:
            _pptx_add_text_box(slide,1.5,4.4,10.3,0.9,subtitle,22,color=(210,210,240),align=_PP_ALIGN.CENTER,font_name=fn)
        line = slide.shapes.add_shape(1,_Inches(3.5),_Inches(5.5),_Inches(6.3),_Emu(35000))
        line.fill.solid(); line.fill.fore_color.rgb = _pptx_rgb(theme["accent"]); line.line.fill.background()
        meta = []
        if cover_data.get("presenter"): meta.append(f"إعداد: {cover_data['presenter']}")
        if cover_data.get("date"): meta.append(cover_data["date"])
        if meta:
            _pptx_add_text_box(slide,1,5.9,11.3,0.7,"  |  ".join(meta),16,color=(200,200,230),align=_PP_ALIGN.CENTER,font_name=fn)

    def _add_title_slide(self, prs, data, theme):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _pptx_rgb(theme["primary"])
        self._add_bottom_strip(slide, theme); fn = self.font_name
        _pptx_add_text_box(slide,1,2.2,11.3,1.8,data.get("title","العنوان"),44,bold=True,color=theme["text_light"],align=_PP_ALIGN.CENTER,font_name=fn)
        subtitle = data.get("subtitle","")
        if subtitle:
            _pptx_add_text_box(slide,2,4.2,9.3,1,subtitle,24,color=(220,220,255),align=_PP_ALIGN.CENTER,font_name=fn)
        line = slide.shapes.add_shape(1,_Inches(4.5),_Inches(4.0),_Inches(4.3),_Emu(40000))
        line.fill.solid(); line.fill.fore_color.rgb = _pptx_rgb(theme["accent"]); line.line.fill.background()

    def _add_content_slide(self, prs, data, theme, index, img_path=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _pptx_rgb(theme["bg"])
        header = slide.shapes.add_shape(1,_Inches(0),_Inches(0),_Inches(13.33),_Inches(1.35))
        header.fill.solid(); header.fill.fore_color.rgb = _pptx_rgb(theme["primary"]); header.line.fill.background()
        _pptx_add_text_box(slide,0.4,0.15,12,1.05,data.get("title",""),27,bold=True,color=theme["text_light"],align=_PP_ALIGN.RIGHT,font_name=self.font_name)
        bullets = data.get("bullets",[])
        has_image = img_path and Path(img_path).exists()
        if has_image:
            content_box = slide.shapes.add_textbox(_Inches(4.8),_Inches(1.55),_Inches(8.1),_Inches(5.7))
            self._fill_bullets(content_box.text_frame, bullets, theme)
            try:
                frame = slide.shapes.add_shape(1,_Inches(0.18),_Inches(1.48),_Inches(4.44),_Inches(5.64))
                frame.fill.solid(); frame.fill.fore_color.rgb = _pptx_rgb(theme["accent"]); frame.line.fill.background()
                slide.shapes.add_picture(img_path,_Inches(0.25),_Inches(1.55),_Inches(4.3),_Inches(5.5))
            except Exception:
                content_box2 = slide.shapes.add_textbox(_Inches(0.5),_Inches(1.55),_Inches(12.3),_Inches(5.7))
                self._fill_bullets(content_box2.text_frame, bullets, theme)
        else:
            content_box = slide.shapes.add_textbox(_Inches(0.5),_Inches(1.55),_Inches(12.3),_Inches(5.7))
            self._fill_bullets(content_box.text_frame, bullets, theme)
        num_box = slide.shapes.add_textbox(_Inches(12.6),_Inches(6.9),_Inches(0.7),_Inches(0.5))
        np_ = num_box.text_frame.paragraphs[0]; np_.alignment = _PP_ALIGN.CENTER
        nr = np_.add_run(); nr.text = str(index+1); nr.font.size = _Pt(11); nr.font.color.rgb = _pptx_rgb(theme["secondary"])

    def _fill_bullets(self, tf, bullets, theme):
        tf.word_wrap = True; fs = self.body_font_size; fn = self.font_name
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.alignment = _PP_ALIGN.RIGHT; _pptx_set_rtl(p); p.space_before = _Pt(10); p.space_after = _Pt(4)
            run = p.add_run(); run.text = f"◆ {bullet}"; run.font.size = _Pt(fs)
            run.font.color.rgb = _pptx_rgb(theme["text_dark"]); _pptx_set_run_font(run, fn)

    def _add_table_slide(self, prs, data, theme):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _pptx_rgb(theme["bg"])
        header = slide.shapes.add_shape(1,_Inches(0),_Inches(0),_Inches(13.33),_Inches(1.35))
        header.fill.solid(); header.fill.fore_color.rgb = _pptx_rgb(theme["primary"]); header.line.fill.background()
        _pptx_add_text_box(slide,0.4,0.15,12,1.05,data.get("title","جدول البيانات"),27,bold=True,color=theme["text_light"],align=_PP_ALIGN.RIGHT,font_name=self.font_name)
        table_data = data.get("table_data",[]); 
        if not table_data: return
        rows = len(table_data); cols = max(len(r) for r in table_data)
        if rows==0 or cols==0: return
        table = slide.shapes.add_table(rows,cols,_Inches(0.5),_Inches(1.6),_Inches(12.3),_Inches(min(5.5,0.5+rows*0.55))).table
        fn = self.font_name
        for r_idx, row in enumerate(table_data):
            for c_idx in range(cols):
                cell = table.cell(r_idx,c_idx); cell.text = str(row[c_idx] if c_idx<len(row) else "")
                tf = cell.text_frame; par = tf.paragraphs[0]; par.alignment = _PP_ALIGN.CENTER; _pptx_set_rtl(par)
                runs = par.runs
                if runs:
                    runs[0].font.size = _Pt(15); runs[0].font.bold = (r_idx==0)
                    runs[0].font.color.rgb = _pptx_rgb(theme["text_light"] if r_idx==0 else theme["text_dark"])
                    _pptx_set_run_font(runs[0], fn)
                _pptx_set_cell_bg(cell, theme["primary"] if r_idx==0 else ((235,238,255) if r_idx%2==0 else (255,255,255)))

    def _add_chart_slide(self, prs, data, theme):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _pptx_rgb(theme["bg"])
        header = slide.shapes.add_shape(1,_Inches(0),_Inches(0),_Inches(13.33),_Inches(1.35))
        header.fill.solid(); header.fill.fore_color.rgb = _pptx_rgb(theme["primary"]); header.line.fill.background()
        _pptx_add_text_box(slide,0.4,0.15,12,1.05,data.get("title","رسم بياني"),27,bold=True,color=theme["text_light"],align=_PP_ALIGN.RIGHT,font_name=self.font_name)
        chart_type = data.get("chart_type","bar"); labels = data.get("chart_labels",[]); values = data.get("chart_values",[])
        if not labels or not values:
            _pptx_add_text_box(slide,1,3,11,2,"لا توجد بيانات كافية",20,color=theme["text_dark"],font_name=self.font_name); return
        try: values_float = [float(v) for v in values]
        except Exception: values_float = [1.0]*len(labels)
        primary_hex = "#{:02x}{:02x}{:02x}".format(*theme["primary"])
        img_path = _pptx_create_chart(chart_type, labels, values_float, data.get("chart_title",""), primary_hex)
        if img_path and Path(img_path).exists():
            slide.shapes.add_picture(img_path,_Inches(1.5),_Inches(1.6),_Inches(10.3),_Inches(5.6))
        else:
            _pptx_add_text_box(slide,1,3,11,2,"تعذّر إنشاء الرسم البياني",20,color=theme["text_dark"],font_name=self.font_name)

    def _add_conclusion_slide(self, prs, data, theme):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _pptx_rgb(theme["secondary"])
        fn = self.font_name
        _pptx_add_text_box(slide,1,1.5,11.3,1.2,data.get("title","الخلاصة"),36,bold=True,color=theme["text_light"],align=_PP_ALIGN.CENTER,font_name=fn)
        line = slide.shapes.add_shape(1,_Inches(3.5),_Inches(2.9),_Inches(6.3),_Emu(35000))
        line.fill.solid(); line.fill.fore_color.rgb = _pptx_rgb(theme["accent"]); line.line.fill.background()
        bullets = data.get("bullets",[]); cont_box = slide.shapes.add_textbox(_Inches(1.5),_Inches(3.2),_Inches(10.3),_Inches(4.0))
        tf = cont_box.text_frame; tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.alignment = _PP_ALIGN.RIGHT; _pptx_set_rtl(p); p.space_before = _Pt(12)
            run = p.add_run(); run.text = bullet; run.font.size = _Pt(self.body_font_size)
            run.font.color.rgb = _RGBColor(220,220,255); _pptx_set_run_font(run, fn)

    def _add_bottom_strip(self, slide, theme):
        rect = slide.shapes.add_shape(1,_Inches(0),_Inches(6.8),_Inches(13.33),_Inches(0.7))
        rect.fill.solid(); rect.fill.fore_color.rgb = _pptx_rgb(theme["secondary"]); rect.line.fill.background()

# ── تطبيق التصميم (design_applier.py) ──────────────────────
class _DesignApplier:
    def __init__(self): self.generator = _PresentationGenerator()
    def apply_design_to_presentation(self, presentation_path, design, slides_data):
        return self.generator.create_presentation(slides_data, theme_color=design.get("theme_color","blue"))

# ── استخراج الملفات (file_extractor.py) ────────────────────
def _extract_content(file_bytes: bytes, filename: str) -> dict:
    ext = Path(filename).suffix.lower()
    if ext in (".docx",".doc"):
        import docx as _docx, io as _io, tempfile as _tmp
        doc = _docx.Document(_io.BytesIO(file_bytes))
        result = {"text_blocks":[],"tables":[],"images":[],"full_text":""}
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text: continue
            result["text_blocks"].append({"type":"paragraph","style":para.style.name if para.style else "","text":text})
            result["full_text"] += text + "\n"
        for tbl in doc.tables:
            rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows]
            if rows: result["tables"].append(rows)
        return result
    elif ext == ".pdf":
        result = {"text_blocks":[],"tables":[],"images":[],"full_text":""}
        try:
            import pdfplumber as _pdfp, io as _io
            with _pdfp.open(_io.BytesIO(file_bytes)) as pdf:
                for pg_num, page in enumerate(pdf.pages, 1):
                    text = page.extract_text() or ""
                    if text.strip():
                        result["text_blocks"].append({"type":"page","style":f"صفحة {pg_num}","text":text.strip()})
                        result["full_text"] += text.strip() + "\n"
                    for tbl in page.extract_tables():
                        if tbl:
                            result["tables"].append([[cell or "" for cell in row] for row in tbl])
        except Exception as e:
            result["full_text"] = f"خطأ في قراءة PDF: {str(e)}"
        return result
    else:
        text = file_bytes.decode("utf-8", errors="ignore")
        return {"text_blocks":[{"type":"paragraph","style":"Normal","text":text}],"tables":[],"images":[],"full_text":text}

# ── محوّل HTML إلى PPTX (html_to_pptx.py) ──────────────────
def _parse_html_slides(html: str):
    from bs4 import BeautifulSoup as _BS
    import re as _re
    soup = _BS(html, "lxml")
    for tag in soup.find_all(["script","noscript"]): tag.decompose()
    slide_cands = (soup.find_all("section") or soup.find_all(class_=_re.compile(r"slide|page|frame",_re.I))
                   or soup.find_all(attrs={"data-slide":True}) or soup.find_all("article"))
    if not slide_cands:
        sections = []; current = None
        for el in (soup.body or soup).children:
            if not hasattr(el,"name"): continue
            if el.name in ("h1","h2","h3"):
                if current: sections.append(current)
                current = [el]
            elif current is not None: current.append(el)
        if current: sections.append(current)
        slide_cands = []
        for group in sections:
            wrapper = _BS("<div></div>","lxml").div
            for el in group: wrapper.append(el.__copy__())
            slide_cands.append(wrapper)
    def _extract(el):
        sd = {"title":"","subtitle":"","bullets":[],"slide_type":"bullets","bg_color":None,"text_color":None,"table_data":None,"is_title_slide":False}
        for htag in ["h1","h2","h3","h4"]:
            h = el.find(htag)
            if h: sd["title"] = h.get_text(strip=True); sd["is_title_slide"] = (htag=="h1"); break
        bullets = [li.get_text(strip=True) for li in el.find_all("li") if li.get_text(strip=True)]
        if not bullets:
            bullets = [p.get_text(strip=True) for p in el.find_all("p") if p.get_text(strip=True) and p.get_text(strip=True) not in (sd["title"],sd["subtitle"]) and len(p.get_text(strip=True))>3]
        sd["bullets"] = bullets[:6]
        tbl = el.find("table")
        if tbl:
            rows = [[td.get_text(strip=True) for td in tr.find_all(["td","th"])] for tr in tbl.find_all("tr") if tr.find_all(["td","th"])]
            if rows: sd["table_data"]=rows; sd["slide_type"]="table"
        classes = " ".join(el.get("class",[]))
        if _re.search(r"title|cover|first|intro",classes,_re.I) or sd["is_title_slide"]: sd["slide_type"]="title"
        elif _re.search(r"end|outro|thank|conclusion",classes,_re.I): sd["slide_type"]="conclusion"
        return sd
    slides = [_extract(el) for el in slide_cands if _extract(el).get("title") or _extract(el).get("bullets")]
    if not slides: slides = [_extract(soup.body or soup)]
    return slides

def _html_to_pptx(html: str, override_title: str = "") -> str:
    from datetime import datetime as _dt
    _HTP_THEME = {"primary":(102,126,234),"secondary":(118,75,162),"accent":(255,255,255),"text_dark":(26,26,46),"text_light":(255,255,255),"bg":(245,247,255)}
    slides_data = _parse_html_slides(html)
    if not slides_data: raise ValueError("لم يتم العثور على محتوى في HTML")
    if override_title and slides_data: slides_data[0]["title"] = override_title
    prs = _Prs(); prs.slide_width = _Inches(13.33); prs.slide_height = _Inches(7.5)
    gen = _PresentationGenerator()
    for i, sd in enumerate(slides_data):
        stype = sd.get("slide_type","bullets")
        if stype=="title": gen._add_title_slide(prs,sd,_HTP_THEME)
        elif stype=="table" and sd.get("table_data"): gen._add_table_slide(prs,sd,_HTP_THEME)
        elif stype=="conclusion": gen._add_conclusion_slide(prs,sd,_HTP_THEME)
        else: gen._add_content_slide(prs,sd,_HTP_THEME,i)
    ts = _dt.now().strftime("%Y%m%d_%H%M%S")
    out = str(_PPTX_OUTPUTS_DIR / f"html_presentation_{ts}.pptx")
    prs.save(out); return out

# ── تهيئة الكائنات ──────────────────────────────────────────
try:
    _pptx_ai  = _AIProcessor()
    _pptx_gen = _PresentationGenerator()
    _pptx_dap = _DesignApplier()
    logger.info("✅ PowerPoint modules loaded successfully (inlined)")
except Exception as _pptx_err:
    logger.error(f"❌ Failed to init PPTX: {_pptx_err}")
    _pptx_ai = _pptx_gen = _pptx_dap = None


@app.route('/tools/pptx/ai_status')
def tools_pptx_ai_status():
    available = _pptx_ai is not None and _pptx_ai.is_ai_available
    return jsonify({'available': available})


@app.route('/tools/pptx/extract_file', methods=['POST'])
def tools_pptx_extract_file():
    if 'file' not in request.files:
        return jsonify({'error': 'لا يوجد ملف'}), 400
    f = request.files['file']
    try:
        result = _extract_content(f.read(), f.filename)
        return jsonify({
            'success': True,
            'full_text': result['full_text'],
            'tables': result.get('tables', []),
            'word_count': len(result['full_text'].split()),
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/tools/pptx/generate', methods=['POST'])
def tools_pptx_generate():
    if _pptx_ai is None or _pptx_gen is None:
        return jsonify({'error': 'وحدات PowerPoint غير متاحة — تحقق من سجلات الخادم'}), 500
    try:
        data = request.json or {}
        text           = data.get('text', '')
        num_slides     = int(data.get('num_slides', 6))
        ptype          = data.get('ptype', 'general')
        title          = data.get('title', '')
        inc_tables     = bool(data.get('inc_tables', False))
        inc_charts     = bool(data.get('inc_charts', False))
        theme_color    = data.get('theme_color', 'blue')
        font_name      = data.get('font_name', 'Traditional Arabic')
        body_font_size = int(data.get('body_font_size', 22))
        cover_data     = data.get('cover_data', None)
        ext_tables     = data.get('extracted_tables', [])

        slides = _pptx_ai.text_to_presentation_structure(
            text=text, num_slides=num_slides, presentation_type=ptype,
            title_override=title, include_tables=inc_tables,
            include_charts=inc_charts, extracted_tables=ext_tables,
        )

        if inc_tables and ext_tables:
            if not any(s.get('slide_type') == 'table' for s in slides):
                for i, tbl in enumerate(ext_tables[:2]):
                    slides.insert(min(2 + i, len(slides) - 1),
                        {'title': f'جدول البيانات {i+1}', 'slide_type': 'table',
                         'table_data': tbl, 'bullets': []})

        if inc_charts and ext_tables:
            from modules.chart_generator import parse_table_for_chart
            if not any(s.get('slide_type') == 'chart' for s in slides):
                for tbl in ext_tables:
                    info = parse_table_for_chart(tbl)
                    if info:
                        slides.insert(min(3, len(slides) - 1),
                            {'title': 'تحليل البيانات', 'slide_type': 'chart',
                             'chart_type': 'bar', 'chart_labels': info['labels'],
                             'chart_values': info['values'],
                             'chart_title': info.get('title', ''), 'bullets': []})
                        break

        groq_client = _pptx_ai.client if (_pptx_ai and _pptx_ai.is_ai_available) else None
        out = _pptx_gen.create_presentation(
            slides_data=slides, theme_color=theme_color,
            cover_data=cover_data, extracted_images=[],
            font_name=font_name, body_font_size=body_font_size,
            ai_images=False, groq_client=groq_client,
        )
        filename = os.path.basename(out)
        return jsonify({'success': True, 'filename': filename, 'slides': slides})
    except Exception as e:
        import traceback
        logger.error(f"PPTX generate error: {e}")
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500


@app.route('/tools/pptx/download/<path:filename>')
def tools_pptx_download(filename):
    outputs_dir = os.path.join(os.path.dirname(__file__), 'pptx_app', 'outputs')
    return send_from_directory(outputs_dir, filename, as_attachment=True)


@app.route('/download/excel/<path:filename>')
def download_excel(filename):
    """تحميل ملفات Excel من مجلد الإخراج"""
    outputs_dir = os.path.join(os.path.dirname(__file__), 'pptx_app', 'outputs')
    file_path = os.path.join(outputs_dir, filename)
    if not os.path.exists(file_path):
        return jsonify({'error': 'الملف غير موجود'}), 404
    return send_file(
        file_path,
        as_attachment=True,
        download_name=filename,
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )


@app.route('/download/excel-list')
def download_excel_list():
    """قائمة ملفات Excel المتاحة للتحميل"""
    outputs_dir = os.path.join(os.path.dirname(__file__), 'pptx_app', 'outputs')
    files = []
    for f in os.listdir(outputs_dir):
        if f.endswith('.xlsx') or f.endswith('.xls'):
            files.append({'name': f, 'url': f'/download/excel/{f}'})
    return jsonify({'files': files})


@app.route('/stats1208')
@app.route('/احصاء1208')
def stats1208_page():
    """صفحة حل اختبار احصاء 1208 التفاعلية"""
    return render_template('stats1208.html')


@app.route('/stats1208.xlsx')
@app.route('/get-stats-excel')
def download_stats1208():
    """تحميل ملف Excel – احصاء 1208 – رابط مباشر"""
    f = os.path.join(os.path.dirname(__file__), 'static', 'stats1208_solution.xlsx')
    if not os.path.exists(f):
        # حاول الإنشاء تلقائياً إن لم يوجد
        try:
            import subprocess, sys
            subprocess.run([sys.executable, 'create_excel_solution.py'], check=True)
            import shutil
            src = os.path.join(os.path.dirname(__file__), 'pptx_app', 'outputs',
                               'احصاء_1208_الحلول_الكاملة.xlsx')
            shutil.copy(src, f)
        except Exception as ex:
            return jsonify({'error': str(ex)}), 500
    return send_file(
        f,
        as_attachment=True,
        download_name='احصاء_1208_الحلول_الكاملة.xlsx',
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )


@app.route('/tools/pptx/chat', methods=['POST'])
def tools_pptx_chat():
    try:
        data     = request.json or {}
        messages = data.get('messages', [])
        groq_key = os.environ.get('GROQ_API_KEY', '').strip()
        if not groq_key:
            return jsonify({'reply': '⚠️ مفتاح GROQ_API_KEY غير موجود في إعدادات الأسرار. يمكنك الحصول على مفتاح مجاني من console.groq.com'})
        from groq import Groq as _Groq
        client = _Groq(api_key=groq_key)
        history = [{'role': 'system', 'content':
            'أنت مساعد ذكي متخصص في إنشاء العروض التقديمية PowerPoint والمحتوى الأكاديمي. أجب باللغة العربية دائماً.'}
        ] + messages
        resp = client.chat.completions.create(
            model='llama-3.3-70b-versatile', messages=history, max_tokens=1500)
        return jsonify({'reply': resp.choices[0].message.content})
    except Exception as e:
        return jsonify({'reply': f'❌ خطأ: {e}'})


@app.route("/academic")
def academic_analysis():
    _groq_key = os.environ.get('GROQ_API_KEY', '').strip()
    resp = make_response(render_template('academic.html', groq_key=_groq_key))
    resp.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate, max-age=0'
    return resp


alert_queue.start()
# تهيئة الحلقة المشتركة مبكراً قبل أي اتصالات gevent
_ensure_shared_login_loop()
load_all_sessions()

def _auto_resume_persistent_tasks():
    def worker():
        time.sleep(3)
        logger.info("🔁 فحص المهام الدائمة لإعادة تشغيلها تلقائياً...")
        with USERS_LOCK:
            user_ids = list(USERS.keys())
        resumed = 0
        for uid in user_ids:
            try:
                settings = load_settings(uid)
                want_monitor = bool(settings.get('monitoring_persistent', False))
                want_rotating = bool(settings.get('rotating_persistent', False))
                if not (want_monitor or want_rotating):
                    continue
                ok = telegram_manager.ensure_client_active(uid)
                if not ok:
                    logger.warning(f"⏭️  لا يمكن استئناف مهام {uid}: العميل غير متاح/غير موثق")
                    continue
                with USERS_LOCK:
                    if uid in USERS and not USERS[uid].get('client_manager'):
                        USERS[uid]['client_manager'] = telegram_manager.get_client_manager(uid)
                if want_monitor:
                    with USERS_LOCK:
                        already = USERS.get(uid, {}).get('is_running', False)
                        if not already:
                            USERS[uid]['is_running'] = True
                    if not already:
                        t = _OSThread(target=monitoring_worker, args=(uid,), daemon=True)
                        t.start()
                        with USERS_LOCK:
                            if uid in USERS:
                                USERS[uid]['thread'] = t
                        logger.info(f"♻️  استُؤنفت المراقبة للحساب {uid}")
                        resumed += 1
                        try:
                            socketio.emit('log_update', {
                                "message": "♻️ تم استئناف المراقبة تلقائياً (مهمة دائمة)"
                            }, to=uid)
                            socketio.emit('monitoring_status', {
                                "monitoring_active": True, "status": "running", "is_running": True
                            }, to=uid)
                            socketio.emit('update_monitoring_buttons', {"is_running": True}, to=uid)
                        except Exception:
                            pass
                if want_rotating:
                    msgs = settings.get('rotating_messages', [])
                    grps = dedupe_groups(settings.get('rotating_groups', []))
                    interval = int(settings.get('rotating_interval', 5))
                    valid_msgs = [m for m in msgs if m and m.strip()]
                    if grps and valid_msgs:
                        is_alive = (uid in rotating_manager.threads
                                    and rotating_manager.threads[uid]
                                    and rotating_manager.threads[uid].is_alive())
                        if not is_alive:
                            def _cb(u, status, group, info):
                                if status == 'success':
                                    socketio.emit('log_update', {"message": f"🔄 [متسلسل] أرسل إلى {group}"}, to=u)
                                else:
                                    socketio.emit('log_update', {"message": f"❌ [متسلسل] فشل إلى {group}: {info}"}, to=u)
                            rotating_manager.start(uid, grps, valid_msgs, interval, _cb)
                            logger.info(f"♻️  استُؤنف الإرسال المتسلسل للحساب {uid}")
                            resumed += 1
                            try:
                                socketio.emit('log_update', {
                                    "message": f"♻️ تم استئناف الإرسال المتسلسل تلقائياً ({len(valid_msgs)} رسائل) كل {interval} دقيقة"
                                }, to=uid)
                            except Exception:
                                pass
                    else:
                        logger.info(f"⏭️  تخطي استئناف الإرسال المتسلسل لـ {uid}: لا توجد رسائل/مجموعات")
            except Exception as e:
                logger.error(f"خطأ أثناء استئناف مهام {uid}: {e}")
        if resumed:
            logger.info(f"✅ تم استئناف {resumed} مهمة دائمة")
        else:
            logger.info("ℹ️  لا توجد مهام دائمة لاستئنافها")
    _OSThread(target=worker, daemon=True, name="AutoResumeTasks").start()

_auto_resume_persistent_tasks()

# ════════════════════════════════════════════════════════════
#  مزامنة الجلسات مع GitHub
# ════════════════════════════════════════════════════════════
def github_upload_session(user_id):
    """رفع ملف الجلسة إلى GitHub للحفظ الاحتياطي"""
    try:
        session_file = os.path.join(SESSIONS_DIR, f"{user_id}_session.session")
        if not os.path.exists(session_file):
            return False
        with open(session_file, 'rb') as f:
            content = f.read()
        import base64 as _b64
        headers = {
            'Authorization': f'token {GITHUB_TOKEN}',
            'Accept': 'application/vnd.github.v3+json',
            'Content-Type': 'application/json',
        }
        rel_path = f"sessions/{user_id}_session.session"
        get_url = f"https://api.github.com/repos/{GITHUB_REPO}/contents/{rel_path}"
        gr = requests.get(get_url, headers=headers, params={'ref': GITHUB_BRANCH}, timeout=15)
        sha = gr.json().get('sha') if gr.ok else None
        body = {
            'message': f'Upload session for {user_id}',
            'content': _b64.b64encode(content).decode(),
            'branch': GITHUB_BRANCH,
        }
        if sha:
            body['sha'] = sha
        pr = requests.put(get_url, headers=headers, json=body, timeout=30)
        if pr.ok:
            logger.info(f"Session for {user_id} uploaded to GitHub")
            return True
        else:
            logger.warning(f"Failed to upload session: {pr.status_code}")
            return False
    except Exception as e:
        logger.error(f"github_upload_session error: {e}")
        return False

def github_delete_session(user_id):
    """حذف ملف الجلسة الملغاة من GitHub"""
    try:
        headers = {
            'Authorization': f'token {GITHUB_TOKEN}',
            'Accept': 'application/vnd.github.v3+json',
            'Content-Type': 'application/json',
        }
        rel_path = f"sessions/{user_id}_session.session"
        get_url = f"https://api.github.com/repos/{GITHUB_REPO}/contents/{rel_path}"
        gr = requests.get(get_url, headers=headers, params={'ref': GITHUB_BRANCH}, timeout=15)
        if not gr.ok:
            return False
        sha = gr.json().get('sha')
        if not sha:
            return False
        body = {
            'message': f'Delete revoked session for {user_id}',
            'sha': sha,
            'branch': GITHUB_BRANCH,
        }
        dr = requests.delete(get_url, headers=headers, json=body, timeout=30)
        if dr.ok:
            logger.info(f"Session for {user_id} deleted from GitHub")
            return True
        return False
    except Exception as e:
        logger.error(f"github_delete_session error: {e}")
        return False

# ════════════════════════════════════════════════════════════
#  مدقق صحة الجلسات الدوري
# ════════════════════════════════════════════════════════════
def start_session_health_checker():
    """يفحص صحة جميع الجلسات كل 120 ثانية"""
    def _checker():
        import time as _t
        _t.sleep(120)  # انتظر دقيقتين قبل الفحص الأول لإعطاء العملاء وقتاً كافياً
        while True:
            try:
                with USERS_LOCK:
                    user_ids = list(USERS.keys())
                for uid in user_ids:
                    try:
                        with USERS_LOCK:
                            user_data = USERS.get(uid, {})
                        client_manager = user_data.get('client_manager')
                        if not client_manager or not client_manager.client:
                            continue
                        # تخطّ الفحص إذا كان الخيط لا يزال يبدأ
                        if client_manager.thread and client_manager.thread.is_alive() and not client_manager.is_ready.is_set():
                            logger.debug(f"Health check skipped for {uid} — client still starting")
                            continue
                        # تخطّ الفحص إذا لم يكن الـ loop يعمل بعد
                        if not client_manager.loop or not client_manager.loop.is_running():
                            continue
                        is_valid = client_manager.check_session_valid_sync()
                        if not is_valid:
                            logger.warning(f"Session health check failed for {uid} — revoking")
                            socketio.emit('session_revoked', {
                                "user_id": uid,
                                "reason": "انتهت صلاحية الجلسة أو تم إلغاؤها"
                            }, to=uid)
                            github_delete_session(uid)
                    except Exception as inner_e:
                        logger.debug(f"Health check error for {uid}: {inner_e}")
            except Exception as e:
                logger.debug(f"Session health checker error: {e}")
            _t.sleep(120)

    # استخدام OS thread حقيقي للـ health checker
    t = _OSThread(target=_checker, daemon=True, name='SessionHealthChecker')
    t.start()
    logger.info("✅ مدقق صحة الجلسات الدوري مُفعّل (كل 120 ثانية)")

start_session_health_checker()

# ════════════════════════════════════════════════════════════
#  نقاط API — فحص الجلسة وإعادة التعيين
# ════════════════════════════════════════════════════════════
@app.route("/api/check_session_valid", methods=["GET"])
def api_check_session_valid():
    user_id = session.get('user_id', 'user_1')
    try:
        with USERS_LOCK:
            user_data = USERS.get(user_id, {})
        client_manager = user_data.get('client_manager')
        if not client_manager or not client_manager.client:
            return jsonify({"success": True, "valid": False, "reason": "العميل غير متصل"})
        is_valid = client_manager.check_session_valid_sync()
        return jsonify({"success": True, "valid": is_valid})
    except Exception as e:
        return jsonify({"success": False, "valid": False, "reason": str(e)})

@app.route("/api/force_reset_session", methods=["POST"])
def api_force_reset_session():
    user_id = session.get('user_id', 'user_1')
    try:
        with USERS_LOCK:
            user_data = USERS.get(user_id, {})
        client_manager = user_data.get('client_manager')
        if client_manager:
            client_manager.run_coroutine(client_manager.force_reset_session())
        github_delete_session(user_id)
        clear_user_session(user_id)
        with USERS_LOCK:
            if user_id in USERS:
                USERS[user_id]['connected'] = False
                USERS[user_id]['authenticated'] = False
                USERS[user_id]['client_manager'] = None
        return jsonify({"success": True, "message": "تم إعادة تعيين الجلسة بنجاح"})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/auto_join/settings", methods=["GET", "POST"])
def api_auto_join_settings():
    user_id = session.get('user_id', 'user_1')
    if request.method == 'POST':
        try:
            data = request.json or {}
            settings = load_settings(user_id)
            if 'links' in data:
                settings['auto_join_links'] = data['links']
            if 'delay' in data:
                settings['auto_join_delay'] = int(data['delay'])
            if 'max_retries' in data:
                settings['auto_join_max_retries'] = int(data['max_retries'])
            save_settings(user_id, settings)
            return jsonify({"success": True})
        except Exception as e:
            return jsonify({"success": False, "message": str(e)})
    settings = load_settings(user_id)
    return jsonify({
        "success": True,
        "links": settings.get('auto_join_links', []),
        "delay": settings.get('auto_join_delay', 3),
        "max_retries": settings.get('auto_join_max_retries', 1)
    })

@app.route("/api/auto_join/stop", methods=["POST"])
def api_auto_join_stop():
    user_id = session.get('user_id', 'user_1')
    try:
        stopped = False
        with USERS_LOCK:
            stop_event = USERS.get(user_id, {}).get('auto_join_stop')
            if stop_event and not stop_event.is_set():
                stop_event.set()
                stopped = True
        msg = "⏹ تم إيقاف الانضمام التلقائي" if stopped else "ℹ️ لا يوجد انضمام نشط حالياً"
        socketio.emit('log_update', {"message": msg}, to=user_id)
        return jsonify({"success": True, "message": msg, "stopped": stopped})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/auto_join/status", methods=["GET"])
def api_auto_join_status():
    """يُرجع الحالة الجارية للانضمام المتقدم (للاستعادة عند تبديل الحسابات)"""
    user_id = session.get('user_id', 'user_1')
    with USERS_LOCK:
        state = USERS.get(user_id, {}).get('auto_join_state', {})
    return jsonify({"success": True, "state": state})

# ════════════════════════════════════════════════════════════
#  نظام الاستمرارية الدائم — يمنع توقف التطبيق تلقائياً
# ════════════════════════════════════════════════════════════
def _start_keepalive():
    """
    يُرسل ping لنفسه كل 4 دقائق لمنع السكون (خاصة على Render Free).
    يعمل كخيط daemon لا يمنع إيقاف التطبيق يدوياً.
    """
    import time as _time
    import socket as _socket

    def _ping_self():
        _time.sleep(30)  # انتظر حتى يكتمل الإقلاع
        _port = int(os.environ.get('PORT', 5000))
        while True:
            try:
                _socket.setdefaulttimeout(10)
                conn = _socket.create_connection(('127.0.0.1', _port), timeout=5)
                conn.close()
            except Exception:
                pass
            _time.sleep(240)  # كل 4 دقائق

    t = _OSThread(target=_ping_self, daemon=True, name='KeepAlive')
    t.start()
    logger.info("🔄 نظام الاستمرارية الدائم مُفعّل (ping كل 4 دقائق)")

_start_keepalive()

# ── نقطة keepalive يمكن استدعاؤها خارجياً (من UptimeRobot مثلاً) ──
@app.route('/keepalive')
@app.route('/ping')
def route_keepalive():
    return jsonify({"status": "alive", "time": datetime.utcnow().isoformat()})


# ════════════════════════════════════════════════════════════
#  إعادة التعيين — مستخدم واحد أو جميع المستخدمين
# ════════════════════════════════════════════════════════════
def _do_reset_user(uid):
    """تنفيذ إعادة تعيين مستخدم واحد (يُستدعى داخلياً)"""
    try:
        with USERS_LOCK:
            user_data = USERS.get(uid, {})
        client_manager = user_data.get('client_manager')
        if client_manager:
            try:
                if client_manager.client and client_manager.loop and client_manager.client.is_connected():
                    future = asyncio.run_coroutine_threadsafe(
                        client_manager.client.log_out(), client_manager.loop
                    )
                    future.result(timeout=8)
            except Exception:
                pass
            try:
                client_manager.stop_flag.set()
            except Exception:
                pass
        telegram_manager.login_managers.pop(uid, None)
        telegram_manager.client_managers.pop(uid, None)
        github_delete_session(uid)
        clear_user_session(uid)
        with USERS_LOCK:
            USERS.pop(uid, None)
    except Exception as e:
        logger.warning(f"_do_reset_user({uid}): {e}")


@app.route("/api/reset_user", methods=["POST"])
def api_reset_user():
    """إعادة تعيين الحساب الحالي فقط — حذف الجلسة محلياً وعلى GitHub"""
    if 'user_id' not in session:
        return jsonify({"success": False, "message": "الجلسة غير صالحة"})
    user_id = session['user_id']
    data = request.get_json(silent=True) or {}
    target_id = data.get('target_user_id', user_id)
    try:
        _do_reset_user(target_id)
        socketio.emit('login_status', {
            "logged_in": False, "connected": False,
            "awaiting_code": False, "awaiting_password": False, "is_running": False,
        }, to=target_id)
        socketio.emit('connection_status', {"status": "disconnected"}, to=target_id)
        socketio.emit('user_reset_done', {"user_id": target_id}, to=target_id)
        return jsonify({"success": True, "message": f"✅ تم إعادة تعيين {target_id} بالكامل"})
    except Exception as e:
        logger.error(f"api_reset_user error for {user_id}: {e}")
        return jsonify({"success": False, "message": str(e)})


@app.route("/api/reset_all", methods=["POST"])
def api_reset_all():
    """إعادة تعيين جميع المستخدمين والجلسات والإعدادات — يُعيد النظام لحالته الأولى"""
    errors = []
    for uid in list(PREDEFINED_USERS.keys()):
        try:
            _do_reset_user(uid)
        except Exception as e:
            errors.append(str(e))
    socketio.emit('force_full_reset', {"message": "تم إعادة تعيين جميع الجلسات والإعدادات"})
    if errors:
        return jsonify({"success": True,
                        "message": f"✅ تمت الإعادة مع ملاحظات: {'; '.join(errors)}"})
    return jsonify({"success": True,
                    "message": "✅ تم إعادة تعيين جميع الجلسات والإعدادات بالكامل"})


@app.route("/api/health", methods=["GET"])
def api_health_status():
    try:
        users_status = {}
        with USERS_LOCK:
            for uid, data in USERS.items():
                cm = data.get('client_manager')
                loop_ok = False
                if cm:
                    lp = getattr(cm, 'loop', None)
                    loop_ok = bool(lp and lp.is_running())

                rot = rotating_manager
                rot_active = (uid in rot.threads and rot.threads[uid] and rot.threads[uid].is_alive())
                next_send_in = None
                if rot_active and uid in rot.next_send_at:
                    remaining = int(rot.next_send_at[uid] - time.time())
                    next_send_in = max(0, remaining)

                users_status[uid] = {
                    "name": PREDEFINED_USERS.get(uid, {}).get("name", uid),
                    "authenticated": data.get("authenticated", False),
                    "connected": data.get("connected", False),
                    "is_running": data.get("is_running", False),
                    "awaiting_code": data.get("awaiting_code", False),
                    "awaiting_password": data.get("awaiting_password", False),
                    "client_loop_running": loop_ok,
                    "rotating_send_active": rot_active,
                    "rotating_next_send_in_seconds": next_send_in,
                    "monitored_keywords": len(cm.monitored_keywords) if cm else 0,
                    "monitored_groups": len(cm.monitored_groups) if cm else 0,
                }

        return jsonify({
            "success": True,
            "timestamp": time.strftime("%Y-%m-%d %H:%M:%S"),
            "users": users_status
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500


# ============================================
# قسم مراقبة الروابط
# ============================================

LINK_MONITORS = {}
LINK_MONITOR_WORKERS = {}

def extract_links_from_text(text):
    """استخراج جميع الروابط من النص"""
    url_pattern = r'https?://[^\s]+|t\.me/[^\s]+|telegram\.me/[^\s]+'
    links = re.findall(url_pattern, text)
    clean_links = []
    for link in links:
        link = re.sub(r'[.,;:!?)]+$', '', link)
        if link and link not in clean_links:
            clean_links.append(link)
    return clean_links

@app.route('/api/link_monitor/start', methods=['POST'])
def start_link_monitor():
    """بدء مراقبة الروابط"""
    user_id = session.get('user_id')
    if user_id not in USERS:
        return jsonify({'error': 'المستخدم غير موجود'}), 404

    data = request.json
    monitor_all = data.get('monitor_all', True)
    specific_chats = data.get('specific_chats', '').strip().split('\n')
    send_to_saved = data.get('send_to_saved', True)

    if user_id in LINK_MONITOR_WORKERS and LINK_MONITOR_WORKERS[user_id].is_alive():
        return jsonify({'error': 'توجد عملية مراقبة جارية بالفعل'}), 400

    if specific_chats and specific_chats[0]:
        specific_chats = [chat.strip() for chat in specific_chats if chat.strip()]
    else:
        specific_chats = []

    LINK_MONITORS[user_id] = {
        'is_running': True,
        'monitor_all': monitor_all,
        'specific_chats': specific_chats,
        'send_to_saved': send_to_saved,
        'links_found': [],
        'start_time': datetime.now().isoformat(),
        'total_links': 0,
        'monitored_chats': 0
    }

    worker = threading.Thread(target=run_link_monitor_worker, args=(user_id,))
    worker.daemon = True
    worker.start()
    LINK_MONITOR_WORKERS[user_id] = worker

    return jsonify({'success': True, 'message': 'بدأت مراقبة الروابط'})

def run_link_monitor_worker(user_id):
    """دالة تعمل في خيط منفصل لمراقبة الروابط"""
    try:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        loop.run_until_complete(link_monitor_worker(user_id))
        loop.close()
    except Exception as e:
        logger.error(f"خطأ في worker مراقبة الروابط: {str(e)}")
        socketio.emit('link_monitor_error', {'error': str(e)}, room=user_id)

async def link_monitor_worker(user_id):
    """الدالة الأساسية لمراقبة الروابط"""
    state = LINK_MONITORS[user_id]
    client = USERS[user_id]['client']

    try:
        from telethon import events as telethon_events

        dialogs = await client.get_dialogs()
        state['monitored_chats'] = len(dialogs)

        if state['monitor_all']:
            chats_to_monitor = None  # None = جميع الدردشات
        else:
            chats_to_monitor = []
            for dialog in dialogs:
                name = getattr(dialog.entity, 'title', None) or getattr(dialog.entity, 'first_name', '') or ''
                if any(s in name for s in state['specific_chats']):
                    chats_to_monitor.append(dialog.entity)

        logger.info(f"بدء مراقبة الروابط — {state['monitored_chats']} دردشة")
        socketio.emit('link_monitor_status_update', {
            'monitored_chats': state['monitored_chats'],
            'is_running': True
        }, room=user_id)

        @client.on(telethon_events.NewMessage(chats=chats_to_monitor))
        async def handle_new_message(event):
            if not state.get('is_running'):
                return
            try:
                message = event.message
                if not message or not message.text:
                    return

                links = extract_links_from_text(message.text)
                if not links:
                    return

                chat = await event.get_chat()
                chat_name = getattr(chat, 'title', None) or getattr(chat, 'first_name', None) or 'Unknown'

                sender = await event.get_sender()
                sender_name = getattr(sender, 'first_name', None) or 'Unknown'

                for link in links:
                    link_data = {
                        'link': link,
                        'chat_name': chat_name,
                        'chat_id': getattr(chat, 'id', 0),
                        'sender_name': sender_name,
                        'message_text': message.text[:200] + ('...' if len(message.text) > 200 else ''),
                        'message_id': message.id,
                        'date': message.date.isoformat() if message.date else datetime.now().isoformat(),
                        'timestamp': datetime.now().isoformat()
                    }

                    state['links_found'].append(link_data)
                    state['total_links'] += 1

                    if state.get('send_to_saved'):
                        try:
                            saved_msg = (
                                f"🔗 رابط جديد!\n\n"
                                f"📌 الرابط: {link}\n"
                                f"💬 الدردشة: {chat_name}\n"
                                f"👤 المرسل: {sender_name}\n"
                                f"📝 النص: {message.text[:100]}\n\n"
                                f"🕐 {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
                            )
                            await client.send_message('me', saved_msg)
                            logger.info(f"تم إرسال الرابط إلى المحفوظات: {link}")
                        except Exception as e:
                            logger.error(f"فشل الإرسال إلى المحفوظات: {str(e)}")

                    socketio.emit('link_monitor_update', {
                        'link': link_data,
                        'total': state['total_links']
                    }, room=user_id)

                    logger.info(f"رابط مكتشف في [{chat_name}]: {link}")

            except Exception as e:
                logger.error(f"خطأ في معالجة رسالة: {str(e)}")

        while state.get('is_running'):
            await asyncio.sleep(1)

        client.remove_event_handler(handle_new_message)

    except Exception as e:
        logger.error(f"خطأ في مراقبة الروابط: {str(e)}")
        socketio.emit('link_monitor_error', {'error': str(e)}, room=user_id)
    finally:
        state['is_running'] = False
        socketio.emit('link_monitor_done', {
            'total_links': state.get('total_links', 0),
            'monitored_chats': state.get('monitored_chats', 0)
        }, room=user_id)
        if user_id in LINK_MONITOR_WORKERS:
            del LINK_MONITOR_WORKERS[user_id]

@app.route('/api/link_monitor/stop')
def stop_link_monitor():
    """إيقاف مراقبة الروابط"""
    user_id = session.get('user_id')
    if user_id not in USERS:
        return jsonify({'error': 'المستخدم غير موجود'}), 404
    if user_id in LINK_MONITORS:
        LINK_MONITORS[user_id]['is_running'] = False
    return jsonify({'success': True, 'message': 'تم إيقاف مراقبة الروابط'})

@app.route('/api/link_monitor/status')
def link_monitor_status():
    """الحصول على حالة مراقبة الروابط"""
    user_id = session.get('user_id')
    if user_id not in USERS:
        return jsonify({'error': 'المستخدم غير موجود'}), 404
    if user_id in LINK_MONITORS:
        state = LINK_MONITORS[user_id]
        is_running = user_id in LINK_MONITOR_WORKERS and LINK_MONITOR_WORKERS[user_id].is_alive()
        return jsonify({
            'is_running': is_running,
            'total_links': state.get('total_links', 0),
            'monitored_chats': state.get('monitored_chats', 0),
            'links_found': state.get('links_found', [])[-50:]
        })
    return jsonify({'is_running': False, 'total_links': 0, 'monitored_chats': 0, 'links_found': []})

@app.route('/api/link_monitor/links')
def get_monitored_links():
    """الحصول على قائمة الروابط المكتشفة"""
    user_id = session.get('user_id')
    if user_id not in USERS:
        return jsonify({'error': 'المستخدم غير موجود'}), 404
    limit = int(request.args.get('limit', 100))
    if user_id in LINK_MONITORS:
        links = LINK_MONITORS[user_id].get('links_found', [])
        links_sorted = sorted(links, key=lambda x: x.get('timestamp', ''), reverse=True)
        return jsonify({'links': links_sorted[:limit], 'total': len(links)})
    return jsonify({'links': [], 'total': 0})

@app.route('/api/link_monitor/clear')
def clear_monitored_links():
    """مسح قائمة الروابط"""
    user_id = session.get('user_id')
    if user_id not in USERS:
        return jsonify({'error': 'المستخدم غير موجود'}), 404
    if user_id in LINK_MONITORS:
        LINK_MONITORS[user_id]['links_found'] = []
        LINK_MONITORS[user_id]['total_links'] = 0
    return jsonify({'success': True, 'message': 'تم مسح الروابط'})


# ============================================================
# قسم الرسائل الذكية — النسخة المطورة مع الجدولة والإرسال الشامل
# ============================================================

SMART_MESSAGES = {}

SCHEDULE_TYPES = {
    'once': 'مرة واحدة',
    'hourly': 'كل ساعة',
    'daily': 'كل يوم',
    'weekly': 'كل أسبوع',
    'custom': 'مخصص (عدد دقائق)'
}

def extract_invite_hash(link):
    """استخراج رمز الدعوة من الرابط"""
    patterns = [
        r't\.me/([a-zA-Z0-9_]+)',
        r'telegram\.me/([a-zA-Z0-9_]+)',
        r'telegram\.dog/([a-zA-Z0-9_]+)',
        r'joinchat/([a-zA-Z0-9_-]+)',
        r'\+([a-zA-Z0-9_-]+)'
    ]
    for pattern in patterns:
        match = re.search(pattern, link)
        if match:
            return match.group(1)
    return None


def should_run_now(schedule_type, last_run_time_iso, custom_interval_minutes=60):
    """
    تحديد ما إذا كان يجب تشغيل المهمة الآن بناءً على نوع الجدولة.
    Returns: (should_run: bool, next_run: datetime|None, wait_seconds: float)
    """
    now = datetime.now()

    # تحويل last_run_time من str إلى datetime
    last_run_time = None
    if last_run_time_iso:
        try:
            last_run_time = datetime.fromisoformat(last_run_time_iso)
        except Exception:
            last_run_time = None

    if schedule_type == 'once':
        if last_run_time is None:
            return True, None, 0
        return False, None, 0

    elif schedule_type == 'hourly':
        if last_run_time is None:
            return True, now + timedelta(hours=1), 0
        next_run = last_run_time + timedelta(hours=1)
        if now >= next_run:
            return True, now + timedelta(hours=1), 0
        return False, next_run, (next_run - now).total_seconds()

    elif schedule_type == 'daily':
        if last_run_time is None:
            return True, now + timedelta(days=1), 0
        next_run = last_run_time + timedelta(days=1)
        if now >= next_run:
            return True, now + timedelta(days=1), 0
        return False, next_run, (next_run - now).total_seconds()

    elif schedule_type == 'weekly':
        if last_run_time is None:
            return True, now + timedelta(days=7), 0
        next_run = last_run_time + timedelta(days=7)
        if now >= next_run:
            return True, now + timedelta(days=7), 0
        return False, next_run, (next_run - now).total_seconds()

    elif schedule_type == 'custom':
        interval = timedelta(minutes=max(1, custom_interval_minutes))
        if last_run_time is None:
            return True, now + interval, 0
        next_run = last_run_time + interval
        if now >= next_run:
            return True, now + interval, 0
        return False, next_run, (next_run - now).total_seconds()

    return False, None, 0


async def _smart_message_worker_async(client_manager, user_id):
    """الدالة الأساسية للرسائل الذكية — مع الجدولة والإرسال الشامل"""
    import time as _time

    state = SMART_MESSAGES.get(user_id)
    if not state or not state.get('is_running'):
        return

    client = client_manager.client
    schedule_type = state.get('schedule_type', 'once')
    custom_interval = int(state.get('custom_interval', 60))
    send_to_all = state.get('send_to_all', False)
    target_links = state.get('links', [])

    # ─── التحقق من الجدولة ───
    should_run, next_run, wait_seconds = should_run_now(
        schedule_type, state.get('last_run_time'), custom_interval
    )

    if not should_run:
        state['next_run_time'] = next_run.isoformat() if next_run else None
        state['wait_seconds'] = int(wait_seconds)
        socketio.emit('smart_message_update', state, room=user_id)
        return

    # ─── تحديد المجموعات المستهدفة ───
    groups = []

    if send_to_all:
        try:
            dialogs = await client.get_dialogs()
            for d in dialogs:
                entity = d.entity
                if hasattr(entity, 'megagroup') or hasattr(entity, 'broadcast') or hasattr(entity, 'left'):
                    username = getattr(entity, 'username', None)
                    title = getattr(entity, 'title', None) or getattr(entity, 'first_name', 'غير معروف')
                    if username:
                        groups.append({'entity': f'@{username}', 'title': title, 'obj': entity})
                    else:
                        groups.append({'entity': entity, 'title': title, 'obj': entity})
            logger.info(f"📊 سيتم الإرسال إلى {len(groups)} مجموعة (جميع المجموعات)")
        except Exception as e:
            logger.error(f"خطأ في جلب جميع المجموعات: {e}")
            state['is_running'] = False
            socketio.emit('smart_message_error', {'error': f'فشل جلب المجموعات: {str(e)}'}, room=user_id)
            return
    else:
        for link in target_links:
            groups.append({'entity': link.strip(), 'title': link.strip(), 'obj': None})
        logger.info(f"📊 سيتم الإرسال إلى {len(groups)} مجموعة (روابط محددة)")

    if not groups:
        state['is_running'] = False
        socketio.emit('smart_message_error', {'error': 'لا توجد مجموعات للإرسال'}, room=user_id)
        return

    # ─── تهيئة إحصائيات الدورة ───
    state['total'] = len(groups)
    state['processed'] = 0
    state['success'] = 0
    state['failed'] = 0
    state['results'] = []

    socketio.emit('smart_message_started', {
        'total': len(groups),
        'send_to_all': send_to_all,
        'schedule_type': schedule_type,
        'next_run': next_run.isoformat() if next_run else None
    }, room=user_id)

    # ─── معالجة كل مجموعة ───
    for group_info in groups:
        if not state.get('is_running', False):
            break

        raw_entity = group_info['entity']
        group_title = group_info['title']
        pre_obj = group_info.get('obj')

        try:
            entity = None

            if pre_obj is not None:
                # entity جاهز من get_dialogs
                entity = pre_obj
            else:
                # تحليل الرابط
                link = str(raw_entity)
                is_private = False
                invite_hash = None
                username = None

                priv_match = re.search(r't(?:elegram)?\.(?:me|dog)/(?:joinchat/|\+)([a-zA-Z0-9_-]+)', link)
                if priv_match:
                    is_private = True
                    invite_hash = priv_match.group(1)
                else:
                    pub_match = re.search(r'(?:t(?:elegram)?\.(?:me|dog)/|@)([a-zA-Z0-9_]+)', link)
                    if pub_match:
                        username = pub_match.group(1)
                    else:
                        username = link.lstrip('@').strip()

                if not invite_hash and not username:
                    state['failed'] += 1
                    state['processed'] += 1
                    state['results'].append({'link': link, 'title': group_title, 'status': 'فشل', 'error': 'رابط غير صالح'})
                    socketio.emit('smart_message_update', state, room=user_id)
                    continue

                try:
                    if is_private and invite_hash:
                        try:
                            result = await client(functions.messages.ImportChatInviteRequest(invite_hash))
                            if hasattr(result, 'chats') and result.chats:
                                entity = result.chats[0]
                        except Exception as join_e:
                            err_lower = str(join_e).lower()
                            if 'already' in err_lower or 'participant' in err_lower:
                                try:
                                    invite_info = await client(functions.messages.CheckChatInviteRequest(invite_hash))
                                    if hasattr(invite_info, 'chat'):
                                        entity = invite_info.chat
                                except Exception:
                                    pass
                            elif 'expired' in err_lower or 'invalid' in err_lower:
                                raise Exception('رابط الدعوة منتهي أو غير صالح')
                            else:
                                raise
                    else:
                        entity = await client.get_entity(username)
                        try:
                            if hasattr(entity, 'megagroup') or hasattr(entity, 'broadcast'):
                                await client(functions.channels.JoinChannelRequest(entity))
                        except Exception:
                            pass

                    if entity is None:
                        raise Exception('تعذّر الوصول إلى المجموعة')

                except (InviteHashExpiredError, InviteHashInvalidError) as inv_e:
                    state['failed'] += 1
                    state['processed'] += 1
                    state['results'].append({'link': link, 'title': group_title, 'status': 'فشل', 'error': f'رابط الدعوة منتهي: {inv_e}'})
                    socketio.emit('smart_message_update', state, room=user_id)
                    continue
                except Exception as join_err:
                    state['failed'] += 1
                    state['processed'] += 1
                    state['results'].append({'link': link, 'title': group_title, 'status': 'فشل', 'error': f'فشل الانضمام: {join_err}'})
                    socketio.emit('smart_message_update', state, room=user_id)
                    continue

            # ─── إرسال الرسالة الأولية ───
            try:
                initial_msg = await client.send_message(entity, state['initial_message'])
                logger.info(f"✅ تم إرسال الرسالة الأولية إلى {group_title}")
            except Exception as send_err:
                state['failed'] += 1
                state['processed'] += 1
                state['results'].append({'title': group_title, 'status': 'فشل', 'error': f'فشل إرسال الرسالة: {send_err}'})
                socketio.emit('smart_message_update', state, room=user_id)
                continue

            # ─── انتظار وصول رسائل من الآخرين ───
            messages_count = 0
            last_message_id = initial_msg.id
            timeout = 300
            t_start = _time.time()

            while messages_count < state['messages_to_wait'] and (_time.time() - t_start) < timeout:
                if not state.get('is_running', False):
                    break
                try:
                    new_msgs = await client.get_messages(entity, limit=20, min_id=last_message_id)
                    others = [m for m in new_msgs if m.id > last_message_id and m.sender_id != initial_msg.sender_id]
                    if others:
                        messages_count += len(others)
                        last_message_id = max(m.id for m in others)
                        logger.info(f"📨 {len(others)} رسائل جديدة في {group_title} (المجموع: {messages_count})")
                    await asyncio.sleep(2)
                except Exception as poll_e:
                    logger.error(f"خطأ في مراقبة الرسائل: {poll_e}")
                    await asyncio.sleep(5)

            # ─── تعديل الرسالة بعد الانتظار ───
            if state.get('is_running', False):
                try:
                    await client.edit_message(entity, initial_msg.id, state['final_message'])
                    state['success'] += 1
                    state['results'].append({'title': group_title, 'status': 'نجاح', 'error': None, 'messages_waited': messages_count})
                    logger.info(f"✏️ تم تعديل الرسالة في {group_title}")
                except Exception as edit_e:
                    state['failed'] += 1
                    state['results'].append({'title': group_title, 'status': 'فشل', 'error': f'فشل التعديل: {edit_e}'})

            state['processed'] += 1
            socketio.emit('smart_message_update', state, room=user_id)
            await asyncio.sleep(3)

        except Exception as e:
            logger.error(f"❌ خطأ في معالجة {group_title}: {e}")
            state['failed'] += 1
            state['processed'] += 1
            state['results'].append({'title': group_title, 'status': 'فشل', 'error': str(e)})
            socketio.emit('smart_message_update', state, room=user_id)

    # ─── تحديث وقت آخر تشغيل ───
    state['last_run_time'] = datetime.now().isoformat()
    state['next_run_time'] = next_run.isoformat() if next_run else None

    if schedule_type == 'once':
        state['is_running'] = False
        state['end_time'] = datetime.now().isoformat()
        socketio.emit('smart_message_done', state, room=user_id)
        if user_id in SMART_MESSAGES:
            del SMART_MESSAGES[user_id]
    else:
        socketio.emit('smart_message_cycle_completed', {
            'next_run': next_run.isoformat() if next_run else None,
            'wait_seconds': int(wait_seconds)
        }, room=user_id)
        socketio.emit('smart_message_update', state, room=user_id)


def run_smart_message_worker(user_id):
    """تشغيل worker الرسائل الذكية في خيط منفصل مع حلقة الجدولة"""
    try:
        with USERS_LOCK:
            user_data = USERS.get(user_id, {})
        client_manager = user_data.get('client_manager')

        if not client_manager or not getattr(client_manager, 'client', None):
            logger.error(f"لا يوجد client_manager للمستخدم {user_id}")
            socketio.emit('smart_message_error', {'error': 'العميل غير متصل'}, room=user_id)
            if user_id in SMART_MESSAGES:
                SMART_MESSAGES[user_id]['is_running'] = False
                del SMART_MESSAGES[user_id]
            return

        state = SMART_MESSAGES.get(user_id, {})
        schedule_type = state.get('schedule_type', 'once')

        while state.get('is_running', False):
            try:
                client_manager.run_coroutine(
                    _smart_message_worker_async(client_manager, user_id)
                )

                if schedule_type == 'once' or not state.get('is_running', False):
                    break

                # انتظار الدورة التالية بفترات قصيرة للسماح بالإيقاف
                wait_seconds = state.get('wait_seconds', 60)
                if wait_seconds > 0:
                    logger.info(f"⏳ الانتظار {wait_seconds:.0f} ثانية للدورة التالية...")
                    elapsed = 0
                    while elapsed < wait_seconds and state.get('is_running', False):
                        time.sleep(min(10, wait_seconds - elapsed))
                        elapsed += 10

            except Exception as cycle_error:
                logger.error(f"خطأ في دورة الجدولة: {cycle_error}")
                time.sleep(60)

    except Exception as e:
        logger.error(f"خطأ في worker الرسائل الذكية: {e}")
        socketio.emit('smart_message_error', {'error': f'خطأ: {str(e)}'}, room=user_id)
        if user_id in SMART_MESSAGES:
            SMART_MESSAGES[user_id]['is_running'] = False


@app.route('/api/smart_message/start', methods=['POST'])
def start_smart_message():
    """بدء عملية الرسائل الذكية مع الجدولة والإرسال الشامل"""
    user_id = session.get('user_id')
    if user_id not in USERS:
        return jsonify({'error': 'المستخدم غير موجود'}), 404

    data = request.json or {}

    initial_message = data.get('initial_message', 'السلام عليكم').strip()
    final_message = data.get('final_message', 'مرحباً بكم في قناتنا الجديدة').strip()
    messages_to_wait = int(data.get('messages_to_wait', 4))
    send_to_all = bool(data.get('send_to_all', False))
    schedule_type = data.get('schedule_type', 'once')
    if schedule_type not in SCHEDULE_TYPES:
        schedule_type = 'once'
    custom_interval = max(1, int(data.get('custom_interval', 60)))
    scheduled_time = data.get('scheduled_time', '')

    links_raw = data.get('links', '')
    links = [l.strip() for l in links_raw.split('\n') if l.strip()]

    if not send_to_all and not links:
        return jsonify({'error': 'الرجاء إدخال روابط المجموعات أو تفعيل الإرسال إلى جميع المجموعات'}), 400

    if not initial_message:
        return jsonify({'error': 'الرجاء إدخال الرسالة التجريبية'}), 400
    if not final_message:
        return jsonify({'error': 'الرجاء إدخال الرسالة المعدلة'}), 400

    if user_id in SMART_MESSAGES and SMART_MESSAGES[user_id].get('is_running', False):
        return jsonify({'error': 'توجد عملية رسائل ذكية جارية بالفعل'}), 400

    with USERS_LOCK:
        user_data = USERS.get(user_id, {})
    client_manager = user_data.get('client_manager')
    if not client_manager or not getattr(client_manager, 'client', None):
        return jsonify({'error': 'العميل غير متصل — يرجى تسجيل الدخول أولاً'}), 400

    try:
        is_auth = client_manager.run_coroutine(client_manager.client.is_user_authorized())
        if not is_auth:
            return jsonify({'error': 'العميل غير مصرح — يرجى تسجيل الدخول مجدداً'}), 400
    except Exception as auth_err:
        return jsonify({'error': f'فشل التحقق من المصادقة: {auth_err}'}), 400

    SMART_MESSAGES[user_id] = {
        'is_running': True,
        'links': links,
        'initial_message': initial_message,
        'final_message': final_message,
        'messages_to_wait': messages_to_wait,
        'send_to_all': send_to_all,
        'schedule_type': schedule_type,
        'custom_interval': custom_interval,
        'scheduled_time': scheduled_time,
        'results': [],
        'start_time': datetime.now().isoformat(),
        'last_run_time': None,
        'next_run_time': None,
        'wait_seconds': 0,
        'total': 0,
        'processed': 0,
        'success': 0,
        'failed': 0
    }

    worker = threading.Thread(target=run_smart_message_worker, args=(user_id,), daemon=True)
    worker.start()
    SMART_MESSAGES[user_id]['thread'] = worker

    schedule_labels = {
        'once': 'مرة واحدة',
        'hourly': 'كل ساعة',
        'daily': 'كل يوم',
        'weekly': 'كل أسبوع',
        'custom': f'كل {custom_interval} دقيقة'
    }
    target_desc = 'جميع المجموعات' if send_to_all else f'{len(links)} مجموعة'
    return jsonify({
        'success': True,
        'message': f'بدأت الرسائل الذكية: {target_desc} • {schedule_labels.get(schedule_type, "")}',
        'send_to_all': send_to_all,
        'schedule_type': schedule_type
    })


@app.route('/api/smart_message/stop')
def stop_smart_message():
    """إيقاف عملية الرسائل الذكية"""
    user_id = session.get('user_id')
    if user_id not in USERS:
        return jsonify({'error': 'المستخدم غير موجود'}), 404
    if user_id in SMART_MESSAGES:
        SMART_MESSAGES[user_id]['is_running'] = False
    return jsonify({'success': True, 'message': 'تم إيقاف الرسائل الذكية'})


@app.route('/api/smart_message/status')
def smart_message_status():
    """الحصول على حالة الرسائل الذكية مع معلومات الجدولة"""
    user_id = session.get('user_id')
    if user_id not in USERS:
        return jsonify({'error': 'المستخدم غير موجود'}), 404
    if user_id in SMART_MESSAGES:
        state = SMART_MESSAGES[user_id].copy()
        state.pop('thread', None)
        return jsonify(state)
    return jsonify({'is_running': False})


# ===========================
# وظيفة رسائلي — تتبع وإدارة الرسائل المرسلة
# ===========================

@app.route("/api/sent_batches")
def api_sent_batches():
    """عرض جميع الدفعات المرسلة (رسائلي)"""
    user_id = session.get('user_id')
    if not user_id or user_id not in USERS:
        return jsonify({"success": False, "message": "غير مسجّل"}), 401
    with USERS_LOCK:
        batches = list(USERS[user_id].get('sent_batches', []))
    result = []
    for b in reversed(batches):
        result.append({
            "id": b["id"],
            "text": b["text"],
            "has_media": b.get("has_media", False),
            "sent_at": b["sent_at"],
            "edited_at": b.get("edited_at"),
            "sent_count": b.get("sent_count", len(b["entries"])),
            "group_count": len(b["entries"]),
            "groups": [{"title": e.get("group", ""), "username": e.get("group", "")} for e in b["entries"]]
        })
    return jsonify({"success": True, "batches": result})


@app.route("/api/edit_batch", methods=["POST"])
def api_edit_batch():
    """تعديل جميع رسائل الدفعة دفعة واحدة"""
    user_id = session.get('user_id')
    if not user_id or user_id not in USERS:
        return jsonify({"success": False, "message": "غير مسجّل"}), 401
    with USERS_LOCK:
        client_manager = USERS[user_id].get('client_manager')
    if not client_manager:
        return jsonify({"success": False, "message": "يجب تسجيل الدخول أولاً"})
    data = request.json or {}
    batch_id = data.get("batch_id", "")
    new_text = data.get("new_text", "")
    if not batch_id or not new_text:
        return jsonify({"success": False, "message": "بيانات ناقصة"})
    def run_edit():
        try:
            client_manager.run_coroutine(
                client_manager._edit_batch_messages(batch_id, new_text)
            )
        except Exception as e:
            socketio.emit('log_update', {"message": f"❌ خطأ في التعديل: {str(e)[:100]}"}, to=user_id)
    _OSThread(target=run_edit, daemon=True).start()
    return jsonify({"success": True, "message": "⏳ جارٍ تعديل الرسائل..."})


@app.route("/api/delete_batch", methods=["POST"])
def api_delete_batch():
    """حذف جميع رسائل الدفعة دفعة واحدة"""
    user_id = session.get('user_id')
    if not user_id or user_id not in USERS:
        return jsonify({"success": False, "message": "غير مسجّل"}), 401
    with USERS_LOCK:
        client_manager = USERS[user_id].get('client_manager')
    if not client_manager:
        return jsonify({"success": False, "message": "يجب تسجيل الدخول أولاً"})
    data = request.json or {}
    batch_id = data.get("batch_id", "")
    if not batch_id:
        return jsonify({"success": False, "message": "batch_id مطلوب"})
    def run_delete():
        try:
            client_manager.run_coroutine(
                client_manager._delete_batch_messages(batch_id)
            )
        except Exception as e:
            socketio.emit('log_update', {"message": f"❌ خطأ في الحذف: {str(e)[:100]}"}, to=user_id)
    _OSThread(target=run_delete, daemon=True).start()
    return jsonify({"success": True, "message": "⏳ جارٍ حذف الرسائل..."})


@app.route("/api/batch_details/<batch_id>")
def api_batch_details(batch_id):
    """عرض تفاصيل دفعة محددة"""
    user_id = session.get('user_id')
    if not user_id or user_id not in USERS:
        return jsonify({"success": False, "message": "غير مسجّل"}), 401
    with USERS_LOCK:
        batches = USERS[user_id].get('sent_batches', [])
        batch = next((b for b in batches if b["id"] == batch_id), None)
    if not batch:
        return jsonify({"success": False, "message": "الدفعة غير موجودة"}), 404
    return jsonify({
        "success": True,
        "batch": {
            "id": batch["id"],
            "text": batch["text"],
            "has_media": batch.get("has_media", False),
            "sent_at": batch["sent_at"],
            "edited_at": batch.get("edited_at"),
            "entries": batch["entries"]
        }
    })


if __name__ == '__main__':
    port = int(os.environ.get("PORT", 5000))
    print(f"🌐 تشغيل الخادم على المنفذ {port}...")
    print(f"🔗 رابط التطبيق: http://0.0.0.0:{port}")
    print("🛡️ نظام الاستمرارية الدائم مُفعل — يعمل حتى الإيقاف اليدوي")
    print("🎓 مركز سرعة انجاز للخدمات الطلابية والأكاديمية - الإصدار المتكامل")
    print("📊 تم دمج النظام الأكاديمي الذكي + منسق المستندات + منشئ العروض")

    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
    )

    network_monitor.start()

    try:
        socketio.run(app, host='0.0.0.0', port=port, debug=False, allow_unsafe_werkzeug=True)
    except Exception as e:
        print(f"❌ خطأ في تشغيل الخادم: {e}")
